"""
OnePilot — Dashboard PowerPoint Generator
Layout : 1 slide unique, grille identique au dashboard consolidé
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import io, json
from datetime import date

# Palette
BG     = RGBColor(0xFF,0xFF,0xFF)
BG2    = RGBColor(0xF0,0xF6,0xFF)
ACCENT = RGBColor(0x00,0x6E,0xB8)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
TEXT1  = RGBColor(0x1A,0x2B,0x3C)
TEXT2  = RGBColor(0x4A,0x6A,0x82)
TEXT3  = RGBColor(0x9A,0xB4,0xC8)
GREEN  = RGBColor(0x00,0x8A,0x50)
RED    = RGBColor(0xC0,0x30,0x30)
BORDER = RGBColor(0xD0,0xE4,0xF4)
CHART_COLORS = [
    RGBColor(0x00,0x6E,0xB8), RGBColor(0x2D,0x72,0xFF),
    RGBColor(0x00,0x8A,0x50), RGBColor(0xD4,0xA0,0x00),
    RGBColor(0xC0,0x30,0x30), RGBColor(0x7C,0x3A,0xED),
]

def I(v): return Inches(v)

def set_bg(slide):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = BG

def rect(slide, x, y, w, h, fill=None, line_color=None, lw=0.3):
    s = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_color: s.line.color.rgb = line_color; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=9, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = "Calibri"
    if color: r.font.color.rgb = color
    return tb

def fmt_val(v):
    if v is None: return "—"
    try:
        n = float(v)
        if abs(n) >= 1e6: return f"{n/1e6:.2f}M €"
        if abs(n) >= 1e3: return f"{n/1e3:.1f}K €"
        return f"{n:,.0f}"
    except: return str(v)[:18]

def add_chart_bar(slide, w, x, y, cw, ch):
    d = w.get("data",{}); labels=(d.get("labels") or [])[:20]
    ds = d.get("datasets") or []
    if not labels or not ds: return
    vals = [float(v or 0) for v in (ds[0].get("data") or [])[:20]]
    cd = ChartData(); cd.categories = labels; cd.add_series(ds[0].get("label",""), vals)
    ct = XL_CHART_TYPE.BAR_CLUSTERED if w.get("chart_type")=="bar_horizontal" else XL_CHART_TYPE.COLUMN_CLUSTERED
    ch_obj = slide.shapes.add_chart(ct, I(x), I(y), I(cw), I(ch), cd).chart
    ch_obj.has_legend = False
    ch_obj.has_title = True; ch_obj.chart_title.text_frame.text = w.get("title","")
    ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT2
    for pt in ch_obj.plots[0].series[0].points:
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = ACCENT

def add_chart_line(slide, w, x, y, cw, ch):
    d = w.get("data",{}); labels = d.get("labels") or []
    ds = d.get("datasets") or []
    if not labels or not ds: return
    cd = ChartData(); cd.categories = labels
    for i, s in enumerate(ds[:3]):
        cd.add_series(s.get("label",f"S{i+1}"), [float(v or 0) for v in s.get("data",[])])
    ch_obj = slide.shapes.add_chart(XL_CHART_TYPE.LINE, I(x), I(y), I(cw), I(ch), cd).chart
    ch_obj.has_legend = len(ds) > 1
    ch_obj.has_title = True; ch_obj.chart_title.text_frame.text = w.get("title","")
    ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT2
    for i, s in enumerate(ch_obj.plots[0].series):
        s.format.line.color.rgb = CHART_COLORS[i%len(CHART_COLORS)]
        s.format.line.width = Pt(2)

def add_chart_pie(slide, w, x, y, cw, ch):
    d = w.get("data",{}); labels=(d.get("labels") or [])[:8]
    ds = d.get("datasets") or []
    if not labels or not ds: return
    vals = [float(v or 0) for v in (ds[0].get("data") or [])[:8]]
    cd = ChartData(); cd.categories = labels; cd.add_series(w.get("title",""), vals)
    ct = XL_CHART_TYPE.DOUGHNUT if w.get("chart_type")=="doughnut" else XL_CHART_TYPE.PIE
    ch_obj = slide.shapes.add_chart(ct, I(x), I(y), I(cw), I(ch), cd).chart
    ch_obj.has_legend = True
    ch_obj.has_title = True; ch_obj.chart_title.text_frame.text = w.get("title","")
    ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT2
    for i, pt in enumerate(ch_obj.plots[0].series[0].points):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = CHART_COLORS[i%len(CHART_COLORS)]

def add_kpi(slide, w, x, y, cw, ch):
    d = w.get("data",{})
    rect(slide, x, y, cw, ch, fill=BG2, line_color=BORDER)
    txt(slide, w.get("title","")[:25], x+0.08, y+0.05, cw-0.16, 0.22,
        size=7, color=TEXT2)
    txt(slide, fmt_val(d.get("value")), x+0.08, y+0.25, cw-0.16, 0.45,
        size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    delta = d.get("delta")
    if delta is not None:
        try:
            dv = float(delta)
            col = GREEN if dv >= 0 else RED
            arrow = "▲" if dv >= 0 else "▼"
            txt(slide, f"{arrow} {abs(dv):.1f}%", x+0.08, y+0.68, cw-0.16, 0.2,
                size=7, color=col, align=PP_ALIGN.CENTER)
        except: pass

def add_table(slide, w, x, y, cw, ch):
    d = w.get("data",{}); headers = d.get("headers") or []; rows=(d.get("rows") or [])[:12]
    if not headers: return
    nc = len(headers); nr = len(rows)+1
    col_w = min(cw/max(nc,1), 2.5)
    tbl = slide.shapes.add_table(nr, nc, I(x), I(y), I(col_w*nc), I(min(ch,nr*0.28+0.05))).table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0,ci); cell.text = str(h)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        r = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else cell.text_frame.paragraphs[0].add_run()
        r.font.bold=True; r.font.size=Pt(7); r.font.color.rgb=WHITE; r.font.name="Calibri"
    for ri, row in enumerate(rows):
        bg = BG2 if ri%2==0 else BG
        for ci, val in enumerate(row[:nc]):
            cell = tbl.cell(ri+1,ci); cell.text = str(val if val is not None else "")
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            r = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else cell.text_frame.paragraphs[0].add_run()
            r.font.size=Pt(7); r.font.color.rgb=TEXT1; r.font.name="Calibri"

def generate_pptx(spec: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # ── Header ──────────────────────────────────────────────────────────
    rect(slide, 0, 0, 13.33, 0.42, fill=ACCENT)
    txt(slide, spec.get("title","Dashboard"), 0.25, 0.05, 9, 0.34,
        size=14, bold=True, color=WHITE)
    txt(slide, f"OnePilot  ·  {date.today().strftime('%d/%m/%Y')}",
        10.5, 0.1, 2.7, 0.25, size=8, color=RGBColor(0xCC,0xEA,0xFF), align=PP_ALIGN.RIGHT)

    # ── Layout grille identique au dashboard consolidé ───────────────────
    widgets = spec.get("widgets", [])
    if not widgets:
        txt(slide, "Aucun widget", 5, 3.5, 3, 0.5, size=14, color=TEXT3, align=PP_ALIGN.CENTER)
        _save_and_return(prs); return

    kpis   = [w for w in widgets if w.get("chart_type") in ("kpi_card","kpi_delta")]
    lines  = [w for w in widgets if w.get("chart_type") in ("line","area","waterfall")]
    bars   = [w for w in widgets if w.get("chart_type") in ("bar","bar_horizontal","column")]
    pies   = [w for w in widgets if w.get("chart_type") in ("pie","doughnut")]
    tables = [w for w in widgets if w.get("chart_type") in ("table","pivot")]
    others = [w for w in widgets if w not in kpis+lines+bars+pies+tables]

    MARGIN = 0.18
    y_cur  = 0.48  # juste après header
    SW     = 13.33 # slide width
    AVAIL_H = 6.78 # hauteur disponible (7.5 - 0.42 header - 0.3 footer)

    # Calcul dynamique des hauteurs selon ce qu'on a
    has_kpi  = len(kpis) > 0
    has_line = len(lines) > 0
    has_pie  = len(pies) > 0
    has_bar  = len(bars) > 0
    has_tbl  = len(tables) > 0

    rows_count = sum([has_kpi, has_line, (has_bar or has_pie), has_tbl])
    if rows_count == 0: rows_count = 1

    KPI_H  = 0.95 if has_kpi else 0
    LINE_H = (AVAIL_H - KPI_H - MARGIN*rows_count) * (0.38 if (has_bar or has_pie or has_tbl) else 0.8) if has_line else 0
    MID_H  = (AVAIL_H - KPI_H - LINE_H - MARGIN*rows_count) * (0.55 if has_tbl else 1.0) if (has_bar or has_pie) else 0
    TBL_H  = AVAIL_H - KPI_H - LINE_H - MID_H - MARGIN*rows_count if has_tbl else 0

    # ── Rangée 1 : KPI cards ──────────────────────────────────────────
    if kpis:
        n = len(kpis)
        kw = (SW - 2*MARGIN - (n-1)*MARGIN) / n
        for i, w in enumerate(kpis):
            kx = MARGIN + i*(kw+MARGIN)
            add_kpi(slide, w, kx, y_cur, kw, KPI_H-0.05)
        y_cur += KPI_H + MARGIN

    # ── Rangée 2 : Line chart ─────────────────────────────────────────
    if lines:
        for w in lines[:1]:  # 1 line chart pleine largeur
            add_chart_line(slide, w, MARGIN, y_cur, SW-2*MARGIN, LINE_H-0.05)
        y_cur += LINE_H + MARGIN

    # ── Rangée 3 : Bar + Pie côte à côte ─────────────────────────────
    if bars or pies:
        x_cur = MARGIN
        all_mid = bars + pies
        n = len(all_mid)
        if n == 1:
            w = all_mid[0]
            if w.get("chart_type") in ("pie","doughnut"):
                add_chart_pie(slide, w, MARGIN, y_cur, SW-2*MARGIN, MID_H-0.05)
            else:
                add_chart_bar(slide, w, MARGIN, y_cur, SW-2*MARGIN, MID_H-0.05)
        else:
            # Bar prend 2/3, pie prend 1/3
            bar_w = (SW - 2*MARGIN - MARGIN) * (2/3 if pies else 1.0)
            pie_w = SW - 2*MARGIN - MARGIN - bar_w
            for w in bars[:1]:
                add_chart_bar(slide, w, x_cur, y_cur, bar_w, MID_H-0.05)
                x_cur += bar_w + MARGIN
            for w in pies[:1]:
                add_chart_pie(slide, w, x_cur, y_cur, pie_w, MID_H-0.05)
        y_cur += MID_H + MARGIN

    # ── Rangée 4 : Table pleine largeur ──────────────────────────────
    if tables and TBL_H > 0.5:
        add_table(slide, tables[0], MARGIN, y_cur, SW-2*MARGIN, TBL_H-0.05)

    # ── Footer ────────────────────────────────────────────────────────
    txt(slide, "OnePilot — Agent Conversationnel ERP",
        0, 7.2, 13.33, 0.25, size=7, color=TEXT3, align=PP_ALIGN.CENTER, italic=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def generate_pptx_multi(specs: list) -> bytes:
    """Génère un PPTX unique avec 1 slide par spec — tout dans une seule présentation."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for spec in specs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)

        widgets = spec.get("widgets", [])
        kpis   = [w for w in widgets if w.get("chart_type") in ("kpi_card","kpi_delta")]
        lines  = [w for w in widgets if w.get("chart_type") in ("line","area","waterfall")]
        bars   = [w for w in widgets if w.get("chart_type") in ("bar","bar_horizontal","column")]
        pies   = [w for w in widgets if w.get("chart_type") in ("pie","doughnut")]
        tables = [w for w in widgets if w.get("chart_type") in ("table","pivot")]

        MARGIN = 0.18; SW = 13.33; AVAIL_H = 6.78
        has_kpi = len(kpis)>0; has_line = len(lines)>0
        has_mid = len(bars)>0 or len(pies)>0; has_tbl = len(tables)>0
        rows_count = sum([has_kpi,has_line,has_mid,has_tbl]) or 1

        KPI_H  = 0.95 if has_kpi else 0
        LINE_H = (AVAIL_H-KPI_H-MARGIN*rows_count)*(0.38 if (has_mid or has_tbl) else 0.85) if has_line else 0
        MID_H  = (AVAIL_H-KPI_H-LINE_H-MARGIN*rows_count)*(0.55 if has_tbl else 1.0) if has_mid else 0
        TBL_H  = AVAIL_H-KPI_H-LINE_H-MID_H-MARGIN*rows_count if has_tbl else 0

        # Header
        rect(slide, 0, 0, SW, 0.42, fill=ACCENT)
        txt(slide, spec.get("title","Dashboard"), 0.25, 0.05, 9, 0.34, size=14, bold=True, color=WHITE)
        txt(slide, f"OnePilot · {date.today().strftime('%d/%m/%Y')}", 10.5, 0.1, 2.7, 0.25,
            size=8, color=RGBColor(0xCC,0xEA,0xFF), align=PP_ALIGN.RIGHT)

        y_cur = 0.48
        if kpis:
            n=len(kpis); kw=(SW-2*MARGIN-(n-1)*MARGIN)/n
            for i,w in enumerate(kpis): add_kpi(slide,w,MARGIN+i*(kw+MARGIN),y_cur,kw,KPI_H-0.05)
            y_cur+=KPI_H+MARGIN
        if lines:
            add_chart_line(slide,lines[0],MARGIN,y_cur,SW-2*MARGIN,LINE_H-0.05)
            y_cur+=LINE_H+MARGIN
        if bars or pies:
            all_mid=bars+pies
            if len(all_mid)==1:
                w=all_mid[0]
                if w.get("chart_type") in ("pie","doughnut"): add_chart_pie(slide,w,MARGIN,y_cur,SW-2*MARGIN,MID_H-0.05)
                else: add_chart_bar(slide,w,MARGIN,y_cur,SW-2*MARGIN,MID_H-0.05)
            else:
                bar_w=(SW-2*MARGIN-MARGIN)*(2/3 if pies else 1.0)
                pie_w=SW-2*MARGIN-MARGIN-bar_w; xc=MARGIN
                for w in bars[:1]: add_chart_bar(slide,w,xc,y_cur,bar_w,MID_H-0.05); xc+=bar_w+MARGIN
                for w in pies[:1]: add_chart_pie(slide,w,xc,y_cur,pie_w,MID_H-0.05)
            y_cur+=MID_H+MARGIN
        if tables and TBL_H>0.5:
            add_table(slide,tables[0],MARGIN,y_cur,SW-2*MARGIN,TBL_H-0.05)

        txt(slide, "OnePilot — Agent Conversationnel ERP", 0, 7.2, SW, 0.25,
            size=7, color=TEXT3, align=PP_ALIGN.CENTER, italic=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    import sys
    with open(sys.argv[1]) as f: spec = json.load(f)
    data = generate_pptx(spec)
    out = sys.argv[2] if len(sys.argv)>2 else "/tmp/out.pptx"
    with open(out,"wb") as f: f.write(data)
    print(f"OK:{out} ({len(data)//1024}KB)")
