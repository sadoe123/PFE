"""
OnePilot – API FastAPI
Universal Data Access Layer – Phase 3
"""
from __future__ import annotations

import asyncio
import subprocess
import tempfile
import csv as csv_module
import io
import json as json_module
import logging
import os
from contextlib import asynccontextmanager
from datetime import datetime
from typing import Any, Dict, List, Optional
from uuid import UUID, uuid4

from fastapi import BackgroundTasks, FastAPI, File, HTTPException, Query, UploadFile, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

from .connection_service import sync_metadata, test_connection
from .database import close_connections, get_pg_pool, get_redis, init_schema
from .relationship_discovery import (
    discover_relationships,
    get_join_paths_for_source,
    get_relations_for_source,
    get_false_positive_pairs,
    remove_false_positive,
    get_validation_stats,
    get_alternative_suggestions,
    record_expert_feedback,
    validate_relation,
)
from .repository import (
    create_source,
    delete_source,
    get_source,
    get_source_with_entities,
    list_sources,
    update_source,
)
from .cdc_engine             import CDCEngine
from .schema_versioner       import SchemaVersioner
from .cdc_file_watcher       import FileCDCEngine
from .cdc_reindexer          import CDCReindexer
from .metadata_enricher      import MetadataEnricher
from .graphql_relation_saver import GraphQLRelationSaver, HATEOASLinkExtractor
from .cdc_db_triggers        import PostgreSQLWALCDC, SQLServerCDC
from .nlu_engine             import get_nlu_pipeline, get_context, clear_context, Intent, ContextManager, ConversationTurn, retrain_fasttext_with_feedback
from .query_engine           import SQLGenerator, UniversalQueryPlanner
from .ambiguity_resolver     import AmbiguityResolver
try:
    from connectors.factory  import ConnectorFactory
except ImportError:
    ConnectorFactory = None

from .schemas import (
    ConnectorType,
    ConnectionTestResult,
    DataSourceCreate,
    DataSourceDetail,
    DataSourceList,
    DataSourceOut,
    DataSourceUpdate,
    MetadataSyncResult,
    SourceCategory,
)

# ── Logging ───────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logger = logging.getLogger(__name__)

def _safe_json(obj):
    """Nettoie un dict/list pour JSON — remplace NaN/Infinity par None."""
    import math
    if isinstance(obj, dict):
        return {k: _safe_json(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_safe_json(v) for v in obj]
    if isinstance(obj, float):
        if math.isnan(obj) or math.isinf(obj):
            return None
        return round(obj, 6)
    return obj

UPLOAD_DIR = os.environ.get("UPLOAD_DIR", "/app/uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)


# ══════════════════════════════════════════════════════════════
# LIFESPAN
# ══════════════════════════════════════════════════════════════

@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("🚀 OnePilot API démarrage...")
    try:
        await get_pg_pool()
        await init_schema()
        logger.info("✅ PostgreSQL connecté et schema initialisé")

        # ── Migration : table conversation_dashboards ─────────────────────
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                await conn.execute("""
                    CREATE TABLE IF NOT EXISTS conversation_dashboards (
                        conv_id      TEXT NOT NULL,
                        dashboard_id TEXT NOT NULL,
                        spec_json    TEXT NOT NULL,
                        updated_at   TIMESTAMPTZ DEFAULT NOW(),
                        PRIMARY KEY (conv_id, dashboard_id)
                    )
                """)
                await conn.execute("""
                    CREATE INDEX IF NOT EXISTS idx_conv_dashboards_conv_id
                    ON conversation_dashboards(conv_id)
                """)
            logger.info("✅ Migration conversation_dashboards OK")
        except Exception as e:
            logger.warning(f"Migration conversation_dashboards (non-bloquant): {e}")
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                await conn.execute("""
                    CREATE TABLE IF NOT EXISTS user_favorites (
                        id           SERIAL PRIMARY KEY,
                        user_id      TEXT NOT NULL DEFAULT 'admin',
                        dashboard_id TEXT NOT NULL,
                        title        TEXT NOT NULL,
                        spec_json    TEXT NOT NULL,
                        source_id    TEXT,
                        conv_id      TEXT,
                        created_at   TIMESTAMPTZ DEFAULT NOW(),
                        UNIQUE(user_id, dashboard_id)
                    )
                """)
                await conn.execute("CREATE INDEX IF NOT EXISTS idx_user_favorites_user_id ON user_favorites(user_id)")
                await conn.execute("ALTER TABLE user_favorites ADD COLUMN IF NOT EXISTS conv_id TEXT")
            logger.info("✅ Migration user_favorites OK")
        except Exception as e:
            logger.warning(f"Migration user_favorites (non-bloquant): {e}")
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                await conn.execute("""
                    CREATE TABLE IF NOT EXISTS shared_dashboards (
                        id           SERIAL PRIMARY KEY,
                        token        TEXT NOT NULL UNIQUE,
                        title        TEXT NOT NULL,
                        spec_json    TEXT NOT NULL,
                        source_id    TEXT,
                        created_by   TEXT DEFAULT 'admin',
                        expires_at   TIMESTAMPTZ NOT NULL,
                        created_at   TIMESTAMPTZ DEFAULT NOW(),
                        view_count   INTEGER DEFAULT 0
                    )
                """)
                await conn.execute("CREATE INDEX IF NOT EXISTS idx_shared_token ON shared_dashboards(token)")
            logger.info("✅ Migration shared_dashboards OK")
        except Exception as e:
            logger.warning(f"Migration shared_dashboards (non-bloquant): {e}")
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                await conn.execute("""
                    ALTER TABLE entity_fields
                    ADD COLUMN IF NOT EXISTS dimension_hierarchy JSONB DEFAULT NULL
                """)
            logger.info("✅ Migration dimension_hierarchy OK")
        except Exception as e:
            logger.warning(f"Migration dimension_hierarchy (non-bloquant): {e}")

        # ── Migration : table cdc_subscriber_log ─────────────────────
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                await conn.execute("""
                    CREATE TABLE IF NOT EXISTS cdc_subscriber_log (
                        id                  SERIAL PRIMARY KEY,
                        source_id           UUID NOT NULL,
                        version             INTEGER NOT NULL,
                        change_count        INTEGER DEFAULT 0,
                        cache_invalidated   BOOLEAN DEFAULT FALSE,
                        reindex_triggered   BOOLEAN DEFAULT FALSE,
                        processed_at        TIMESTAMPTZ DEFAULT NOW(),
                        UNIQUE (source_id, version)
                    )
                """)
                await conn.execute("""
                    CREATE INDEX IF NOT EXISTS idx_cdc_subscriber_source
                    ON cdc_subscriber_log(source_id, processed_at DESC)
                """)
            logger.info("✅ Migration cdc_subscriber_log OK")
        except Exception as e:
            logger.warning(f"Migration cdc_subscriber_log (non-bloquant): {e}")

        # ── Migration : table chat_feedback ──────────────────────────
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                await conn.execute("""
                    CREATE TABLE IF NOT EXISTS chat_feedback (
                        id              SERIAL PRIMARY KEY,
                        conversation_id TEXT,
                        source_id       UUID,
                        message_id      TEXT,
                        question        TEXT,
                        answer          TEXT,
                        intent          TEXT,
                        confidence      FLOAT,
                        nlu_method      TEXT,
                        feedback_type   TEXT NOT NULL,  -- 'like' ou 'dislike'
                        used_for_training BOOLEAN DEFAULT FALSE,
                        created_at      TIMESTAMPTZ DEFAULT NOW()
                    )
                """)
                await conn.execute("""
                    CREATE INDEX IF NOT EXISTS idx_chat_feedback_source
                    ON chat_feedback(source_id, feedback_type, created_at DESC)
                """)
                await conn.execute("""
                    CREATE INDEX IF NOT EXISTS idx_chat_feedback_training
                    ON chat_feedback(used_for_training, feedback_type)
                """)
            logger.info("✅ Migration chat_feedback OK")
        except Exception as e:
            logger.warning(f"Migration chat_feedback (non-bloquant): {e}")

    except Exception as e:
        logger.error(f"❌ PostgreSQL erreur: {e}")
    try:
        await get_redis()
        logger.info("✅ Redis connecté")
    except Exception as e:
        logger.warning(f"⚠️  Redis non disponible: {e}")

    # ── Démarrage WAL CDC listener ────────────────────────────────────
    _wal_task = None
    try:
        from .cdc_db_triggers import PostgreSQLWALCDC
        pool = await get_pg_pool()
        redis = await get_redis()
        wal_cdc = PostgreSQLWALCDC(pool, redis)
        # Vérifie si le slot existe
        # Vérifie directement si le slot existe dans PostgreSQL
        pool = await get_pg_pool()
        async with pool.acquire() as _wal_conn:
            _slot = await _wal_conn.fetchrow(
                "SELECT slot_name, active FROM pg_replication_slots WHERE slot_name = 'onepilot_cdc'"
            )
        if _slot:
            # Démarre le listener en background
            async def _wal_loop():
                while True:
                    try:
                        # Lit les changements WAL via connexion interne
                        _pool2 = await get_pg_pool()
                        async with _pool2.acquire() as _wconn:
                            await _wconn.fetch(f"""
                                SELECT lsn, xid, data
                                FROM pg_logical_slot_peek_changes(
                                    'onepilot_cdc', NULL, 10,
                                    'include-schemas', 'true'
                                )
                            """)
                    except Exception as _e:
                        logger.debug(f"[WAL] loop error: {_e}")
                    await asyncio.sleep(10)
            _wal_task = asyncio.create_task(_wal_loop())
            logger.info("✅ WAL CDC listener démarré (slot: onepilot_cdc)")
        else:
            logger.info("ℹ️  WAL slot non trouvé — WAL CDC désactivé")
    except Exception as e:
        logger.warning(f"⚠️  WAL CDC non disponible: {e}")

    # ── Reprise automatique des descriptions LLM ────────────────────
    # Si des sources ont des tables sans description après un redémarrage,
    # relancer automatiquement le générateur en arrière-plan
    async def _auto_resume_descriptions():
        await asyncio.sleep(30)  # Attendre 30s que tout soit initialisé
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                sources_to_enrich = await conn.fetch("""
                    SELECT ds.id, ds.name, ds.connector_type, COUNT(*) as undesc_count
                    FROM data_sources ds
                    JOIN source_entities se ON se.source_id = ds.id
                    WHERE ds.status = 'active'
                      AND se.is_visible = TRUE
                      AND (se.description IS NULL OR se.description = ''
                           OR se.description NOT LIKE '[LLM]%')
                    GROUP BY ds.id, ds.name, ds.connector_type
                    HAVING COUNT(*) > 10
                    ORDER BY undesc_count DESC
                """)

            if sources_to_enrich:
                logger.info(f"[DescGen Auto] {len(sources_to_enrich)} sources à enrichir au démarrage")
                from .llm_description_generator import run_description_generation_bg
                for src in sources_to_enrich:
                    src_type = src['connector_type'] if src['connector_type'] else 'unknown'
                    logger.info(f"[DescGen Auto] Reprise: {src['name']} ({src['id']})")
                    asyncio.create_task(run_description_generation_bg(
                        source_id   = src['id'],
                        source_name = src['name'],
                        source_type = src_type,
                        pg_pool     = pool,
                        limit       = 2000,
                    ))
            else:
                logger.info("[DescGen Auto] Toutes les sources sont déjà enrichies")
        except Exception as e:
            logger.warning(f"[DescGen Auto] Erreur reprise auto (non-bloquant): {e}")

    asyncio.create_task(_auto_resume_descriptions())
    logger.info("[DescGen Auto] Reprise automatique programmée dans 30s")

    # ── CDC Breaking Change Subscriber ──────────────────────────────
    # Thread séparé avec redis sync — listen() bloque l'event loop asyncio
    _cdc_thread = None
    _cdc_stop   = False
    _loop_ref   = asyncio.get_event_loop()

    def _cdc_thread_fn():
        """Thread dédié Pub/Sub — redis sync, pas asyncio."""
        import os, redis as _redis_sync, json as _json, logging as _log
        _tlog = _log.getLogger(__name__)

        REDIS_HOST = os.environ.get("REDIS_HOST", "onepilot_redis")
        REDIS_PORT = int(os.environ.get("REDIS_PORT", "6379"))

        while not _cdc_stop:
            _r = None
            try:
                _r  = _redis_sync.Redis(host=REDIS_HOST, port=REDIS_PORT, decode_responses=True)
                _ps = _r.pubsub()
                _ps.psubscribe("cdc:breaking:*")
                _tlog.info("[CDC Thread] Subscriber démarré (redis sync)")

                for msg in _ps.listen():
                    if _cdc_stop:
                        break
                    if msg.get("type") not in ("pmessage", "message"):
                        continue
                    try:
                        data          = _json.loads(msg["data"])
                        source_id_str = data.get("source_id", "")
                        version       = data.get("version", 0)
                        changes       = data.get("changes", [])
                        count         = data.get("count", 0)
                        if not source_id_str:
                            continue

                        _tlog.warning(
                            f"[CDC Subscriber] ⚠️  Breaking change reçu — "
                            f"source={source_id_str} v{version} ({count} changements)"
                        )

                        # Délègue les actions async à l'event loop principal
                        async def _handle(sid=source_id_str, ver=version, chg=changes, cnt=count):
                            try:
                                from uuid import UUID as _UUID
                                _redis = await get_redis()
                                _pool  = await get_pg_pool()

                                # 1. Invalide cache Redis
                                dash_k = await _redis.keys(f"onepilot:dashboard:{sid}:*")
                                step_k = await _redis.keys("onepilot:step:*")
                                prof_k = await _redis.keys(f"onepilot:profile:{sid}:*")
                                all_k  = dash_k + step_k + prof_k
                                if all_k:
                                    await _redis.delete(*all_k)
                                _tlog.info(
                                    f"[CDC Subscriber] {len(all_k)} clés cache "
                                    f"invalidées pour source {sid}"
                                )

                                # 2. Réindexation MeiliSearch
                                from .cdc_reindexer import CDCReindexer
                                reindexer = CDCReindexer(_pool, _redis)
                                res = await reindexer.reindex_after_breaking_change(
                                    source_id=_UUID(sid), version=ver, changes=chg
                                )
                                _tlog.info(
                                    f"[CDC Subscriber] Réindexation — "
                                    f"deleted={res.get('deleted',0)}, "
                                    f"reindexed={res.get('reindexed',0)}"
                                )

                                # 3. Log en base
                                async with _pool.acquire() as _conn:
                                    await _conn.execute("""
                                        INSERT INTO cdc_subscriber_log
                                            (source_id, version, change_count,
                                             cache_invalidated, reindex_triggered, processed_at)
                                        VALUES ($1,$2,$3,TRUE,TRUE,NOW())
                                        ON CONFLICT (source_id, version) DO UPDATE
                                        SET cache_invalidated=TRUE, reindex_triggered=TRUE,
                                            processed_at=NOW()
                                    """, _UUID(sid), ver, cnt)
                            except Exception as _ex:
                                _tlog.error(f"[CDC Subscriber] Handle error: {_ex}")

                        asyncio.run_coroutine_threadsafe(_handle(), _loop_ref)

                    except Exception as _me:
                        _tlog.warning(f"[CDC Thread] Message error: {_me}")

            except Exception as _e:
                if not _cdc_stop:
                    _tlog.warning(f"[CDC Thread] Reconnexion dans 5s ({_e})")
                    import time; time.sleep(5)
            finally:
                try:
                    if _r: _r.close()
                except Exception:
                    pass

    try:
        import threading
        _cdc_thread = threading.Thread(
            target=_cdc_thread_fn, daemon=True, name="cdc-subscriber"
        )
        _cdc_thread.start()
        logger.info("✅ CDC Breaking Change Subscriber programmé")
    except Exception as e:
        logger.warning(f"⚠️  CDC Subscriber non démarré: {e}")


    await asyncio.sleep(0.2)  # Laisse les tasks démarrer
    yield

    # ── Arrêt CDC Subscriber (thread) ───────────────────────────────
    _cdc_stop = True
    if _cdc_thread and _cdc_thread.is_alive():
        _cdc_thread.join(timeout=2.0)
        logger.info("[CDC Thread] Arrêté")

    # ── Arrêt WAL listener ────────────────────────────────────────────
    if _wal_task:
        _wal_task.cancel()
        try: await _wal_task
        except asyncio.CancelledError: pass

    await close_connections()
    logger.info("👋 OnePilot API arrêté")


# ══════════════════════════════════════════════════════════════
# APP
# ══════════════════════════════════════════════════════════════

app = FastAPI(redirect_slashes=False, 
    title="OnePilot API",
    description="Universal Data Access Layer – Phase 3",
    version="3.0.0",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ══════════════════════════════════════════════════════════════
# SCHÉMAS PYDANTIC
# ══════════════════════════════════════════════════════════════

class FkRelation(BaseModel):
    source_entity: str
    source_field:  str
    target_entity: str
    target_field:  str


class ImportFkRequest(BaseModel):
    relations: List[FkRelation]


class ValidateRelationRequest(BaseModel):
    is_confirmed:  bool
    validated_by:  Optional[str] = "expert"
    reject_reason: Optional[str] = None


class FeedbackRequest(BaseModel):
    """PATCH P6 — Feedback expert sur une relation détectée."""
    source_entity:  str
    source_field:   str
    target_entity:  str
    target_field:   str
    feedback:       str            # 'confirmed' | 'rejected' | 'alternative'
    user_id:        Optional[str]  = None
    comment:        Optional[str]  = None
    alternative_target_entity: Optional[str] = None
    alternative_target_field:  Optional[str] = None


# ══════════════════════════════════════════════════════════════
# HEALTH
# ══════════════════════════════════════════════════════════════

@app.get("/", tags=["Health"])
async def root():
    return {"service": "OnePilot", "version": "3.0.0", "status": "running", "docs": "/docs"}


@app.get("/health", tags=["Health"])
async def health():
    checks = {}
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.fetchval("SELECT 1")
        checks["postgres"] = "ok"
    except Exception as e:
        checks["postgres"] = f"error: {e}"
    try:
        r = await get_redis()
        await r.ping()
        checks["redis"] = "ok"
    except Exception as e:
        checks["redis"] = f"unavailable: {e}"

    all_ok = all(v == "ok" for v in checks.values())
    return JSONResponse(
        status_code=200 if all_ok else 207,
        content={"status": "healthy" if all_ok else "degraded", "checks": checks},
    )


# ══════════════════════════════════════════════════════════════
# META — CONNECTOR TYPES
# ══════════════════════════════════════════════════════════════

@app.get("/connector-types", tags=["Meta"])
async def get_connector_types():
    return {
        "database": {
            "label": "Base de données directe",
            "description": "Connexion directe via driver SQL",
            "icon": "database",
            "types": [
                {"id": "postgresql", "label": "PostgreSQL",  "icon": "🐘", "default_port": 5432},
                {"id": "mysql",      "label": "MySQL",       "icon": "🐬", "default_port": 3306},
                {"id": "mssql",      "label": "SQL Server",  "icon": "🪟", "default_port": 1433},
                {"id": "sqlite",     "label": "SQLite",      "icon": "📁", "default_port": None},
                {"id": "sage_100",   "label": "SAGE 100",    "icon": "🟢", "default_port": 1433},
            ],
        },
        "webservice": {
            "label": "Web Service / ERP",
            "description": "Connexion via protocole HTTP ou ERP",
            "icon": "globe",
            "types": [
                {"id": "rest",        "label": "REST API",            "icon": "FK", "default_port": None},
                {"id": "odata",       "label": "OData",               "icon": "⚡", "default_port": None},
                {"id": "graphql",     "label": "GraphQL",             "icon": "◈",  "default_port": None},
                {"id": "soap",        "label": "SOAP / WSDL",         "icon": "📮", "default_port": None},
                {"id": "sap_rfc",     "label": "SAP RFC/BAPI",        "icon": "🔷", "default_port": 3300},
                {"id": "sap_odata",   "label": "SAP OData (S/4HANA)", "icon": "🔶", "default_port": 443},
                {"id": "dynamics365", "label": "Dynamics 365",        "icon": "🟦", "default_port": None},
                {"id": "sage_x3",     "label": "SAGE X3",             "icon": "🟩", "default_port": None},
                {"id": "sage_cloud",  "label": "SAGE Business Cloud", "icon": "☁️", "default_port": None},
            ],
        },
        "file": {
            "label": "Fichiers",
            "description": "Import de fichiers locaux ou upload",
            "icon": "file",
            "types": [
                {"id": "file_csv",     "label": "CSV",                    "icon": "📄", "default_port": None},
                {"id": "file_excel",   "label": "Excel (.xls)",            "icon": "📊", "default_port": None},
                {"id": "file_xlsx",    "label": "Excel multi-feuilles",    "icon": "📗", "default_port": None},
                {"id": "file_json",    "label": "JSON",                    "icon": "📋", "default_port": None},
                {"id": "file_parquet", "label": "Parquet (Big Data)",      "icon": "🗃", "default_port": None},
                {"id": "file_avro",    "label": "Avro (Streaming/Kafka)",  "icon": "🗄", "default_port": None},
            ],
        },
    }


# ══════════════════════════════════════════════════════════════
# SOURCES CRUD
# ══════════════════════════════════════════════════════════════

@app.post("/sources", response_model=DataSourceOut, status_code=201, tags=["Sources"])
async def create_data_source(data: DataSourceCreate):
    try:
        return await create_source(data)
    except ValueError as e:
        raise HTTPException(422, str(e))
    except Exception as e:
        logger.error(f"[API] create_source: {e}", exc_info=True)
        raise HTTPException(500, str(e))


@app.get("/sources", response_model=DataSourceList, tags=["Sources"])
async def list_data_sources(
    category: Optional[str] = Query(None),
    status:   Optional[str] = Query(None),
    search:   Optional[str] = Query(None),
):
    sources = await list_sources(category=category, status=status, search=search)
    return DataSourceList(total=len(sources), sources=sources)


@app.get("/sources/{source_id}", tags=["Sources"])
async def get_data_source(
    source_id: UUID,
    page:      int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=200),
    search:    str = Query(""),
):
    source = await get_source_with_entities(
        source_id, page=page, page_size=page_size, search=search
    )
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    pagination = source.__dict__.get("_pagination", {})
    data = source.model_dump()
    data["pagination"] = pagination
    return data


@app.patch("/sources/{source_id}", response_model=DataSourceOut, tags=["Sources"])
async def update_data_source(source_id: UUID, data: DataSourceUpdate):
    source = await update_source(source_id, data)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    return source


@app.delete("/sources/{source_id}", status_code=204, tags=["Sources"])
async def delete_data_source(source_id: UUID):
    deleted = await delete_source(source_id)
    if not deleted:
        raise HTTPException(404, f"Source {source_id} introuvable")


# ══════════════════════════════════════════════════════════════
# TEST + SYNC
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/test", response_model=ConnectionTestResult, tags=["Connections"])
async def test_source_connection(source_id: UUID):
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    result = await test_connection(source_id)
    return ConnectionTestResult(
        source_id=source_id,
        success=result["success"],
        message=result["message"],
        latency_ms=result.get("latency_ms", -1),
        tested_at=result.get("tested_at") or datetime.utcnow(),
    )


async def _run_sync_background(source_id: UUID):
    """Sync complet en background — ne bloque pas la réponse HTTP."""
    logger.info(f"[Sync BG] Démarrage background sync — source {source_id}")
    try:
        result = await sync_metadata(source_id)
        entity_count = result.get("entity_count", 0)
        logger.info(f"[Sync BG] sync_metadata terminé — {entity_count} entités, success={result.get('success')}")
        if result.get("success") and entity_count > 0:
            # ── Étape 2 : Découverte des relations ───────────────────────────
            try:
                discovery = await discover_relationships(source_id)
                logger.info(f"[Sync BG] Relations découvertes: {discovery.get('relations_found', 0)}")
            except Exception as e:
                logger.warning(f"[Sync BG] Relationship discovery failed (non-bloquant): {e}")

            # ── Étape 3 : Enrichissement sémantique BERT ─────────────────────
            try:
                from .semantic_enricher import enrich_source as _enrich_source
                source = await get_source(source_id)
                if source:
                    enrich_result = await _enrich_source(
                        source_id,
                        source.name,
                        source.connector_type.value,
                    )
                    logger.info(
                        f"[Sync BG] Enrichissement sémantique: "
                        f"{enrich_result.get('enriched', 0)} entités | "
                        f"MeiliSearch: {enrich_result.get('meili_indexed', 0)} | "
                        f"Domaines: {enrich_result.get('domain_stats', {})}"
                    )
            except Exception as e:
                logger.warning(f"[Sync BG] Enrichissement sémantique failed (non-bloquant): {e}")

        logger.info(f"[Sync BG] ✅ Pipeline complet — source {source_id}: {entity_count} entités")
    except Exception as e:
        logger.error(f"[Sync BG] ❌ Erreur — source {source_id}: {e}", exc_info=True)


@app.post("/sources/{source_id}/sync", status_code=202, tags=["Connections"])
async def sync_source_metadata(source_id: UUID, background_tasks: BackgroundTasks):
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    background_tasks.add_task(_run_sync_background, source_id)
    logger.info(f"[Sync] Tâche background enregistrée pour source {source_id}")

    return {
        "source_id":      str(source_id),
        "success":        True,
        "message":        "Sync lancé en arrière-plan — rafraîchis dans 2-3 minutes",
        "entity_count":   0,
        "field_count":    0,
        "relation_count": 0,
        "duration_ms":    0,
    }


@app.get("/sources/{source_id}/entities", tags=["Metadata"])
async def get_source_entities(
    source_id: UUID,
    page:      int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=200),
    search:    str = Query(""),
):
    source = await get_source_with_entities(
        source_id, page=page, page_size=page_size, search=search
    )
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    pagination = source.__dict__.get("_pagination", {})
    return {
        "source_id":  source_id,
        "entities":   source.entities,
        "pagination": pagination,
    }


@app.get("/sources/{source_id}/db-objects", tags=["Metadata"])
async def get_db_objects(
    source_id: UUID,
    object_type: Optional[str] = Query(None, description="filter: stored_procedure, function, trigger, sequence, index"),
):
    """
    Retourne les objets DB avancés d'une source :
    procédures stockées, fonctions, triggers, séquences, et index.
    Disponible pour MSSQL et PostgreSQL après sync.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT id, name, entity_type, description, metadata
            FROM source_entities
            WHERE source_id = $1
              AND entity_type IN (
                  'stored_procedure','function','trigger','sequence',
                  'aggregate','window'
              )
            ORDER BY entity_type, name
        """, source_id)

        # Index : stockés dans metadata de chaque table
        idx_rows = await conn.fetch("""
            SELECT name, indexes
            FROM source_entities
            WHERE source_id = $1
              AND entity_type IN ('table','view','materialized_view')
              AND indexes IS NOT NULL
              AND indexes != '[]'
        """, source_id)

    objects = []
    for r in rows:
        etype = r["entity_type"]
        if object_type and etype != object_type:
            continue
        meta = {}
        try:
            meta = json_module.loads(r["metadata"] or "{}")
        except Exception:
            pass
        objects.append({
            "id":          r["id"],
            "name":        r["name"],
            "type":        etype,
            "description": r["description"],
            "metadata":    meta,
        })

    indexes = []
    if not object_type or object_type == "index":
        for r in idx_rows:
            try:
                idx_list = json_module.loads(r["indexes"] or "[]")
                for idx in idx_list:
                    indexes.append({
                        "table":   r["name"],
                        "name":    idx.get("name"),
                        "type":    idx.get("type"),
                        "unique":  idx.get("unique", False),
                        "pk":      idx.get("pk", False),
                        "columns": idx.get("columns", ""),
                    })
            except Exception:
                pass

    counts = {}
    for o in objects:
        counts[o["type"]] = counts.get(o["type"], 0) + 1
    if indexes:
        counts["index"] = len(indexes)

    return {
        "source_id":  str(source_id),
        "counts":     counts,
        "objects":    objects,
        "indexes":    indexes,
        "total":      len(objects) + len(indexes),
    }


# ══════════════════════════════════════════════════════════════
# RELATIONS
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/discover", tags=["Relations"])
async def discover_source_relations(source_id: UUID):
    """Lance la détection automatique des relations pour une source."""
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    return await discover_relationships(source_id)


@app.get("/sources/{source_id}/relations/stats", tags=["Relations"])
async def get_relations_stats(source_id: UUID):
    """Statistiques de validation des relations."""
    return await get_validation_stats(source_id)


@app.get("/sources/{source_id}/relations", tags=["Relations"])
async def get_source_relations(source_id: UUID):
    """Retourne toutes les relations détectées pour une source."""
    try:
        source = await get_source(source_id)
        if not source:
            raise HTTPException(404, f"Source {source_id} introuvable")
        relations = await get_relations_for_source(source_id)
        return {
            "source_id": str(source_id),
            "total":     len(relations),
            "relations": relations,
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[Relations GET] {source_id}: {e}", exc_info=True)
        raise HTTPException(500, f"Erreur chargement relations: {str(e)}")


@app.post("/relations/{relation_id}/validate", tags=["Relations"])
async def validate_single_relation(relation_id: int, req: ValidateRelationRequest):
    """
    Valide ou rejette une relation détectée.
    Body: { is_confirmed: bool, validated_by?: str, reject_reason?: str }
    Enregistre aussi dans relation_feedback pour le blacklisting des faux positifs.
    """
    ok = await validate_relation(
        relation_id=relation_id,
        confirmed=req.is_confirmed,
        validated_by=req.validated_by or "expert",
        reject_reason=req.reject_reason,
    )
    if not ok:
        raise HTTPException(404, f"Relation {relation_id} introuvable")

    # Enregistrer dans relation_feedback pour tracking faux positifs récurrents
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rel = await conn.fetchrow(
                "SELECT source_id, source_entity, source_field, target_entity, target_field "
                "FROM entity_relations WHERE id=$1", relation_id)
            if rel:
                feedback_type = "confirmed" if req.is_confirmed else "rejected"
                await conn.execute("""
                    INSERT INTO relation_feedback
                        (source_id, relation_id, source_entity, source_field,
                         target_entity, target_field, feedback, user_id, comment)
                    VALUES ($1,$2,$3,$4,$5,$6,$7,$8,$9)
                """, rel["source_id"], relation_id,
                    rel["source_entity"], rel["source_field"],
                    rel["target_entity"], rel["target_field"],
                    feedback_type, req.validated_by or "expert",
                    req.reject_reason)
    except Exception as _fe:
        logger.warning(f"[validate] feedback insert warning: {_fe}")

    return {
        "success":      True,
        "relation_id":  relation_id,
        "is_confirmed": req.is_confirmed,
        "message":      "✅ Confirmée" if req.is_confirmed else "❌ Rejetée",
    }


@app.patch("/sources/{source_id}/relations/validate-bulk", tags=["Relations"])
@app.post("/sources/{source_id}/relations/validate-bulk", tags=["Relations"])
async def validate_bulk_relations(source_id: UUID, body: dict):
    """
    Valide plusieurs relations en une fois.
    Body: { validations: [{id, is_confirmed, reject_reason?}], validated_by?: str }
    """
    validations  = body.get("validations", [])
    validated_by = body.get("validated_by", "expert")
    results = []
    for v in validations:
        ok = await validate_relation(
            relation_id=v["id"],
            confirmed=bool(v["is_confirmed"]),
            validated_by=validated_by,
            reject_reason=v.get("reject_reason"),
        )
        results.append({"id": v["id"], "success": ok})
    return {"validated": len(results), "results": results}


@app.get("/sources/{source_id}/join-paths", tags=["Relations"])
async def find_join_paths(
    source_id:  UUID,
    from_table: str = Query(...),
    to_table:   str = Query(...),
    max_depth:  int = Query(3, ge=1, le=5),
):
    """Trouve les chemins de jointure entre deux tables."""
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    paths = await get_join_paths_for_source(source_id, from_table, to_table, max_depth)
    return {"from": from_table, "to": to_table, "paths": paths}


# ══════════════════════════════════════════════════════════════
# IMPORT FK DEPUIS SSMS
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/import-fk", tags=["Relations"])
async def import_fk_from_ssms(source_id: UUID, req: ImportFkRequest):
    """
    Importe des FK réelles exportées depuis SSMS.
    Chaque relation est insérée avec confidence=1.0, method=explicit_fk, is_confirmed=TRUE.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        imported = updated = skipped = 0

        for rel in req.relations:
            src_exists = await conn.fetchval(
                "SELECT 1 FROM source_entities WHERE source_id=$1 AND name=$2",
                source_id, rel.source_entity,
            )
            tgt_exists = await conn.fetchval(
                "SELECT 1 FROM source_entities WHERE source_id=$1 AND name=$2",
                source_id, rel.target_entity,
            )
            if not src_exists or not tgt_exists:
                skipped += 1
                continue

            result = await conn.execute(
                """
                INSERT INTO entity_relations
                    (source_id, source_entity, source_field,
                     target_entity, target_field,
                     relation_type, confidence, detection_method, is_confirmed)
                VALUES ($1, $2, $3, $4, $5, 'many_to_one', 1.0, 'explicit_fk', TRUE)
                ON CONFLICT (source_id, source_entity, source_field, target_entity)
                DO UPDATE SET
                    target_field     = EXCLUDED.target_field,
                    confidence       = 1.0,
                    detection_method = 'explicit_fk',
                    is_confirmed     = TRUE
                """,
                source_id,
                rel.source_entity, rel.source_field,
                rel.target_entity, rel.target_field,
            )
            if "INSERT 0 1" in result:
                imported += 1
            else:
                updated += 1

    return {
        "success":  True,
        "imported": imported,
        "updated":  updated,
        "skipped":  skipped,
        "message": (
            f"{imported} FK importées, {updated} mises à jour, "
            f"{skipped} ignorées (tables inconnues)"
        ),
    }




# ══════════════════════════════════════════════════════════════
# FILE UPLOAD
# ══════════════════════════════════════════════════════════════

@app.post("/upload", tags=["Files"])
async def upload_file(file: UploadFile = File(...)):
    """
    Upload un fichier et retourne un aperçu des données.
    Formats supportés : CSV, JSON, TXT, Excel (.xlsx/.xls), Parquet, Avro.
    Pour Excel multi-feuilles, retourne aussi sheets[] pour l'UI.
    """
    allowed = {".csv", ".json", ".txt", ".xlsx", ".xls", ".parquet", ".parq", ".avro"}
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in allowed:
        raise HTTPException(400, f"Type non supporté : {ext}. Acceptés : {', '.join(sorted(allowed))}")

    raw = await file.read()
    dest = os.path.join(UPLOAD_DIR, file.filename)
    with open(dest, "wb") as f:
        f.write(raw)

    try:
        # ── Excel (.xlsx / .xls) ─────────────────────────────────
        if ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb     = openpyxl.load_workbook(dest, read_only=True, data_only=True)
                sheets = wb.sheetnames
                # Aperçu de la première feuille
                ws      = wb[sheets[0]]
                rows_raw = list(ws.iter_rows(values_only=True))
                wb.close()
                headers  = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(rows_raw[0])] if rows_raw else []
                preview  = [
                    {headers[i]: row[i] for i in range(min(len(headers), len(row)))}
                    for row in rows_raw[1:6]
                ]
                return {
                    "filename":      file.filename,
                    "uploaded_path": dest,
                    "size":          len(raw),
                    "format":        "excel",
                    "sheets":        sheets,
                    "sheet_count":   len(sheets),
                    "columns":       headers,
                    "column_count":  len(headers),
                    "preview":       preview,
                    "message":       f"Excel uploadé — {len(sheets)} feuille(s) · {len(headers)} colonnes",
                }
            except ImportError:
                raise HTTPException(500, "openpyxl non installé : pip install openpyxl")

        # ── Parquet ──────────────────────────────────────────────
        elif ext in (".parquet", ".parq"):
            try:
                import pyarrow.parquet as pq
                pf       = pq.ParquetFile(dest)
                schema   = pf.schema_arrow
                meta     = pf.metadata
                columns  = [schema.field(i).name for i in range(len(schema))]
                col_types = {schema.field(i).name: str(schema.field(i).type) for i in range(len(schema))}
                # Aperçu des 5 premières lignes
                preview_tbl = pf.read_row_group(0).slice(0, 5).to_pydict()
                preview = [
                    {c: preview_tbl[c][j] for c in columns if j < len(preview_tbl.get(c, []))}
                    for j in range(min(5, meta.row_group(0).num_rows))
                ]
                return {
                    "filename":      file.filename,
                    "uploaded_path": dest,
                    "size":          len(raw),
                    "format":        "parquet",
                    "columns":       columns,
                    "column_count":  len(columns),
                    "col_types":     col_types,
                    "row_count":     meta.num_rows,
                    "row_groups":    meta.num_row_groups,
                    "preview":       preview,
                    "message":       f"Parquet uploadé — {meta.num_rows:,} lignes · {len(columns)} colonnes",
                }
            except ImportError:
                raise HTTPException(500, "pyarrow non installé : pip install pyarrow")

        # ── Avro ─────────────────────────────────────────────────
        elif ext == ".avro":
            try:
                import fastavro  # type: ignore[import]
                with open(dest, "rb") as avf:
                    reader  = fastavro.reader(avf)
                    schema  = reader.writer_schema or {}
                    records = []
                    for i, rec in enumerate(reader):
                        if i >= 5: break
                        records.append(rec)
                columns = list(records[0].keys()) if records else []
                if isinstance(schema, dict) and schema.get("fields"):
                    columns = [f["name"] for f in schema["fields"]]
                return {
                    "filename":      file.filename,
                    "uploaded_path": dest,
                    "size":          len(raw),
                    "format":        "avro",
                    "columns":       columns,
                    "column_count":  len(columns),
                    "schema_name":   schema.get("name", "") if isinstance(schema, dict) else "",
                    "preview":       records,
                    "message":       f"Avro uploadé — {len(columns)} champs · schéma intégré",
                }
            except ImportError:
                raise HTTPException(500, "fastavro non installé : pip install fastavro")

        # ── CSV / TXT ────────────────────────────────────────────
        elif ext in (".csv", ".txt"):
            try:
                text = raw.decode("utf-8")
            except UnicodeDecodeError:
                text = raw.decode("latin-1")
            reader  = csv_module.DictReader(io.StringIO(text))
            rows    = [dict(row) for i, row in enumerate(reader) if i < 5]
            columns = list(rows[0].keys()) if rows else []
            return {
                "filename":      file.filename,
                "uploaded_path": dest,
                "size":          len(raw),
                "format":        "csv",
                "columns":       columns,
                "column_count":  len(columns),
                "preview":       rows,
                "message":       f"CSV uploadé — {len(columns)} colonnes détectées",
            }

        # ── JSON ─────────────────────────────────────────────────
        else:
            try:
                text = raw.decode("utf-8")
            except UnicodeDecodeError:
                text = raw.decode("latin-1")
            data = json_module.loads(text)
            if isinstance(data, list):
                preview = data[:5]
                columns = list(preview[0].keys()) if preview else []
            else:
                preview = [data]
                columns = list(data.keys())
            return {
                "filename":      file.filename,
                "uploaded_path": dest,
                "size":          len(raw),
                "format":        "json",
                "columns":       columns,
                "column_count":  len(columns),
                "preview":       preview,
                "message":       f"JSON uploadé — {len(columns)} colonnes détectées",
            }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[Upload] {file.filename}: {e}", exc_info=True)
        raise HTTPException(500, f"Erreur traitement fichier : {str(e)}")


@app.get("/uploads", tags=["Files"])
async def list_uploads():
    """Liste les fichiers uploadés disponibles."""
    if not os.path.exists(UPLOAD_DIR):
        return {"files": []}
    files = []
    for fname in os.listdir(UPLOAD_DIR):
        fpath = os.path.join(UPLOAD_DIR, fname)
        if os.path.isfile(fpath):
            files.append({"filename": fname, "path": fpath, "size": os.path.getsize(fpath)})
    return {"files": files}




class ImportSchemaRequest(BaseModel):
    mode:          str             # "sql", "csv", "entities", "odata", "json", "graphql", "wsdl"
    schema_sql:    Optional[str]   = None   # DDL SQL (CREATE TABLE ...)
    schema_csv:    Optional[List[dict]] = None  # [{table_name, column_name, data_type, is_pk, is_fk}]
    entities:      Optional[List[str]]  = None  # noms de tables/entités (SAP, Dynamics, SAGE X3)
    odata_url:     Optional[str]   = None
    schema_json:   Optional[str]   = None
    schema_graphql:Optional[str]   = None
    wsdl_url:      Optional[str]   = None
    wsdl_xml:      Optional[str]   = None


@app.post("/sources/{source_id}/import-schema", tags=["Metadata"])
async def import_schema(source_id: UUID, req: ImportSchemaRequest):
    """
    Importe un schéma décrivant la structure d'une source sans connexion directe.
    Modes : sql (DDL), csv (table exportée), entities (noms SAP/Dynamics/SAGE),
            odata (URL $metadata), json (exemple réponse), graphql (SDL), wsdl.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    entities = []

    # ── Mode SQL DDL ─────────────────────────────────────────────────
    if req.mode == "sql" and req.schema_sql:
        entities = _parse_ddl_to_entities(req.schema_sql)

    # ── Mode CSV ─────────────────────────────────────────────────────
    elif req.mode == "csv" and req.schema_csv:
        from collections import defaultdict
        tables: dict = defaultdict(list)
        for row in req.schema_csv:
            tbl = row.get("table_name") or row.get("TABLE_NAME") or "unknown"
            tables[tbl].append({
                "name":        row.get("column_name") or row.get("COLUMN_NAME") or "",
                "type":        row.get("data_type")   or row.get("DATA_TYPE")   or "string",
                "native_type": row.get("data_type")   or "csv_import",
                "nullable":    str(row.get("is_nullable", "YES")).upper() != "NO",
                "primary_key": str(row.get("is_pk", "0")) in ("1", "YES", "True", "PRI"),
                "foreign_key": str(row.get("is_fk", "0")) in ("1", "YES", "True", "MUL"),
            })
        for tbl_name, fields in tables.items():
            entities.append({"name": tbl_name, "entity_type": "table", "fields": fields})

    # ── Mode Entités (SAP ABAP / Dynamics / SAGE X3) ─────────────────
    elif req.mode == "entities" and req.entities:
        for ent_name in req.entities:
            entities.append({
                "name":        ent_name.strip(),
                "entity_type": "entity",
                "fields":      [{"name": "id", "type": "string", "native_type": "auto",
                                 "nullable": False, "primary_key": True, "foreign_key": False}],
            })

    # ── Mode OData URL ────────────────────────────────────────────────
    elif req.mode == "odata" and req.odata_url:
        try:
            from .connection_service import _parse_odata_metadata
            entities = await _parse_odata_metadata(req.odata_url, {})
        except Exception as e:
            raise HTTPException(400, f"Erreur parsing OData $metadata : {e}")

    # ── Mode JSON exemple ─────────────────────────────────────────────
    elif req.mode == "json" and req.schema_json:
        import json as _json
        try:
            data = _json.loads(req.schema_json)
            from .connection_service import _parse_json
            name = source.name.replace(" ", "_").lower()
            entities = _parse_json(req.schema_json, name)
        except Exception as e:
            raise HTTPException(400, f"Erreur parsing JSON : {e}")

    else:
        raise HTTPException(400, f"Mode '{req.mode}' invalide ou données manquantes")

    if not entities:
        return {"success": False, "message": "Aucune entité extraite — vérifiez le contenu fourni", "entity_count": 0}

    from .repository import save_metadata
    entity_count = await save_metadata(source_id, entities)
    field_count  = sum(len(e.get("fields", [])) for e in entities)

    return {
        "success":      True,
        "entity_count": entity_count,
        "field_count":  field_count,
        "message":      f"{entity_count} entité(s) importée(s) · {field_count} champs",
        "mode":         req.mode,
    }


def _parse_ddl_to_entities(ddl: str) -> List[dict]:
    """Parse du DDL SQL (CREATE TABLE) et extrait les entités et leurs champs."""
    import re
    entities = []

    # Nettoyer les commentaires
    ddl_clean = re.sub(r'--[^\n]*', '', ddl)
    ddl_clean = re.sub(r'/\*.*?\*/', '', ddl_clean, flags=re.DOTALL)

    # Trouver les CREATE TABLE
    pattern = re.compile(
        r'CREATE\s+TABLE\s+(?:IF\s+NOT\s+EXISTS\s+)?'
        r'(?:\[?[\w\.]+\]?\.)?\[?([\w]+)\]?'
        r'\s*\(([^;]+?)\)',
        re.IGNORECASE | re.DOTALL
    )

    TYPE_MAP = {
        "int": "integer", "integer": "integer", "bigint": "integer",
        "smallint": "integer", "tinyint": "integer", "serial": "integer",
        "decimal": "decimal", "numeric": "decimal", "money": "decimal",
        "float": "float", "real": "float", "double": "float",
        "varchar": "string", "nvarchar": "string", "char": "string",
        "nchar": "string", "text": "string", "ntext": "string",
        "date": "date", "datetime": "datetime", "timestamp": "datetime",
        "datetime2": "datetime", "time": "string",
        "bit": "boolean", "bool": "boolean", "boolean": "boolean",
        "uuid": "uuid", "uniqueidentifier": "uuid",
        "json": "json", "jsonb": "json", "xml": "string",
    }

    for match in pattern.finditer(ddl_clean):
        table_name = match.group(1)
        body       = match.group(2)
        fields     = []
        pks        = set()

        # Chercher PRIMARY KEY inline ou en contrainte
        pk_constraint = re.search(r'PRIMARY\s+KEY\s*\(([^)]+)\)', body, re.IGNORECASE)
        if pk_constraint:
            for col in pk_constraint.group(1).split(","):
                pks.add(col.strip().strip("[]`\"'"))

        fk_refs = {}
        for fk_match in re.finditer(
            r'FOREIGN\s+KEY\s*\(([^)]+)\)\s*REFERENCES\s+[\[\`]?[\w\.]+[\]\`]?\s*\(([^)]+)\)',
            body, re.IGNORECASE
        ):
            for col in fk_match.group(1).split(","):
                fk_refs[col.strip().strip("[]`\"'")] = True

        # Parser chaque ligne de colonne
        col_pat = re.compile(
            r'^\s*\[?([\w]+)\]?\s+([\w]+)(?:\([^)]*\))?'
            r'(.*?)(?:,|$)',
            re.IGNORECASE
        )
        for line in body.split("\n"):
            line = line.strip().rstrip(",")
            if not line or re.match(r'(PRIMARY|FOREIGN|UNIQUE|CHECK|INDEX|KEY|CONSTRAINT)\b', line, re.IGNORECASE):
                continue
            m = col_pat.match(line)
            if not m:
                continue
            col_name  = m.group(1)
            col_type  = m.group(2).lower()
            col_rest  = m.group(3).upper()
            is_pk     = col_name in pks or "PRIMARY KEY" in col_rest
            is_fk     = col_name in fk_refs
            nullable  = "NOT NULL" not in col_rest and not is_pk

            if is_pk:
                pks.add(col_name)

            fields.append({
                "name":        col_name,
                "type":        TYPE_MAP.get(col_type, "string"),
                "native_type": col_type,
                "nullable":    nullable,
                "primary_key": is_pk,
                "foreign_key": is_fk,
            })

        if fields:
            entities.append({
                "name":        table_name,
                "entity_type": "table",
                "description": f"Importé depuis DDL SQL",
                "fields":      fields,
            })

    return entities

# ══════════════════════════════════════════════════════════════
# PROFILING
# ══════════════════════════════════════════════════════════════

@app.get("/sources/{source_id}/profile", tags=["Profiling"])
async def get_source_profile(source_id: UUID):
    """Résumé de profiling pour toutes les tables déjà profilées."""
    try:
        from .data_profiler import profile_source_summary
        return await profile_source_summary(source_id)
    except Exception as e:
        logger.error(f"[Profile] {source_id}: {e}", exc_info=True)
        raise HTTPException(500, str(e))


@app.get("/sources/{source_id}/profile/{table_name}", tags=["Profiling"])
async def get_table_profile(source_id: UUID, table_name: str, refresh: bool = False):
    """
    Profile une table spécifique (avec cache TTL intelligent).
    TTL : référentiel=24h | transactionnel=1h | inconnu=6h
    Passe refresh=true pour forcer le recalcul immédiat.
    """
    try:
        from .data_profiler import get_cached_profile, profile_entity
        if not refresh:
            cached = await get_cached_profile(source_id, table_name)
            if cached:
                return cached
        return await profile_entity(source_id, table_name)
    except Exception as e:
        logger.error(f"[Profile] {source_id}/{table_name}: {e}", exc_info=True)
        raise HTTPException(500, str(e))


@app.get("/sources/{source_id}/profile/cache/stats", tags=["Profiling"])
async def get_profile_cache_stats(source_id: UUID):
    """
    Retourne les statistiques du cache de profiling pour une source.
    Indique quelles tables sont fraîches, expirées ou jamais profilées.
    """
    from .database import get_pg_pool
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT entity_name,
                   profiled_at,
                   profile_data->>'table_class' AS table_class,
                   EXTRACT(EPOCH FROM (NOW() - profiled_at))/3600 AS age_hours
            FROM entity_profiles
            WHERE source_id = $1
            ORDER BY profiled_at DESC
        """, source_id)

    TTL = {"reference": 24, "transactional": 1}
    stats = {"fresh": 0, "stale": 0, "total": len(rows), "tables": []}

    for row in rows:
        tc      = row["table_class"] or "unknown"
        ttl     = TTL.get(tc, 6)
        age     = round(float(row["age_hours"]), 1)
        is_fresh = age <= ttl
        if is_fresh:
            stats["fresh"] += 1
        else:
            stats["stale"] += 1
        stats["tables"].append({
            "entity":      row["entity_name"],
            "table_class": tc,
            "age_hours":   age,
            "ttl_hours":   ttl,
            "fresh":       is_fresh,
            "profiled_at": row["profiled_at"].isoformat(),
        })

    stats["fresh_pct"] = round(stats["fresh"] / stats["total"] * 100, 1) if stats["total"] else 0
    return stats


@app.post("/sources/{source_id}/profile/batch", tags=["Profiling"])
async def batch_profile_tables(source_id: UUID, body: dict):
    """Profile plusieurs tables en parallèle (max 10 par appel)."""
    from .data_profiler import profile_entity
    tables = (body.get("tables") or [])[:10]
    if not tables:
        raise HTTPException(400, "tables[] requis")
    results = await asyncio.gather(
        *[profile_entity(source_id, t) for t in tables],
        return_exceptions=True,
    )
    return {
        "results": [
            r if not isinstance(r, Exception) else {"error": str(r), "table": tables[i]}
            for i, r in enumerate(results)
        ]
    }


@app.post("/sources/{source_id}/profile/all", tags=["Profiling"])
async def profile_all_entities(
    source_id:   UUID,
    sample_size: int = Query(1000, ge=100, le=10000, description="Lignes à échantillonner par table"),
    batch_size:  int = Query(5,    ge=1,   le=20,    description="Tables en parallèle par micro-batch"),
):
    """
    Profile TOUTES les entités de la source en séquence par micro-batches.
    - Reprend automatiquement là où il s'est arrêté (skip entités déjà profilées)
    - Générique : fonctionne pour SQL Server, MySQL, PostgreSQL, OData, CSV...
    - Retourne un résumé : total, success, errors
    """
    from .data_profiler import profile_source_all
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    result = await profile_source_all(
        source_id,
        sample_size=sample_size,
        batch_size=batch_size,
    )
    return {"source_id": str(source_id), **result}




@app.post("/sources/{source_id}/circuit-breaker/reset", tags=["Connections"])
async def reset_circuit_breaker(source_id: UUID):
    """
    Réinitialise manuellement le circuit breaker pour une source.
    Utile après avoir résolu un problème de connexion.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    from .connection_service import CircuitBreaker
    cb = CircuitBreaker.get(str(source_id))
    cb.reset()
    return {
        "source_id": str(source_id),
        "success":   True,
        "message":   "🟢 Circuit Breaker réinitialisé — état : CLOSED",
        "state":     "closed",
    }


@app.get("/sources/{source_id}/circuit-breaker/status", tags=["Connections"])
async def get_circuit_breaker_status(source_id: UUID):
    """
    Retourne l'état actuel du circuit breaker pour une source.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    from .connection_service import CircuitBreaker
    cb = CircuitBreaker.get(str(source_id), source.options or {})
    import time
    seconds_until_reset = max(0, int(cb.reset_timeout - (time.time() - cb._last_failure_at))) if cb.state.value == "open" else 0
    return {
        "source_id":           str(source_id),
        "state":               cb.state.value,
        "state_label":         cb.state_label,
        "failure_count":       cb._failure_count,
        "failure_threshold":   cb.failure_threshold,
        "reset_timeout_s":     cb.reset_timeout,
        "seconds_until_reset": seconds_until_reset,
        "enabled":             cb.enabled,
    }

# ══════════════════════════════════════════════════════════════
# ML — TRAIN / PREDICT / STATUS
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/ml/train", tags=["ML"])
async def ml_train(source_id: UUID):
    """
    Entraîne un modèle ML pour détecter les relations de la source.
    Nécessite ml_detector.py OU utiliser ml_relation_detector_v6.ipynb directement.
    """
    try:
        from .ml_detector import train_ml_model  # type: ignore[import]
    except ImportError:
        raise HTTPException(503, (
            "ml_detector.py non disponible. "
            "Utiliser ml_relation_detector_v6.ipynb pour générer best_model_xgboost.pkl, "
            "puis POST /discover utilisera le modèle automatiquement."
        ))
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    try:
        pool = await get_pg_pool()
        result = await train_ml_model(source_id, pool)

        # ── Auto-refresh : supprime les anciennes prédictions ML + redécouvre ──
        try:
            async with pool.acquire() as conn:
                deleted = await conn.fetchval(
                    "DELETE FROM entity_relations WHERE source_id=$1 AND detection_method='ml_predicted' RETURNING COUNT(*)",
                    source_id
                )
            logger.info(f"[ML Train] {deleted or 0} anciennes prédictions ML supprimées")
            # Relance la découverte avec le nouveau modèle
            discovery = await discover_relationships(source_id)
            result["auto_rediscovery"] = {
                "deleted_old": deleted or 0,
                "new_relations": discovery.get("relations_found", 0),
            }
            logger.info(f"[ML Train] Auto-rediscovery: {discovery.get('relations_found', 0)} nouvelles relations")
        except Exception as _e:
            logger.warning(f"[ML Train] Auto-rediscovery error (non-bloquant): {_e}")

        return _safe_json({"source_id": str(source_id), **result})
    except ValueError as e:
        raise HTTPException(400, str(e))
    except FileNotFoundError as e:
        raise HTTPException(404, str(e))
    except Exception as e:
        msg = str(e)
        logger.error(f"[ML Train] {source_id}: {msg}", exc_info=True)
        # Message lisible pour les cas fréquents
        if "insufficient_data" in msg or "assez de relations" in msg.lower():
            raise HTTPException(400, {
                "status": "insufficient_data",
                "message": "Pas assez de relations pour entraîner le modèle. Lancez d'abord ↻ Relancer pour détecter les relations explicites."
            })
        raise HTTPException(500, msg)


@app.post("/sources/{source_id}/ml/predict", tags=["ML"])
async def ml_predict(source_id: UUID):
    """
    Prédit les nouvelles relations ML. Alternative : POST /discover charge le .pkl auto.
    """
    try:
        from .ml_detector import predict_ml_relations  # type: ignore[import]
    except ImportError:
        raise HTTPException(503, (
            "ml_detector.py non disponible. "
            "Utiliser POST /discover qui charge best_model_xgboost.pkl automatiquement."
        ))
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    try:
        pool = await get_pg_pool()
        result = await predict_ml_relations(source_id, pool)
        return _safe_json({"source_id": str(source_id), **result})
    except FileNotFoundError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        logger.error(f"[ML Predict] {source_id}: {e}", exc_info=True)
        raise HTTPException(500, str(e))


@app.get("/sources/{source_id}/ml/status", tags=["ML"])
async def ml_status(source_id: UUID):
    """
    Retourne le statut du modèle ML pour la source :
    - available: bool
    - model_name, trained_at, threshold, metrics (F1, AUC...)
    """
    from .ml_detector import get_ml_status
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    return _safe_json({"source_id": str(source_id), **get_ml_status(source_id)})


@app.get("/sources/{source_id}/relations/false-positives", tags=["Relations"])
async def list_false_positives(
    source_id: UUID,
    min_rejections: int = Query(3, ge=1, le=20)
):
    """Retourne les paires rejetées >= min_rejections fois (blacklistées)."""
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    pairs = await get_false_positive_pairs(source_id, min_rejections)
    return {"source_id": str(source_id), "threshold": min_rejections, "total": len(pairs), "pairs": pairs}


@app.delete("/sources/{source_id}/relations/false-positives", tags=["Relations"])
async def rehabilitate_pair(
    source_id:     UUID,
    source_entity: str = Query(...),
    source_field:  str = Query(...),
    target_entity: str = Query(...),
):
    """Réhabilite une paire blacklistée."""
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    return await remove_false_positive(source_id, source_entity, source_field, target_entity)


@app.post("/sources/{source_id}/ml/retrain", tags=["ML"])
async def ml_retrain_feedback(source_id: UUID):
    """
    PATCH P12 — Re-entraîne le modèle ML avec les feedbacks experts accumulés.
    Déclenché automatiquement tous les 20 feedbacks, ou manuellement ici.
    """
    try:
        from .ml_detector import retrain_with_feedback  # type: ignore
    except ImportError:
        raise HTTPException(503, "ml_detector.py non disponible.")
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    try:
        pool   = await get_pg_pool()
        result = await retrain_with_feedback(source_id, pool)
        return {"source_id": str(source_id), **result}
    except Exception as e:
        logger.error(f"[ML Retrain] {source_id}: {e}", exc_info=True)
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════════════════════════════
# PATCH P5/P6 — Alternatives + Feedback expert
# ══════════════════════════════════════════════════════════════════════

@app.get("/sources/{source_id}/relations/alternatives", tags=["Relations"])
async def get_relation_alternatives(
    source_id:      UUID,
    source_entity:  str = Query(...),
    source_field:   str = Query(...),
    current_target: str = Query(...),
    top_k:          int = Query(3, ge=1, le=10),
):
    """
    PATCH P5 — Retourne les K meilleures relations alternatives pour un champ donné.
    Utilisé par l'UI de validation pour proposer des corrections.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    return await get_alternative_suggestions(
        source_id, source_entity, source_field, current_target, top_k
    )


@app.post("/sources/{source_id}/relations/{relation_id}/feedback", tags=["Relations"])
async def submit_relation_feedback(
    source_id:   UUID,
    relation_id: UUID,
    req:         FeedbackRequest,
):
    """
    PATCH P6 — Enregistre le feedback expert (confirmé/rejeté/alternative).
    Si 20 feedbacks accumulés → retourne should_retrain=True.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    result = await record_expert_feedback(
        source_id=source_id,
        relation_id=relation_id,
        source_entity=req.source_entity,
        source_field=req.source_field,
        target_entity=req.target_entity,
        target_field=req.target_field,
        feedback=req.feedback,
        user_id=req.user_id,
        comment=req.comment,
        alternative_target_entity=req.alternative_target_entity,
        alternative_target_field=req.alternative_target_field,
    )
    # Auto-trigger re-training si seuil atteint
    if result.get("should_retrain"):
        try:
            from .ml_detector import retrain_with_feedback  # type: ignore
            pool = await get_pg_pool()
            asyncio.create_task(retrain_with_feedback(source_id, pool))
            logger.info(f"[Feedback] Re-training ML déclenché en background — {source_id}")
        except Exception as _rt:
            logger.warning(f"[Feedback] Auto-retrain failed (non-bloquant): {_rt}")
    return result


# ══════════════════════════════════════════════════════════════════════
# PATCH P7/P8 — Histogrammes DB + Dépendances objets
# ══════════════════════════════════════════════════════════════════════

@app.get("/sources/{source_id}/entities/{entity_name}/columns/{column_name}/histogram",
         tags=["Profiling"])
async def get_column_histogram_endpoint(
    source_id:   UUID,
    entity_name: str,
    column_name: str,
):
    """
    PATCH P7 — Retourne l'histogramme d'une colonne depuis les stats moteur DB.
    PostgreSQL : pg_stats | MSSQL : sys.dm_db_stats_histogram
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    from .data_profiler import get_column_histogram
    return await get_column_histogram(source_id, entity_name, column_name)


@app.get("/sources/{source_id}/dependencies", tags=["Metadata"])
async def get_source_dependencies(source_id: UUID):
    """
    PATCH P8 — Retourne les dépendances entre objets DB (vues→tables, procédures, triggers).
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    from .data_profiler import get_object_dependencies
    return await get_object_dependencies(source_id)


@app.get("/sources/{source_id}/excel-metadata/{entity_name}", tags=["Metadata"])
async def get_excel_metadata_endpoint(source_id: UUID, entity_name: str):
    """
    PATCH P9 — Retourne les métadonnées avancées d'un fichier Excel
    (formules, auteur, headers/footers, plages nommées).
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")
    # Récupérer le path du fichier depuis les options de la source
    opts      = source.options or {}
    file_path = opts.get("file_path") or source.base_url or ""
    if not file_path:
        raise HTTPException(400, "Chemin de fichier non configuré pour cette source")
    if not file_path.lower().endswith((".xlsx", ".xlsm", ".xls")):
        raise HTTPException(400, "Cette source n'est pas un fichier Excel")
    from .data_profiler import extract_excel_metadata
    return extract_excel_metadata(file_path)


# ══════════════════════════════════════════════════════════════════════
# PATCH CROSS-SOURCES — Détection + Path-finding multi-sources
# ══════════════════════════════════════════════════════════════════════

@app.post("/cross-source/detect", tags=["Cross-Source"])
async def detect_cross_source(body: dict):
    """
    Phase 3 — Détecte les mappings entre colonnes de deux sources différentes.
    Body: { source_id_a: UUID, source_id_b: UUID, min_overlap: float = 0.30 }
    """
    try:
        from .cross_source_mapper import (
            detect_cross_source_mappings, save_cross_source_mappings
        )
    except ImportError:
        raise HTTPException(503, "cross_source_mapper.py non disponible")

    sid_a = UUID(body.get("source_id_a", ""))
    sid_b = UUID(body.get("source_id_b", ""))
    min_overlap = float(body.get("min_overlap", 0.30))

    src_a = await get_source(sid_a)
    src_b = await get_source(sid_b)
    if not src_a:
        raise HTTPException(404, f"Source A {sid_a} introuvable")
    if not src_b:
        raise HTTPException(404, f"Source B {sid_b} introuvable")

    mappings = await detect_cross_source_mappings(sid_a, sid_b, min_overlap)
    saved    = await save_cross_source_mappings(mappings)
    return {
        "source_id_a": str(sid_a),
        "source_id_b": str(sid_b),
        "mappings_found": len(mappings),
        "mappings_saved": saved,
        "top_mappings":   mappings[:10],
    }


@app.post("/cross-source/path", tags=["Cross-Source"])
async def find_cross_source_path(body: dict):
    """
    Phase 3 — Trouve les chemins de jointure entre entités de sources différentes.
    Body: {
      source_ids: [UUID, ...],
      from: { source_id: UUID, entity: str },
      to:   { source_id: UUID, entity: str },
      max_depth: int = 4
    }
    """
    try:
        from .cross_source_mapper import build_multi_source_graph
    except ImportError:
        raise HTTPException(503, "cross_source_mapper.py non disponible")

    source_ids = [UUID(s) for s in body.get("source_ids", [])]
    from_info  = body.get("from", {})
    to_info    = body.get("to",   {})
    max_depth  = int(body.get("max_depth", 4))

    if not source_ids or not from_info or not to_info:
        raise HTTPException(400, "source_ids, from et to requis")

    graph = await build_multi_source_graph(source_ids)
    paths = graph.find_paths(
        str(from_info.get("source_id", "")), from_info.get("entity", ""),
        str(to_info.get("source_id",   "")), to_info.get("entity",   ""),
        max_depth,
    )
    return {
        "from":       from_info,
        "to":         to_info,
        "paths_found": len(paths),
        "paths":       paths,
        "graph_stats": {"nodes": graph.node_count, "edges": graph.edge_count},
    }


@app.get("/cross-source/mappings", tags=["Cross-Source"])
async def list_cross_source_mappings(
    source_id_a:    UUID = Query(...),
    source_id_b:    Optional[UUID] = Query(None),
    min_confidence: float = Query(0.30, ge=0.0, le=1.0),
):
    """Phase 3 — Liste les mappings cross-sources existants."""
    try:
        from .cross_source_mapper import get_cross_source_mappings
    except ImportError:
        raise HTTPException(503, "cross_source_mapper.py non disponible")
    return await get_cross_source_mappings(source_id_a, source_id_b, min_confidence)


@app.put("/cross-source/mappings/{mapping_id}/validate", tags=["Cross-Source"])
async def validate_cross_source_mapping(mapping_id: UUID, body: dict):
    """
    Valide ou rejette un mapping cross-source.
    Body: { validated: bool, validated_by: str }
    """
    validated    = bool(body.get("validated", True))
    validated_by = body.get("validated_by", "user")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        result = await conn.execute("""
            UPDATE cross_source_mappings
            SET is_validated = $1,
                validated_by = $2
            WHERE id = $3
        """, validated, validated_by, mapping_id)

    if result == "UPDATE 0":
        raise HTTPException(404, f"Mapping {mapping_id} introuvable")

    return {
        "success":      True,
        "mapping_id":   str(mapping_id),
        "validated":    validated,
        "validated_by": validated_by,
        "message":      f"Mapping {'validé' if validated else 'rejeté'} par {validated_by}",
    }


# ══════════════════════════════════════════════════════════════
# CHAT IA ENDPOINT §2.4
# ══════════════════════════════════════════════════════════════

class ChatRequest(BaseModel):
    source_id: str
    question: str
    history: Optional[List[dict]] = []

@app.post("/chat")
async def chat_endpoint(req: ChatRequest):
    """Endpoint Chat IA — fallback local en attendant Ollama (Phase 2)."""
    try:
        answer = await _build_chat_answer(req.source_id, req.question)
        return {"answer": answer}
    except Exception as e:
        logger.error(f"Chat error: {e}")
        return {"answer": f"Erreur: {str(e)}"}


# ══════════════════════════════════════════════════════════════
# CONVERSATIONS CHAT IA
# ══════════════════════════════════════════════════════════════

class Message(BaseModel):
    role: str  # "user" ou "bot"
    content: str

class ConversationIn(BaseModel):
    id: str
    source_id: str
    title: str
    messages: List[Message]
    created_at: str

@app.post("/conversations")
async def save_conversation(conv: ConversationIn):
    """Sauvegarder ou mettre à jour une conversation (PAS d'ID UUID requis)"""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            # Vérifier si la conversation existe déjà (par son mongo-like ID string)
            existing_id = await conn.fetchval(
                "SELECT id FROM conversations WHERE id::text = $1",
                conv.id
            )
            
            messages_json = json_module.dumps([m.model_dump() for m in conv.messages])
            
            if existing_id:
                # Mise à jour
                await conn.execute(
                    """UPDATE conversations 
                       SET title=$2, messages=$3, updated_at=NOW()
                       WHERE id::text = $1""",
                    conv.id, conv.title, messages_json
                )
                return {"status": "ok", "id": str(existing_id)}
            else:
                # Insertion - générer UUID si nécessaire
                try:
                    conv_uuid = UUID(conv.id)
                except ValueError:
                    conv_uuid = uuid4()
                
                await conn.execute(
                    """INSERT INTO conversations (id, source_id, title, messages)
                       VALUES ($1::uuid, $2::uuid, $3, $4)
                       ON CONFLICT (id) DO UPDATE
                       SET title=$3, messages=$4, updated_at=NOW()""",
                    conv_uuid, 
                    conv.source_id, 
                    conv.title, 
                    messages_json
                )
                
                return {"status": "ok", "id": str(conv_uuid)}
    except Exception as e:
        logger.error(f"Conversation save error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/conversations")
async def get_conversations(limit: int = 200, offset: int = 0):
    """Récupérer toutes les conversations avec pagination"""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch(
                "SELECT id, source_id, title, created_at, updated_at FROM conversations ORDER BY COALESCE(updated_at, created_at) DESC LIMIT $1 OFFSET $2",
                min(limit, 500), offset
            )
            
            conversations = []
            for row in rows:
                conversations.append({
                    "id": str(row['id']),
                    "source_id": str(row['source_id']),
                    "title": row['title'],
                    "created_at": row['created_at'].isoformat() if row['created_at'] else None,
                    "updated_at": row['updated_at'].isoformat() if row['updated_at'] else None,
                })
            
            total = await conn.fetchval("SELECT COUNT(*) FROM conversations")
            return {"conversations": conversations, "total": total}
    except Exception as e:
        logger.error(f"Get conversations error: {e}")
        return {"conversations": []}



@app.get("/conversations/pinned", tags=["Conversations"])
async def get_pinned_conversations(user_id: str = "admin"):
    """Retourne la liste des conv_ids épinglées pour un utilisateur."""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT conv_id::text, pinned_at
                FROM pinned_conversations
                WHERE user_id = $1
                ORDER BY pinned_at DESC
            """, user_id)
        return {"pinned": [{"conv_id": str(r["conv_id"]), "pinned_at": r["pinned_at"].isoformat()} for r in rows]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/conversations/{conv_id}")
async def get_conversation(conv_id: str):
    """Récupérer une conversation spécifique"""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            # Essayer de parser en UUID, sinon utiliser comme string
            try:
                uuid_val = UUID(conv_id)
            except ValueError:
                uuid_val = None
            
            if uuid_val:
                row = await conn.fetchrow(
                    "SELECT id, source_id, title, messages FROM conversations WHERE id = $1 OR id::text = $2",
                    uuid_val, conv_id
                )
            else:
                row = await conn.fetchrow(
                    "SELECT id, source_id, title, messages FROM conversations WHERE id::text = $1",
                    conv_id
                )
            
            if not row:
                raise HTTPException(status_code=404, detail="Conversation not found")
            
            messages = json_module.loads(row['messages']) if row['messages'] else []
            
            return {
                "id": str(row['id']),
                "source_id": str(row['source_id']),
                "title": row['title'],
                "messages": messages
            }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Get conversation error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ══════════════════════════════════════════════════════════════
# CHAT & CHAT STREAMING ENDPOINTS
# ══════════════════════════════════════════════════════════════


async def _fix_slots(slots, question: str, schema: dict, known_entities: list,
                      source_id_str: str = "", pg_pool=None):
    """
    Corrige les intents + résout les entités via MeiliSearch/synonymes.
    Appelée dans _build_chat_answer ET /nlu/generate-sql.
    """
    import re as _re_fix
    q_low = question.lower()

    if not slots:
        return slots

    # Injections raw_text
    slots.raw_text = question

    # ── Résolution sémantique : mots métier → vraies tables via MeiliSearch ──
    # Bloqué pour SHOW_DASHBOARD (routing propre dans dashboard_engine)
    # Bloqué si la question sera gérée par AgentRAG (matching flou SXA_DIRECT_SQL)
    # Sinon _fix_slots résout "affiche-moi" → TH_SECFR (faux positif sémantique)
    _AGENT_PREFIXES = (
        "montre-moi", "affiche-moi", "donne-moi", "montrer les", "afficher les",
        "quelles sont les", "quels sont les", "liste les", "lister les",
        "taux de change", "cours de change", "cours marché", "forex",
        "intégration bancaire", "flux de trésorerie", "journal des flux",
        "solde bancaire", "si bancaire", "rapprochement bancaire",
        "groupe de sociétés", "liste les sociétés", "liste les banques",
    )
    _is_agent_question = any(q_low.startswith(p) or p in q_low for p in _AGENT_PREFIXES)

    if not slots.table_names and source_id_str and slots.intent != Intent.SHOW_DASHBOARD and not _is_agent_question:
        # ── Tables prioritaires SXA pour questions financières ────────────────
        _SXA_PRIORITY = {
            "solde": "Dernière integration bancaire",
            "trésorerie": "Dernière integration bancaire",
            "tresorerie": "Dernière integration bancaire",
            "bancaire": "Dernière integration bancaire",
            "banque": "Dernière integration bancaire",
            "compte": "Comptes",
            "comptes": "Comptes",
            "bancaires": "Dernière integration bancaire",
            "disponibles": "Comptes",
            "financement": "FINANCEMENT_BI",
            "financements": "FINANCEMENT_BI",
            "amortissements": "Tableaux d'amortissement",
            "amortissement": "Tableaux d'amortissement",
            "mouvement": "SI_Bancaire",
            "transaction": "SI_Bancaire",
            "virement": "SI_Bancaire",
        }
        q_words = q_low.replace("é","e").replace("è","e").replace("ê","e").split()
        for word in q_words:
            if word in _SXA_PRIORITY:
                tbl = _SXA_PRIORITY[word]
                if tbl in known_entities:
                    slots.table_names = [tbl]
                    logger.info(f"[_fix_slots] SXA priority: '{word}' → '{tbl}'")
                    break
        
        # Si toujours pas de table → MeiliSearch sémantique
        if not slots.table_names:
            try:
                from .semantic_enricher import semantic_search as _sem_search, SYNONYMS, _SYNONYM_INDEX
                concepts = []
                for word in q_low.split():
                    if word in _SYNONYM_INDEX:
                        concepts.append(_SYNONYM_INDEX[word])
                search_terms = set(concepts) if concepts else {q_low.split()[0]}
                for term in list(search_terms)[:3]:
                    results = await _sem_search(
                        query      = term,
                        source_ids = [source_id_str],
                        limit      = 3,
                        use_vector = False,
                    )
                    if results:
                        best = results[0]
                        if best["relevance"] > 0.3 and best["name"] in known_entities:
                            slots.table_names = [best["name"]]
                            logger.info(f"[_fix_slots] Résolution sémantique : '{term}' → '{best['name']}' (score={best['relevance']})")
                            break
            except Exception as e:
                logger.warning(f"[_fix_slots] Semantic search error: {e}")

    # 1. DDL → bloquer (géré ailleurs)

    # 2. HAVING : "ayant / plus de N fois / commandes"
    _having_kw = ["ayant", "having", "au moins", "au plus", "dont le total"]
    _filter_kw = ["prix", "price", "montant", "amount", "cout", "cost", "salaire"]
    _count_ctx  = ["fois", "commandes", "orders", "produits", "articles", "lignes", "enregistrements", "transactions"]
    _is_filter  = any(kw in q_low for kw in _filter_kw) and not any(kw in q_low for kw in _count_ctx)
    _is_having  = any(kw in q_low for kw in _having_kw) or (
        any(kw in q_low for kw in ["plus de", "moins de"]) and not _is_filter
    )
    if _is_having and not _is_filter and slots.intent != Intent.SHOW_DASHBOARD:
        if slots.intent != Intent.SHOW_DASHBOARD:
            slots.intent = Intent.GENERATE_AGG
        # Extrait valeur numérique
        _num = _re_fix.search(r"(?:plus de|moins de|au moins|au plus|ayant)\s+(\d+)", q_low)
        if _num:
            _op = "lt" if any(k in q_low for k in ["moins de", "inferieur", "inférieur"]) else "gt"
            slots.amount_filter = {"op": _op, "value": int(_num.group(1))}
        elif slots.amount_filter and slots.amount_filter.get("op") in ("eq", "="):
            slots.amount_filter["op"] = "gt"
        # Extrait le sujet "banques ayant..." → group_by = "banque"
        if not slots.group_by:
            _hav_m = _re_fix.match(r"(\w+)\s+(?:ayant|having|dont|with)", q_low)
            if _hav_m:
                import unicodedata as _ud3
                _subj = _hav_m.group(1)
                _subj_n = _ud3.normalize("NFD", _subj).encode("ascii","ignore").decode().lower()
                if _subj_n.endswith("es"): _subj_n = _subj_n[:-1]
                elif _subj_n.endswith("s"): _subj_n = _subj_n[:-1]
                if _subj_n not in ("le","la","les","l","un","une","des","du"):
                    slots.group_by = _subj_n
        # Sinon injecte group_by depuis le schéma (champ non-ID)
        if not slots.group_by and slots.table_names:
            tbl = slots.table_names[0]
            tbl_fields = schema.get(tbl, [])
            gby = next(
                (f for f in tbl_fields if any(k in f.lower() for k in ["name","nom","banque","societe","company","label"])
                 and "id" not in f.lower()),
                next((f for f in tbl_fields if "id" not in f.lower() and "code" not in f.lower()),
                     tbl_fields[0] if tbl_fields else None)
            )
            if gby:
                slots.group_by = gby

    # 3. Filtre prix/montant
    elif _is_filter and any(kw in q_low for kw in ["plus de", "superieur", "supérieur", "greater", "inferieur", "inférieur", "moins de"]):
        slots.intent = Intent.GENERATE_FILTER

    # 4. "montre / affiche + table" → SQL
    _show_kw = ["montre", "affiche", "donne", "show", "display"]
    _data_kw = ["premiers", "derniers", "chers", "recents", "first", "last", "expensive"]
    if (slots.intent == Intent.LIST_ENTITIES and slots.table_names and
            (any(kw in q_low for kw in _show_kw) or any(kw in q_low for kw in _data_kw))):
        slots.intent = Intent.GENERATE_SQL

    # 5. AGG FR+EN : "total X by Y", "nombre de X par Y", "sum of X by Y"
    _agg_trg = ["nombre de", "number of", "total", "sum of", "count of",
                "average of", "avg of", "total des", "total par",
                # ES
                "total de", "suma de", "promedio de", "cantidad de", "numero de",
                "cuantos", "ventas por", "importe por",
                # DE
                "gesamt", "summe", "anzahl", "durchschnitt", "umsatz nach"]
    _per_trg = ["par", "by", "per", "pour chaque", "for each", "grouped by",
                # ES
                "por", "para cada", "agrupado por",
                # DE
                "nach", "pro", "je", "gruppiert nach"]
    if (any(kw in q_low for kw in _agg_trg) and
            any(kw in q_low for kw in _per_trg) and
            slots.intent not in (Intent.GENERATE_AGG, Intent.SHOW_DASHBOARD)):
        if slots.intent != Intent.SHOW_DASHBOARD:
            slots.intent = Intent.GENERATE_AGG
        if not slots.metric:
            if any(kw in q_low for kw in ["count", "number of", "nombre de"]):
                slots.metric = "COUNT"
            elif any(kw in q_low for kw in ["average", "avg", "moyenne"]):
                slots.metric = "AVG"
            else:
                slots.metric = "SUM"
        # Extraire group_by depuis "by/par/por/nach/pro/per X" — générique toutes langues
        import re as _re_gb
        _by_m = _re_gb.search(
            r"(?:groupe\s+par|grouped\s+by|agrupado\s+por|gruppiert\s+nach"
            r"|by|par|por|nach|per|pro|je|pour\s+chaque|for\s+each)\s+(\w+)", q_low
        )
        if _by_m and not slots.group_by:
            slots.group_by = _by_m.group(1)
        # Si toujours pas de group_by, cherche le mot APRES "par/por/by/nach"
        if not slots.group_by:
            _sim = _re_gb.search(r"(?:par|by|por|nach)\s+(\w{3,})", q_low)
            if _sim:
                slots.group_by = _sim.group(1)

    # 6. LAG / mois précédent
    if any(kw in q_low for kw in ["mois precedent", "mois précédent", "par rapport au", "comparaison mois"]) and slots.intent != Intent.SHOW_DASHBOARD:
        if slots.intent != Intent.SHOW_DASHBOARD:
            slots.intent = Intent.GENERATE_AGG

    # 7. top N par groupe → window
    # top N X par/por/by/nach Y → AGG générique toutes langues
    _top_par = _re_fix.search(
        r"top\s+(\d+)\s+(\w+)\s+(?:par|by|por|nach|per|pro)\s+(\w+)", q_low
    )
    if _top_par:
        slots.top_n  = int(_top_par.group(1))
        if slots.intent != Intent.SHOW_DASHBOARD:
            if slots.intent != Intent.SHOW_DASHBOARD:
                slots.intent = Intent.GENERATE_AGG
        _entity_word = _top_par.group(2)  # ex: "clientes", "clients", "kunden"
        _metric_word = _top_par.group(3)  # ex: "importe", "montant", "betrag"
        # group_by = le mot ENTITE (avant par/por/by/nach)
        if not slots.group_by:
            slots.group_by = _entity_word
        # metric = SUM sur le mot METRIQUE (après par/por/by/nach)
        if not slots.metric:
            slots.metric = "SUM"
        # Stocker le champ métrique pour le SQL Generator
        if not getattr(slots, 'numeric_field', None):
            slots.numeric_field = _metric_word

    # 8. Mapping mots métier → table
    _entity_kw_map = {
        "fournisseur": ["Suppliers","SUPPLIER","BPSUPPLIER"],
        "supplier":    ["Suppliers","SUPPLIER"],
        "produit":     ["Products","PRODUCT","ITMMATER"],
        "product":     ["Products","PRODUCT"],
        "client":      ["Customers","CUSTOMER","BPCUSTOMER"],
        "customer":    ["Customers","CUSTOMER"],
        "commande":    ["Orders","ORDER","SORDERQ"],
        "order":       ["Orders","SORDER"],
        "employe":     ["Employees","EMPLOYEE"],
        "employee":    ["Employees","EMPLOYEE"],
        "categorie":   ["Categories","CATEGORY"],
        "category":    ["Categories","CATEGORY"],
    }
    if not slots.table_names:
        for kw, candidates in _entity_kw_map.items():
            if kw in q_low:
                for c in candidates:
                    matched = next((e for e in known_entities if e.lower() == c.lower()), None)
                    if matched:
                        slots.table_names = [matched]
                        break
                if slots.table_names:
                    break

    # 9. count_entities + table connue → AGG
    if slots.intent == Intent.COUNT_ENTITIES and slots.table_names and slots.table_names[0] in schema:
        if slots.intent != Intent.SHOW_DASHBOARD:
            slots.intent = Intent.GENERATE_AGG
        if not slots.metric:
            slots.metric = "COUNT"

    return slots

async def _build_chat_answer(source_id_str: str, question: str) -> str:
    """
    Pipeline complet : NLU → SQLGenerator → SQLValidator → réponse formatée.
    §2.3.3 — remplace l'ancien système if/elif par le vrai moteur.
    """
    import time
    t0 = time.time()

    # ── Validation source ─────────────────────────────────────────────
    try:
        source_id = UUID(source_id_str)
    except (ValueError, AttributeError):
        return "ID source invalide."

    source = await get_source(source_id)
    if not source:
        return "Source introuvable."

    name         = source.name
    entity_count = source.entity_count or 0

    # ── Récupère entités et relations (pour intents catalogue) ────────
    try:
        detail   = await get_source_with_entities(source_id, page=1, page_size=20)
        entities = detail.entities if detail else []
    except Exception:
        entities = []

    try:
        relations = await get_relations_for_source(source_id)
        relations = relations[:5] if relations else []
    except Exception:
        relations = []

    # ── Récupère le schéma complet pour le SQL Generator ─────────────
    schema: dict = {}
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT se.name AS table_name, ef.name AS field_name
                FROM source_entities se
                JOIN entity_fields ef ON ef.entity_id = se.id
                WHERE se.source_id = $1 AND se.is_visible = TRUE
                ORDER BY se.name, ef.position
                LIMIT 20000
            """, source_id)
        for r in rows:
            schema.setdefault(r["table_name"], []).append(r["field_name"])
    except Exception as e:
        logger.warning(f"[Chat] Schema fetch error: {e}")

    # Fix : filtrer les tables infra du schema ET de known_entities
    _INFRA_KE = ("QRTZ_","qrtz_","sys","SYS","dt_","DT_","MSreplication","sysdiagram","__")
    schema = {k: v for k, v in schema.items() if not any(k.startswith(p) for p in _INFRA_KE)}
    known_entities = list(schema.keys())

    # Ajoute aussi les entités sans champs (vues, etc.) dans known_entities
    # pour que directly_mentioned puisse les détecter
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            all_names = await conn.fetch("""
                SELECT name FROM source_entities
                WHERE source_id = $1 AND is_visible = TRUE
                  AND name NOT IN (SELECT DISTINCT se.name
                                   FROM source_entities se
                                   JOIN entity_fields ef ON ef.entity_id = se.id
                                   WHERE se.source_id = $1)
            """, source_id)
            for r in all_names:
                if r["name"] not in known_entities:
                    known_entities.append(r["name"])
    except Exception:
        pass

    # ── Détection directe des tables mentionnées dans la question ─────
    # Normalise : retire accents, espaces pour la comparaison
    import unicodedata
    def _norm(s):
        return unicodedata.normalize('NFD', s).encode('ascii', 'ignore').decode().upper().replace(' ', '_').replace('-', '_')

    q_norm = _norm(question)
    directly_mentioned = [t for t in known_entities if _norm(t) in q_norm]

    # Fuzzy match: "SOINVOICE" peut matcher "SOI_INVOICE" ou "SO_INVOICE" etc.
    if not directly_mentioned:
        q_words = q_norm.split()
        for word in q_words:
            if len(word) >= 4:  # Ignore mots courts
                for entity in known_entities:
                    e_norm = _norm(entity)
                    # Match si l'un contient l'autre (sans underscore)
                    if (word in e_norm.replace('_','') or
                        e_norm.replace('_','') in word or
                        word in e_norm):
                        if entity not in directly_mentioned:
                            directly_mentioned.append(entity)
                            logger.info(f"[Chat] Fuzzy match: '{word}' → '{entity}'")

    # ── Blocage DDL/DML avant tout (sécurité) ────────────────────────
    import re as _re
    _DDL = _re.compile(
        r"\b(DROP|INSERT|UPDATE|DELETE|TRUNCATE|ALTER|CREATE|EXEC|EXECUTE)\b",
        _re.IGNORECASE
    )
    if _DDL.search(question):
        return "**Requête refusée** — instruction interdite.\n\n*Seules les questions en langage naturel sont autorisées.*"

    # ── Blocage questions hors-périmètre ERP ─────────────────────────
    _OUT_OF_SCOPE = _re.compile(
        r"\b(mot de passe|password|météo|meteo|weather"
        r"|verrou.*système|accord.*accès|option.*accès"
        r"|accès.*système|login.*admin|admin.*login)\b",
        _re.IGNORECASE
    )
    if _OUT_OF_SCOPE.search(question):
        return (
            "❌ **Question hors périmètre** — OnePilot répond uniquement aux questions "
            f"sur les données de **{name}**.\n\n"
            "*Exemples : solde bancaire, transactions, financements, comptes...*"
        )

    # ── Détection réponse à une clarification en attente ─────────────
    # Si le contexte a une clarification pendante ET la question est courte
    # (≤ 5 mots) → c'est probablement une réponse à la clarification
    context_check = get_context(source_id_str)
    if context_check.pending_clarification and context_check.pending_slots:
        pending = context_check.pending_clarification
        opts_lower = [o.lower().strip() for o in pending.get("options", [])]
        q_stripped = question.strip().lower()
        # Vérifie si la réponse correspond à une option (exacte ou partielle)
        matched_option = None
        for opt in pending.get("options", []):
            if (opt.lower().strip() == q_stripped or
                q_stripped in opt.lower() or
                opt.lower() in q_stripped or
                # Match sur les premiers mots (ex: "Comptes" match "Comptes courants")
                any(word == q_stripped for word in opt.lower().split())):
                matched_option = opt
                break
        # Aussi vérifie si la question est une entité connue
        if not matched_option and q_stripped in [e.lower() for e in known_entities]:
            matched_option = next((e for e in known_entities if e.lower() == q_stripped), None)

        if matched_option:
            logger.info(f"[Chat] Clarification résolue: '{matched_option}' pour slot='{pending.get('slot_key')}'")
            # Applique la clarification aux slots en attente
            from .ambiguity_resolver import AmbiguityResolver as _AR
            _resolver = _AR()
            resolved_slots = _resolver.apply_clarification(
                context_check.pending_slots,
                pending.get("slot_key", "table_names"),
                matched_option,
            )
            # Réinitialise l'état de clarification
            context_check.pending_clarification = None
            context_check.pending_slots = None
            orig_question = context_check.pending_question or question
            context_check.pending_question = None
            # Continue avec les slots résolus et la question originale
            slots = resolved_slots
            question = orig_question  # Utilise la question originale pour le contexte
            goto_sql = True
        else:
            goto_sql = False
    else:
        goto_sql = False

    # ── Pipeline NLU ──────────────────────────────────────────────────
    if not goto_sql:
        try:
            nlu     = get_nlu_pipeline()
            context = get_context(source_id_str)
            slots   = nlu.process(question, context, known_entities)
        except Exception as e:
            logger.warning(f"[Chat] NLU error: {e}")
            slots = None
    # Si goto_sql=True → slots déjà définis par la clarification résolue

    # ── Log NLU dans nlu_query_log (immédiatement après NLU) ─────────
    if slots:
        try:
            _pool_log = await get_pg_pool()
            async with _pool_log.acquire() as _conn_log:
                await _conn_log.execute("""
                    INSERT INTO nlu_query_log
                        (conversation_id, source_id, question, intent, confidence,
                         tables_detected, slots_json, needs_clarification, clarification_json, response_ms)
                    VALUES ($1,$2,$3,$4,$5,$6,$7,$8,$9,$10)
                """,
                    "chat",
                    UUID(source_id_str) if source_id_str else None,
                    question,
                    slots.intent,
                    float(slots.confidence),
                    slots.table_names or [],
                    json_module.dumps({"metric": slots.metric, "group_by": slots.group_by}),
                    False,
                    json_module.dumps({}),
                    0,
                )
        except Exception as _log_err:
            logger.warning(f"[Chat] nlu_query_log INSERT error: {_log_err}")

    # Injecte les tables détectées directement si le NLU les a ratées
    if slots and directly_mentioned:
        for t in directly_mentioned:
            if t not in slots.table_names:
                slots.table_names.insert(0, t)

    # Inject last_table from context if no table detected
    # Exception: LIST_FIELDS PK/FK global → ne pas injecter last_table
    context = get_context(source_id_str)
    _pk_fk_kws = ["primaire","primary","etrangere","étrangère","étrangères","etrangeres","foreign","pk","fk","clé","cle","index"]
    import unicodedata as _ucd_gf
    def _norm_gf(s): return ''.join(c for c in _ucd_gf.normalize('NFD',s.lower()) if _ucd_gf.category(c)!='Mn')
    _is_global_fields = (slots and slots.intent == Intent.LIST_FIELDS and
                         any(kw in _norm_gf(question) for kw in ["primaire","primary","etrangere","foreign","pk","fk","cle","index"]))
    # Si global fields, forcer table_names vide même si le NLU en a trouvé
    if slots and slots.intent == Intent.LIST_FIELDS and _is_global_fields:
        slots.table_names = []
    if slots and not slots.table_names and context.last_table and not _is_global_fields:
        if context.last_table in schema:
            slots.table_names = [context.last_table]
            logger.info(f"[Chat] Context injection: last_table='{context.last_table}'")

    # Correction intents + raw_text via _fix_slots()
    q_low = question.lower()
    if slots:
        _original_intent = slots.intent  # Sauvegarder avant _fix_slots
        slots = await _fix_slots(slots, question, schema, known_entities, source_id_str=source_id_str, pg_pool=None)
        # ── Protéger SHOW_DASHBOARD — priorité absolue ──────────────────
        # Si le NLU a détecté show_dashboard, _fix_slots ne peut pas l'écraser
        if _original_intent == Intent.SHOW_DASHBOARD:
            slots.intent = Intent.SHOW_DASHBOARD
            logger.info(f"[Chat] SHOW_DASHBOARD intent protégé contre _fix_slots")
        # Protéger LIST_FIELDS contre les surcharges de _fix_slots
        if (_original_intent == Intent.LIST_FIELDS and
            slots.intent != Intent.LIST_FIELDS and
            any(kw in question.lower() for kw in [
                "primaire", "primary", "étrangère", "etrangere", "foreign",
                "cle", "clé", "index", "champs", "colonnes", "fields", "columns"
            ])):
            slots.intent = Intent.LIST_FIELDS
            logger.info(f"[Chat] LIST_FIELDS intent protégé contre surcharge de _fix_slots")
    if slots:
        # "ayant / having / plus de N fois" → forcer GENERATE_AGG + injecter group_by
        _having_kw = ["ayant", "having", "au moins", "au plus", "dont le total"]
        # "plus de / moins de" → FILTER si sur un champ (prix, montant), AGG si sur un compte
        _filter_kw = ["prix", "price", "montant", "amount", "cout", "cost", "salaire", "wage"]
        _count_ctx  = ["fois", "commandes", "orders", "produits", "products", "articles", "lignes"]
        _is_filter  = any(kw in q_low for kw in _filter_kw) and not any(kw in q_low for kw in _count_ctx)
        _is_having  = any(kw in q_low for kw in _having_kw) or (
            any(kw in q_low for kw in ["plus de", "moins de"]) and not _is_filter
        )
        if _is_having and not _is_filter:
            if slots.intent != Intent.SHOW_DASHBOARD:
                slots.intent = Intent.GENERATE_AGG
            # Extraire le sujet de la phrase HAVING comme group_by
            # "banques ayant..." → group_by = "banque"
            import re as _re_hav
            _hav_subj = _re_hav.match(r"(\w+)\s+(?:ayant|having|dont|with)", q_low)
            if _hav_subj and not slots.group_by:
                import unicodedata as _ud2
                _subj = _hav_subj.group(1)
                _subj = _ud2.normalize('NFD', _subj).encode('ascii','ignore').decode().lower()
                if _subj.endswith('es'): _subj = _subj[:-1]
                elif _subj.endswith('s'): _subj = _subj[:-1]
                if _subj not in ("le","la","les","l","un","une","des","du"):
                    slots.group_by = _subj
        elif _is_filter and any(kw in q_low for kw in ["plus de", "superieur", "supérieur", "greater", "inferieur", "inférieur", "moins de"]):
            slots.intent = Intent.GENERATE_FILTER
            # group_by = première table mentionnée si pas déjà défini
            if not slots.group_by and slots.table_names:
                # Cherche un champ ID ou nom dans le schéma de la table
                tbl = slots.table_names[0]
                tbl_fields = schema.get(tbl, [])
                # Priorité : champ name/nom/code, sinon premier champ non-PK
                gby = next(
                    (f for f in tbl_fields if any(k in f.lower() for k in ["name","nom","code","company","label"])),
                    next((f for f in tbl_fields if "id" not in f.lower()), tbl_fields[0] if tbl_fields else None)
                )
                if gby:
                    slots.group_by = gby
            # Extrait la valeur numérique HAVING depuis la question ("plus de 5" → 5)
            import re as _re2
            _num = _re2.search(r"(?:plus de|moins de|au moins|au plus|ayant)\s+(\d+)", q_low)
            if _num:
                _op = "lt" if any(k in q_low for k in ["moins de", "inferieur", "inférieur"]) else "gt"
                # Toujours écraser — notre regex est plus fiable que le NLU
                slots.amount_filter = {"op": _op, "value": int(_num.group(1))}
            elif slots.amount_filter:
                # Corrige op=eq → gt par défaut pour les HAVING
                if slots.amount_filter.get("op") in ("eq", "="):
                    slots.amount_filter["op"] = "gt"

        # "montre / affiche / liste les X" + table connue → forcer GENERATE_SQL
        _show_kw = ["montre", "affiche", "donne", "show", "display"]
        _data_kw = ["premiers", "derniers", "chers", "recents", "first", "last", "expensive"]
        if (slots.intent == Intent.LIST_ENTITIES and slots.table_names and
                (any(kw in q_low for kw in _show_kw) or any(kw in q_low for kw in _data_kw))):
            slots.intent = Intent.GENERATE_SQL

        # "nombre de X par Y" → forcer COUNT + GENERATE_AGG
        _count_kw = ["nombre de", "number of", "combien de", "count of"]
        _per_kw   = ["par", "by", "per", "pour chaque", "for each"]
        if (any(kw in q_low for kw in _count_kw) and
                any(kw in q_low for kw in _per_kw) and
                slots.intent not in (Intent.GENERATE_AGG,)):
            if slots.intent != Intent.SHOW_DASHBOARD:
                slots.intent = Intent.GENERATE_AGG
            if not slots.metric:
                slots.metric = "COUNT"

        # "mois précédent / par rapport" → forcer window LAG
        _lag_kw = ["mois precedent", "mois précédent", "par rapport au", "comparaison mois", "vs mois"]
        if any(kw in q_low for kw in _lag_kw):
            if slots.intent != Intent.SHOW_DASHBOARD:
                slots.intent = Intent.GENERATE_AGG

        # "top N X par Y" ou "rang des X par Y" → window ROW_NUMBER
        import re as _re3
        # Pattern: "top 5 banques par montant" ou "rang des banques par montant"
        _top_pat = _re3.search(r"(?:top\s+\d+\s+|rang\s+des?\s+)(\w+)\s+par\s+(\w+)", q_low)
        if _top_pat and slots.table_names:
            if slots.intent != Intent.SHOW_DASHBOARD:
                slots.intent = Intent.GENERATE_AGG
            slots.raw_text = question
            # Effacer amount_filter si c est top_n mal interprete
            if slots.amount_filter:
                try:
                    # Effacer si op=top_n OU si value=top_n avec op=eq
                    af_op = slots.amount_filter.get("op","")
                    af_val = float(slots.amount_filter.get("value",-1))
                    if (af_op == "top_n" or
                        (af_op in ("eq","=") and abs(af_val - float(slots.top_n)) < 0.001)):
                        slots.amount_filter = None
                except Exception:
                    pass
            group_entity = _top_pat.group(1)   # "banques", "sociétés"...
            sort_metric  = _top_pat.group(2)   # "montant", "ventes"...
            _skip = ("le","la","les","l","du","de","des","sur","dans","un","une")
            # group_by = entité (mot avant "par") — FORCER même si NLU a mis autre chose
            if group_entity and group_entity not in _skip:
                # Normalise pluriel → singulier : "banques"→"banque", "sociétés"→"société"
                import unicodedata as _ud
                def _norm_word(w):
                    # Normalise accents
                    w_ascii = _ud.normalize('NFD', w).encode('ascii','ignore').decode().lower()
                    # Pluriel → singulier sur version originale (avec accents)
                    w_orig = w.lower()
                    # "sociétés" → "société", "banques" → "banque"
                    if w_orig.endswith('és') or w_orig.endswith('es'):
                        return w_orig[:-1]  # enlève juste le s
                    elif w_orig.endswith('s'):
                        return w_orig[:-1]
                    return w_orig
                slots.group_by = _norm_word(group_entity)  # force override NLU
            # metric = SUM si sort_metric est un mot de montant
            _amount_kw = ["montant","amount","vente","ventes","total","prix","chiffre","valeur"]
            if any(kw in sort_metric for kw in _amount_kw):
                if not slots.metric:
                    slots.metric = "SUM"

        # Mots métier → mapping table implicite si aucune table détectée
        _entity_kw_map = {
            "fournisseur": ["Suppliers","SUPPLIER","BPSUPPLIER"],
            "supplier":    ["Suppliers","SUPPLIER"],
            "produit":     ["Products","PRODUCT","ITMMATER"],
            "product":     ["Products","PRODUCT"],
            "client":      ["Customers","CUSTOMER","BPCUSTOMER"],
            "customer":    ["Customers","CUSTOMER"],
            "commande":    ["Orders","ORDER","SORDERQ"],
            "order":       ["Orders","SORDER"],
            "employe":     ["Employees","EMPLOYEE"],
            "employee":    ["Employees","EMPLOYEE"],
            "categorie":   ["Categories","CATEGORY"],
            "category":    ["Categories","CATEGORY"],
            # Fix SXA — tables métier manquantes
            "transaction": ["Transactions bancaires","Transactions"],
            "journal":     ["Journal"],
            "financement": ["Ligne de financement","FINANCEMENT_BI"],
            "amortissement":["Tableaux d'amortissement"],
            "integration": ["Dernière integration bancaire","SI_Bancaire"],
        }
        if not slots.table_names:
            for kw, candidates in _entity_kw_map.items():
                if kw in q_low:
                    for c in candidates:
                        matched = next((e for e in known_entities if e.lower() == c.lower()), None)
                        if matched:
                            slots.table_names = [matched]
                            break
                    if slots.table_names:
                        break

    # ── Intents catalogue (pas de SQL) ────────────────────────────────
    q = question.lower()

    if slots and slots.intent == Intent.COUNT_ENTITIES:
        # Si une table est identifiée → COUNT SQL sur cette table
        if slots.table_names and slots.table_names[0] in schema:
            if slots.intent != Intent.SHOW_DASHBOARD:
                slots.intent = Intent.GENERATE_AGG
            if not slots.metric:
                slots.metric = "COUNT"
        # Sinon → réponse catalogue
        else:
            sample = ", ".join(f"`{e.name}`" for e in entities[:6])
            extra  = f"\n\nExemples : {sample}{'…' if entity_count > 6 else ''}" if sample else ""
            return f"**{name}** contient **{entity_count} entités** indexées.{extra}"

    # ── Sprint 7C : Routing conceptuel en amont — avant tout pipeline SQL ──
    # Vérification AVANT LIST_ENTITIES pour éviter que "quels sont les ERP ?"
    # soit intercepté comme LIST_ENTITIES et retourne la liste des tables SXA.
    try:
        from .rag_engine import is_conceptual_question as _is_conceptual_early
        if _is_conceptual_early(question):
            logger.info(f"[CRAG] Routing conceptuel précoce → '{question[:60]}'")
            try:
                import httpx as _httpx_early
                _cp = (
                    f"Tu es un expert ERP. Réponds en français à cette question "
                    f"de manière concise et informative.\n"
                    f"Question : {question}\n\n"
                    f"Réponds directement sans générer de SQL. "
                    f"Donne une réponse claire en 3-5 phrases."
                )
                _cr = _httpx_early.post(
                    f"{OLLAMA_HOST}/api/generate",
                    json={"model": OLLAMA_MODEL, "prompt": _cp, "stream": False,
                          "options": {"temperature": 0.3, "num_predict": 400}},
                    timeout=30,
                )
                if _cr.status_code == 200:
                    _ans = _cr.json().get("response", "").strip()
                    if _ans:
                        return _ans
            except Exception as _ce_early:
                logger.warning(f"[CRAG] Réponse conceptuelle précoce échouée : {_ce_early}")
            return "Cette question est de nature conceptuelle. Précisez si vous souhaitez une requête SQL sur les données SXA."
    except ImportError:
        pass

    # ── Sprint8 : Bypass LIST_ENTITIES pour agent ────────────────────
    # Élargi Sprint 8.5 : couvre tous les préfixes synonymes + nouveaux patterns
    _bypass_kw = [
        # Patterns directs SXA_DIRECT_SQL
        "codes pays", "liste les pays", "codes iso",
        "liste les devises", "devises disponibles",
        "taux de change", "cours de change", "cours marché", "forex",
        "liste les sociétés", "toutes les sociétés",
        "liste les banques", "toutes les banques",
        "intégration bancaire", "flux de trésorerie", "journal des flux",
        "solde bancaire", "si bancaire", "rapprochement bancaire",
        "groupe de sociétés",
        # Préfixes synonymes — interceptés avant résolution de table
        "affiche-moi", "montre-moi", "donne-moi",
        "afficher les", "montrer les", "affiche les",
        "quelles sont les", "quels sont les",
        "liste des", "lister les", "lister des",
        "utilisateurs bloqués", "utilisateurs avec", "utilisateurs ayant",
        "comptes actifs", "comptes par",
        "quelle banque", "quel est le montant", "combien de",
        "nombre de comptes", "comptes par banque",
        # Questions complexes
        "utilisateurs avec", "jointure entre",
        "liste les utilisateurs", "affiche les utilisateurs",
    ]
    if slots and slots.intent == Intent.LIST_ENTITIES and not slots.table_names:
        if any(kw in question.lower() for kw in _bypass_kw):
            slots.intent = Intent.GENERATE_SQL
            goto_sql = True
            logger.info(f"[Sprint8] LIST_ENTITIES bypasse pour '{question[:50]}'")

    if slots and slots.intent == Intent.LIST_ENTITIES and not slots.table_names:
        # Si top_n présent → l'utilisateur veut des données, pas la liste des tables
        if slots.top_n and schema:
            # Redirige vers SQL sur la première table du schéma
            # Fix : jamais prendre une table infra comme table par défaut
            _INFRA = ("QRTZ_","qrtz_","sys","SYS","dt_","DT_","MSreplication","sysdiagram","__")
            # Exclure les tables audit (_A, _FL, _RPT, _LOG, _HIST) et préférer GS_/AA_
            _AUDIT_SFX = ("_A", "_FL", "_RPT", "_LOG", "_HIST", "_AUD")
            first_table = next(
                (t for t in schema.keys()
                 if not any(t.startswith(p) for p in _INFRA)
                 and not any(t.endswith(s) for s in _AUDIT_SFX)
                 and (t.startswith("GS_") or t.startswith("AA_") or t.startswith("CS_"))),
                next(
                    (t for t in schema.keys()
                     if not any(t.startswith(p) for p in _INFRA)
                     and not any(t.endswith(s) for s in _AUDIT_SFX)),
                    next(
                        (t for t in schema.keys() if not any(t.startswith(p) for p in _INFRA)),
                        list(schema.keys())[0] if schema else "table"
                    )
                )
            )
            slots.table_names = [first_table]
            slots.intent = Intent.GENERATE_SQL
            # Continue vers le pipeline SQL ci-dessous
        else:
            if not entities:
                return f"Aucune entité synchronisée pour **{name}**.\n\nCliquez sur **↻ Sync**."
            rows = "\n".join(f"{i+1}. `{e.name}` — {e.field_count or 0} champs" for i, e in enumerate(entities[:10]))
            more = f"\n\n*… et {entity_count - 10} autres entités.*" if entity_count > 10 else ""
            return f"Entités de **{name}** :\n\n{rows}{more}"

    if slots and slots.intent == Intent.GET_RELATIONS:
        if not relations:
            return f"Aucune relation détectée pour **{name}**.\n\nLancez la découverte depuis **Relations → ↻ Relancer**."
        rows = "\n".join(
            f"- `{r.get('source_entity','?')}.{r.get('source_field','?')}` → "
            f"`{r.get('target_entity','?')}.{r.get('target_field','?')}` "
            f"*({r.get('detection_method','?')}, {round((r.get('confidence') or 1)*100)}%)*"
            for r in relations
        )
        return f"Relations détectées pour **{name}** :\n\n{rows}\n\nVoir onglet **Relations**."

    if slots and slots.intent == Intent.DESCRIBE_ENTITY:
        host_info = f"`{source.host}:{source.port}`" if source.host else "—"
        return (
            f"**{name}**\n\n"
            f"- **Type** : {source.connector_type}\n"
            f"- **Hôte** : {host_info}\n"
            f"- **Base** : `{source.database_name or source.base_url or '—'}`\n"
            f"- **Entités** : {entity_count}\n"
            f"- **Relations** : {len(relations)}\n"
            f"- **Statut** : {source.status}"
        )

    if slots and slots.intent == Intent.LIST_FIELDS:
        q_low_lf = question.lower()
        is_pk_query = any(kw in q_low_lf for kw in [
            "primaire", "primary", "clé", "cle", "pk", "index", "foreign", "étrangère"
        ])

        if is_pk_query and not slots.table_names:
            # "Quels champs sont des clés primaires ?" → cherche les PK dans toutes les entités
            try:
                pool_lf = await get_pg_pool()
                async with pool_lf.acquire() as conn_lf:
                    pk_flag = "primary" in q_low_lf or "primaire" in q_low_lf or "pk" in q_low_lf
                    import unicodedata as _ucd2
                    def _n2(s): return ''.join(ch for ch in _ucd2.normalize('NFD',s.lower()) if _ucd2.category(ch)!='Mn')
                    q_norm_lf = _n2(q_low_lf)
                    fk_flag = "foreign" in q_norm_lf or "etrangere" in q_norm_lf or "fk" in q_norm_lf

                    # Requête ciblée selon le type de clé demandé
                    if pk_flag and not fk_flag:
                        where_key = "ef.is_primary_key = TRUE"
                        title = "clés primaires"
                        icon = "PK"
                    elif fk_flag and not pk_flag:
                        where_key = "ef.is_foreign_key = TRUE AND ef.is_primary_key = FALSE"
                        title = "clés étrangères"
                        icon = "FK"
                    else:
                        where_key = "(ef.is_primary_key = TRUE OR ef.is_foreign_key = TRUE)"
                        title = "clés (PK + FK)"
                        icon = "PK"

                    pk_rows = await conn_lf.fetch(f"""
                        SELECT se.name AS table_name, ef.name AS field_name,
                               ef.data_type, ef.is_primary_key, ef.is_foreign_key
                        FROM source_entities se
                        JOIN entity_fields ef ON ef.entity_id = se.id
                        WHERE se.source_id = $1
                          AND se.is_visible = TRUE
                          AND {where_key}
                        ORDER BY se.name, ef.position
                        LIMIT 500
                    """, source_id)
                    # Compte total — dans le même bloc async with
                    count_row = await conn_lf.fetchrow(f"""
                        SELECT COUNT(*) as total
                        FROM source_entities se
                        JOIN entity_fields ef ON ef.entity_id = se.id
                        WHERE se.source_id = $1 AND se.is_visible = TRUE
                          AND {where_key}
                    """, source_id)
                    total_count = count_row["total"] if count_row else len(pk_rows)
                if pk_rows:

                    rows_txt = "\n".join(
                        f"- `{r['table_name']}`.`{r['field_name']}` ({r['data_type']}) "
                        f"{'PK' if r['is_primary_key'] and not r['is_foreign_key'] else 'FK' if r['is_foreign_key'] and not r['is_primary_key'] else 'PK+FK'}"
                        for r in pk_rows[:30]
                    )
                    more = f"\n\n*… et {total_count - 30} autres.*" if total_count > 30 else ""
                    return f"{icon} **{title.capitalize()} dans {name}** ({total_count} trouvées) :\n\n{rows_txt}{more}"
                else:
                    return f"Aucune {title} détectée dans **{name}**. Les métadonnées de clés dépendent du connecteur source."
            except Exception as e_lf:
                logger.warning(f"[Chat] PK query error: {e_lf}")

        if entities:
            e = entities[0]
            flds  = e.fields[:15] if e.fields else []
            rows  = "\n".join(
                f"- `{f.name}` ({f.data_type})"
                f"{' PK' if f.is_primary_key else ''}{' FK' if f.is_foreign_key else ''}"
                for f in flds
            )
            pk_count = sum(1 for f in flds if f.is_primary_key)
            fk_count = sum(1 for f in flds if f.is_foreign_key)
            meta = f" — {pk_count} PK, {fk_count} FK" if (pk_count or fk_count) else ""
            return f"**Champs de `{e.name}`**{meta} :\n\n{rows}"
        return "Aucune entité chargée. Lancez une synchronisation d'abord."

    # ── Dashboard intent → génération automatique ────────────
    if slots and slots.intent == Intent.SHOW_DASHBOARD:
        try:
            from .dashboard_engine import get_dashboard_generator
            generator = get_dashboard_generator()
            spec = await generator.generate(
                question          = question,
                slots             = slots,
                schema            = schema,
                source_id         = source_id_str,
                pg_pool           = await get_pg_pool(),
                redis             = await get_redis(),
                connector_factory = ConnectorFactory,
            )
            spec_dict = spec.to_dict()
            n_widgets = len(spec_dict.get("widgets", []))
            # Retourne uniquement le bloc dashboard JSON — le frontend gère l'affichage
            return f"```dashboard\n{json_module.dumps(spec_dict, default=str)}\n```"
        except Exception as e:
            logger.error(f"[Chat/Dashboard] {e}")
            return f"Erreur génération dashboard : {e}"

    if slots and slots.intent in (Intent.GREETING, Intent.HELP):
        return (
            f"Bonjour, je suis **OnePilot**, votre assistant données.\n\n"
            f"Source active : **{name}** ({entity_count} entités)\n\n"
            f"**Exemples de questions :**\n"
            f"- *Total des ventes par client*\n"
            f"- *Top 10 commandes ce mois*\n"
            f"- *Clients ayant commandé plus de 5 fois*\n"
            f"- *Total cumulatif des ventes par date*\n"
            f"- *Jointure entre commandes et clients*"
        )

    # ── LLM_EXPLAIN : questions conceptuelles → Ollama ───────────────
    if slots and slots.intent == "llm_explain":
        try:
            import httpx as _httpx
            OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://host.docker.internal:11434")
            # Contexte source pour enrichir la réponse
            domain_ctx = f"Source active: {name} ({entity_count} entités ERP)." if name else ""
            prompt = f"""{domain_ctx}
Tu es OnePilot, un assistant expert en systèmes ERP et données d'entreprise.
Réponds en français de manière claire et structurée à la question suivante :

{question}

Donne une réponse concise (5-10 lignes max), pratique et orientée métier."""

            # Routing intelligent : mistral pour questions NL, qwen pour SQL/code
            q_lower = question.lower()
            is_technical = any(k in q_lower for k in [
                'sql','requête','code','script','query','table','colonne',
                'jointure','select','insert','update','index','schéma'
            ])
            llm_model = "qwen2.5-coder:3b" if is_technical else "mistral:latest"
            logger.info(f"[LLM Explain] model={llm_model} (technical={is_technical})")

            async with _httpx.AsyncClient(timeout=120) as _cli:
                resp = await _cli.post(
                    f"{OLLAMA_HOST}/api/generate",
                    json={"model": llm_model, "prompt": prompt, "stream": False,
                          "options": {"num_predict": 400, "temperature": 0.3}}
                )
                if resp.status_code == 200:
                    data = resp.json()
                    answer = data.get("response", "").strip()
                    if answer:
                        model_label = "Mistral 7B" if "mistral" in llm_model else "Qwen2.5-Coder 3B"
                        return f"*LLM — via {model_label}*\n\n{answer}"
        except Exception as _e:
            import traceback
            logger.warning(f"[LLM Explain] {_e}\n{traceback.format_exc()}")
        # Fallback si LLM indisponible
        return f"❓ Je ne peux pas répondre à cette question conceptuelle pour le moment. Essayez de poser une question sur les données de **{name}**."

    # ── Intents SQL : pipeline complet ───────────────────────────────
    if not schema:
        return f"Aucun schéma disponible pour **{name}**. Lancez une synchronisation d'abord."

    # Vérification ambiguïtés bloquantes
    if slots and not goto_sql:
        try:
            resolver  = AmbiguityResolver()
            questions = resolver.analyze(slots, known_entities, schema)
            if questions and questions[0].required:
                clarif = resolver.build_clarification_response(questions)
                opts   = "\n".join(f"- {o}" for o in clarif.get("options", [])[:5])
                # Stocke l'état de clarification dans le contexte
                ctx_store = get_context(source_id_str)
                ctx_store.pending_clarification = clarif
                ctx_store.pending_slots = slots
                ctx_store.pending_question = question
                logger.info(f"[Chat] Clarification stockée — slot={clarif.get('slot_key')} options={clarif.get('options', [])[:3]}")
                return f"❓ **{clarif['question']}**\n\n{opts}"
        except Exception as e:
            logger.warning(f"[Chat] Ambiguity resolver error: {e}")

    # Dialecte SQL
    dialect_map = {
        "mssql": "mssql", "sage_100": "mssql", "sage_x3": "mssql",
        "mysql": "mysql", "postgresql": "postgresql", "sqlite": "sqlite",
    }
    dialect = dialect_map.get(source.connector_type.value if hasattr(source.connector_type, 'value') else str(source.connector_type), "mssql")

    # Génération SQL
    try:
        sql_gen = SQLGenerator()
        result  = sql_gen.generate(slots, schema, dialect) if slots else {
            "sql": None, "explanation": "", "warnings": [],
            "validation": {"valid": False, "errors": ["NLU indisponible"], "warnings": [], "score": 0},
            "complexity": "simple",
        }
    except Exception as e:
        logger.warning(f"[Chat] SQLGenerator error: {e}")
        return f"Erreur génération SQL : {e}"

    sql        = result.get("sql", "")
    validation = result.get("validation", {})
    complexity = result.get("complexity", "simple")
    explanation = result.get("explanation", "").replace("*", "×")
    warnings   = result.get("warnings", [])
    ms         = int((time.time() - t0) * 1000)

    # Met à jour le contexte avec la table utilisée
    if slots and slots.table_names:
        try:
            from .nlu_engine import ConversationTurn
            turn = ConversationTurn(
                question  = question,
                intent    = slots.intent,
                slots     = slots,
                answer    = sql or "",
            )
            context.add_turn(turn)
        except Exception:
            pass

    # ── Sprint 7A : LLM+RAG si SQLGenerator donne mauvais résultat ─────
    score_val = validation.get("score", 1.0)

    # Mots-clés qui nécessitent le vrai schéma RAG (SXA financier)
    _financial_kw = [
        "solde", "trésorerie", "tresorerie", "bancaire", "banque",
        "montant", "clôture", "cloture", "devise", "financement",
        "compte", "virement", "flux", "balance", "sum", "somme",
        # Fix : mots SXA manquants
        "transaction", "journal", "ligne", "société", "societe",
        "groupe", "integration", "closingbalance", "amortissement",
    ]
    _needs_rag = any(kw in question.lower() for kw in _financial_kw)

    # ── Force LLM+RAG pour questions interceptées par list_entities/count_entities ──
    _force_rag_kw = [
        "ont accès", "ont acces", "accès à", "acces a", "droits d", "permissions",
        "utilisateurs", "utilisateur", "liste les utilisateurs",
        "codes pays", "liste les pays", "codes iso",
        "jointure", "relation entre",
    ]
    if any(kw in question.lower() for kw in _force_rag_kw):
        _needs_rag = True

    # ── Sprint8 : Agent direct pour questions connues ─────────────────
    _direct_kw = ["codes pays","liste les pays","codes iso","liste les devises",
                  "devises disponibles","utilisateurs avec","jointure entre",
                  "liste les utilisateurs","affiche les utilisateurs"]
    if any(kw in question.lower() for kw in _direct_kw):
        try:
            from .agentic_rag import run_agentic_rag as _run_ag
            _pg = await get_pg_pool()
            _sr = await _pg.fetchrow(
                "SELECT connector_type,host,port,database_name,username,base_url FROM data_sources WHERE id=$1",
                UUID(source_id))
            if _sr:
                _sd = dict(_sr); _sd["id"]=source_id; _sd["password"]=""
                try:
                    async with _pg.acquire() as _c:
                        _s = await _c.fetchrow(
                            "SELECT secret_value FROM connection_secrets WHERE source_id=$1 AND secret_key='password'",
                            UUID(source_id))
                        if _s: _sd["password"]=_s["secret_value"]
                except Exception: pass
                _ar = await _run_ag(question=question,source_id=UUID(source_id),pg_pool=_pg,source_dict=_sd,dialect=dialect)
                if _ar.success and _ar.sql:
                    logger.info(f"[Sprint8] Agent direct OK — {_ar.method}")
                    return _ar.sql
        except Exception as _e:
            logger.warning(f"[Sprint8] Agent direct echec: {_e}")

    _use_llm_rag = (
        not sql
        or score_val < 0.7
        or complexity in ("cte", "advanced", "window")
        or _needs_rag
    )

    # ── Sprint 8 : Agentic RAG — activé pour questions complexes ──────────
    AGENTIC_RAG_ENABLED = True
    _q_lower = question.lower()

    # Mots-clés multi-entités — déclenchent l'agent indépendamment du SQLGenerator
    _multi_entity_kw = [
        "avec leur", "avec sa", "avec son", "et leur", "et sa", "et son",
        "et les", "avec les", "ainsi que",
        "ont accès", "ont acces", "accès à", "acces a",
        "jointure", "relation entre",
        "par compte", "par devise",
        "montant total", "total des virements", "total par",
        "virement", "paiement",
        "employé", "employe", "personnel",
    ]
    _is_multi_entity = any(kw in _q_lower for kw in _multi_entity_kw)

    _use_agentic = (
        AGENTIC_RAG_ENABLED
        and (
            # Cas 1 : questions complexes via pipeline LLM
            (_use_llm_rag and complexity in ("cte", "advanced", "window", "join", "aggregate"))
            # Cas 2 : questions multi-entités — toujours actif
            or _is_multi_entity
        )
    )

    # ── Sprint 8 : Agentic RAG ────────────────────────────────────────────
    if _use_agentic:
        try:
            from .agentic_rag import run_agentic_rag, agent_result_to_dict
            _pool_agent = await get_pg_pool()
            # Construire source_dict pour ConnectorFactory
            _src_dict = {
                "id":             source_id_str,
                "connector_type": source.connector_type.value if hasattr(source.connector_type, "value") else str(source.connector_type),
                "host":           source.host,
                "port":           source.port,
                "database_name":  source.database_name,
                "username":       source.username,
                "base_url":       source.base_url,
            }
            # Récupérer le mot de passe depuis les secrets
            try:
                async with _pool_agent.acquire() as _conn_sec:
                    _sec = await _conn_sec.fetchrow(
                        "SELECT secret_value FROM connection_secrets WHERE source_id=$1 AND secret_key='password'",
                        source_id
                    )
                    if _sec:
                        _src_dict["password"] = _sec["secret_value"]
            except Exception:
                _src_dict["password"] = ""

            _agent_result = await run_agentic_rag(
                question    = question,
                source_id   = source_id,
                pg_pool     = _pool_agent,
                source_dict = _src_dict,
                dialect     = dialect,
            )
            if _agent_result.success and _agent_result.sql:
                sql         = _agent_result.sql
                explanation = f"Agentic RAG ({_agent_result.iterations} itérations)"
                complexity  = "agentic_rag"
                if _agent_result.warnings:
                    warnings.extend(_agent_result.warnings)
                logger.info(
                    f"[AgentRAG] Succès — {_agent_result.iterations} itérations "
                    f"| {_agent_result.duration_ms}ms"
                )
                # Bypass du pipeline LLM standard si on a un bon résultat
                if _agent_result.result is not None:
                    _use_llm_rag = False
        except Exception as _ae:
            logger.warning(f"[AgentRAG] Échec, fallback vers LLM+RAG : {_ae}")

    if _use_llm_rag:
        try:
            from .llm_engine import generate_sql_with_llm
            _pool_rag = await get_pg_pool()
            llm_result = await generate_sql_with_llm(
                question    = question,
                schema      = schema,
                dialect     = dialect,
                table_names = slots.table_names if slots else None,
                pg_pool     = _pool_rag,
                source_id   = source_id_str,
            )

            # ── Sprint 7C : Routing question conceptuelle ─────────────────────
            if llm_result.get("is_conceptual"):
                logger.info(f"[CRAG] Question conceptuelle → réponse LLM texte")
                # Laisser tomber vers le LLM explain normal (intent llm_explain)
                # en forçant la réponse via Ollama en texte libre
                try:
                    import httpx as _httpx
                    _conceptual_prompt = (
                        f"Tu es un expert ERP. Réponds en français à cette question de manière concise et informative.\n"
                        f"Question : {question}\n\n"
                        f"Réponds directement sans générer de SQL. Donne une réponse claire en 2-4 phrases."
                    )
                    _r = _httpx.post(
                        f"{OLLAMA_HOST}/api/generate",
                        json={"model": OLLAMA_MODEL, "prompt": _conceptual_prompt, "stream": False,
                              "options": {"temperature": 0.3, "num_predict": 300}},
                        timeout=30,
                    )
                    if _r.status_code == 200:
                        return _r.json().get("response", "").strip()
                except Exception as _ce:
                    logger.warning(f"[CRAG] Réponse conceptuelle échouée : {_ce}")
                return f"Cette question est de nature conceptuelle. Précisez si vous souhaitez une requête SQL spécifique."

            if llm_result.get("sql"):
                sql         = llm_result["sql"]
                explanation = "LLM+RAG"
                complexity  = "llm_rag"
                # Ajouter warnings colonnes invalides
                col_warnings = llm_result.get("warnings", [])
                if col_warnings:
                    warnings.extend(col_warnings)
                # ── Sauvegarde SQL pour correction interactive ────────────
                if sql:
                    try:
                        _sv_pool = await get_pg_pool()
                        async with _sv_pool.acquire() as _sv_conn:
                            # Update la ligne existante
                            _upd = await _sv_conn.execute("""
                                UPDATE nlu_query_log
                                SET    sql_generated = $1
                                WHERE  id = (
                                    SELECT id FROM nlu_query_log
                                    WHERE  source_id = $2
                                      AND  question  = $3
                                    ORDER  BY created_at DESC
                                    LIMIT  1
                                )
                            """, sql, UUID(source_id), question)
                            # Si aucune ligne → INSERT
                            if _upd == "UPDATE 0":
                                await _sv_conn.execute("""
                                    INSERT INTO nlu_query_log
                                        (source_id, question, intent, confidence,
                                         sql_generated, created_at)
                                    VALUES ($1,$2,$3,$4,$5,NOW())
                                """, UUID(source_id), question,
                                   slots.intent if slots else "generate_aggregate",
                                   float(slots.confidence) if slots and slots.confidence else 0.9,
                                   sql)
                        logger.info(f"[Chat] sql_generated sauvegardé — {len(sql)} chars")
                    except Exception as _sv_err:
                        logger.warning(f"[Chat] sql save error: {_sv_err}")
        except Exception as _e:
            logger.warning(f"[Chat] LLM+RAG échoué : {_e}")

    # ── Réponse si SQL invalide (injection / DDL) ─────────────────────
    if not validation.get("valid", True):
        errs = "\n".join(f"- {e}" for e in validation.get("errors", []))
        return (
            f"**Requête refusée** — validation échouée :\n\n{errs}\n\n"
            f"*Seules les requêtes SELECT sont autorisées.*"
        )

    # ── Réponse avec SQL généré ───────────────────────────────────────
    if not sql:
        return f"Impossible de générer une requête SQL pour cette question.\n\nEssayez de préciser la table ou l'action souhaitée."

    # Badge complexité
    complexity_badge = {
        "window":   "Window",
        "having":   "HAVING",
        "cte":      "CTE",
        "advanced": "Avancé",
        "simple":   "",
        "moderate": "",
        "llm_rag": "LLM+RAG",
    }.get(complexity, "")

    # Warnings
    warn_text = ""
    if warnings:
        warn_text = "\n\n" + " | ".join(warnings[:2])

    # Score validation
    score = validation.get("score", 1.0)
    score_badge = f" score={score}" if score >= 0.9 else f" score={score}"

    badge_line = f"\n*{complexity_badge}{score_badge} — {ms}ms*" if complexity_badge else f"\n*{score_badge} — {ms}ms*"

    return (
        f"**{explanation}**\n\n"
        f"```sql\n{sql}\n```"
        f"{warn_text}"
        f"{badge_line}"
    )




@app.delete("/conversations/{conv_id}", status_code=204)
async def delete_conversation(conv_id: str):
    """Supprimer une conversation et ses dashboards associés"""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            try:
                uuid_val = UUID(conv_id)
            except ValueError:
                uuid_val = None
            if uuid_val:
                result = await conn.execute(
                    "DELETE FROM conversations WHERE id = $1 OR id::text = $2",
                    uuid_val, conv_id
                )
                # Supprimer aussi les dashboards associés
                await conn.execute(
                    "DELETE FROM conversation_dashboards WHERE conv_id = $1 OR conv_id = $2",
                    str(uuid_val), conv_id
                )
            else:
                result = await conn.execute(
                    "DELETE FROM conversations WHERE id::text = $1",
                    conv_id
                )
                await conn.execute(
                    "DELETE FROM conversation_dashboards WHERE conv_id = $1",
                    conv_id
                )
            if result == "DELETE 0":
                raise HTTPException(status_code=404, detail="Conversation not found")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Delete conversation error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.patch("/conversations/{conv_id}")
async def patch_conversation(conv_id: str, body: dict):
    """Renommer une conversation"""
    title = body.get("title", "").strip()
    if not title:
        raise HTTPException(status_code=400, detail="Title required")
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            try:
                uuid_val = UUID(conv_id)
            except ValueError:
                uuid_val = None
            if uuid_val:
                result = await conn.execute(
                    "UPDATE conversations SET title=$1, updated_at=NOW() WHERE id=$2 OR id::text=$3",
                    title, uuid_val, conv_id
                )
            else:
                result = await conn.execute(
                    "UPDATE conversations SET title=$1, updated_at=NOW() WHERE id::text=$2",
                    title, conv_id
                )
            if result == "UPDATE 0":
                raise HTTPException(status_code=404, detail="Conversation not found")
        return {"success": True, "title": title}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Patch conversation error: {e}")
        raise HTTPException(status_code=500, detail=str(e))



# ── Épinglage de conversations §UI ────────────────────────────────────────────

@app.post("/conversations/{conv_id}/pin", tags=["Conversations"])
async def pin_conversation(conv_id: str, user_id: str = "admin"):
    """Épingle une conversation (Option A — direct, sans dashboard requis)."""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.execute("""
                INSERT INTO pinned_conversations (conv_id, user_id)
                VALUES ($1::uuid, $2)
                ON CONFLICT (conv_id, user_id) DO NOTHING
            """, conv_id, user_id)
        return {"success": True, "conv_id": conv_id, "pinned": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/conversations/{conv_id}/pin", status_code=200, tags=["Conversations"])
async def unpin_conversation(conv_id: str, user_id: str = "admin"):
    """Désépingle une conversation."""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            result = await conn.execute("""
                DELETE FROM pinned_conversations
                WHERE conv_id = $1::uuid AND user_id = $2
            """, conv_id, user_id)
        return {"success": True, "conv_id": conv_id, "pinned": False}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/dashboard/conv/{conv_id}/delete", status_code=204, tags=["Dashboard"])
async def delete_conv_dashboards(conv_id: str):
    """Supprimer tous les dashboards d'une conversation"""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.execute(
                "DELETE FROM conversation_dashboards WHERE conv_id = $1",
                conv_id
            )
    except Exception as e:
        logger.warning(f"Delete conv dashboards: {e}")

@app.post("/chat/stream")
async def chat_stream(req: ChatRequest):
    """
    Vrai streaming SSE — pipeline complet token par token.
    Sprint Point 2 — Phase 1: NLU+préparation, Phase 2: LLM tokens, Phase 3: SQL final.

    Format événements SSE :
      data: {"type":"thinking", "content":"🔍 Analyse…"}
      data: {"type":"token",    "content":"SELECT"}
      data: {"type":"sql",      "content":"SELECT TOP 100 …", "method":"llm+rag"}
      data: {"type":"dashboard","spec":{…}}
      data: {"type":"done",     "duration_ms":3200}
      data: [DONE]
    """
    from fastapi.responses import StreamingResponse
    import time as _time

    source_id_str = req.source_id
    question      = req.question

    if not source_id_str or not question:
        raise HTTPException(status_code=400, detail="Missing source_id or question")

    async def _sse(obj: dict) -> str:
        """Sérialise un dict en ligne SSE."""
        return "data: " + json_module.dumps(obj, default=str) + "\n\n"

    async def _stream_pipeline():
        import re as _re
        import unicodedata as _ucd
        t0 = _time.time()

        # ── Vérification DDL sécurité ──────────────────────────────────
        _DDL = _re.compile(r"\b(DROP|INSERT|UPDATE|DELETE|TRUNCATE|ALTER|CREATE|EXEC|EXECUTE)\b", _re.IGNORECASE)
        if _DDL.search(question):
            yield await _sse({"type": "token", "content": "**Requête refusée** — instruction interdite.\n\n*Seules les questions en langage naturel sont autorisées.*"})
            yield await _sse({"type": "done", "duration_ms": 0})
            yield "data: [DONE]\n\n"
            return

        # ── Phase 1 : NLU + préparation schéma ────────────────────────
        yield await _sse({"type": "thinking", "content": "nlu_analysis"})

        try:
            source_id = UUID(source_id_str)
        except (ValueError, AttributeError):
            yield await _sse({"type": "token", "content": "ID source invalide."})
            yield await _sse({"type": "done", "duration_ms": 0})
            yield "data: [DONE]\n\n"
            return

        source = await get_source(source_id)
        if not source:
            yield await _sse({"type": "token", "content": "Source introuvable."})
            yield await _sse({"type": "done", "duration_ms": 0})
            yield "data: [DONE]\n\n"
            return

        # Charge le schéma
        schema: dict = {}
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                rows = await conn.fetch("""
                    SELECT se.name AS table_name, ef.name AS field_name
                    FROM source_entities se
                    JOIN entity_fields ef ON ef.entity_id = se.id
                    WHERE se.source_id = $1 AND se.is_visible = TRUE
                    ORDER BY se.name, ef.position LIMIT 20000
                """, source_id)
            for r in rows:
                schema.setdefault(r["table_name"], []).append(r["field_name"])
        except Exception as _e:
            logger.warning(f"[Stream] Schema fetch error: {_e}")

        # Dialecte SQL
        dialect_map = {
            "mssql": "mssql", "sage_100": "mssql", "sage_x3": "mssql",
            "mysql": "mysql", "postgresql": "postgresql", "sqlite": "sqlite",
        }
        dialect = dialect_map.get(
            source.connector_type.value if hasattr(source.connector_type, "value")
            else str(source.connector_type), "mssql"
        )

        # ── Détection intents non-SQL : fallback vers _build_chat_answer ─
        # Pour greeting, list_entities, list_fields, dashboard, get_relations etc.
        # on délègue au pipeline complet (pas de streaming token, mais cohérence)
        try:
            nlu = get_nlu_pipeline()
            context_nlu = get_context(source_id_str)
            known_entities = list(schema.keys())
            slots = nlu.process(question, context_nlu, known_entities)
        except Exception as _e:
            logger.warning(f"[Stream] NLU error: {_e}")
            slots = None

        # ── CORRECTION SQL INTERACTIVE — Sprint 13 ───────────────────────
        # Détecte si l'utilisateur veut corriger le dernier SQL généré
        _CORRECTION_PHRASES = [
            "c'est faux", "c'est pas bon", "c'est incorrect", "c'est mauvais",
            "corrige", "corriger", "correction", "il manque", "manque un",
            "manque les", "ajoute un", "ajoute la", "ajoute le", "ajoute les",
            "remplace", "change la", "change le", "change les", "modifie",
            "enlève", "enleve", "supprime le filtre", "mauvaise colonne",
            "mauvaise table", "erreur dans", "le sql est", "ca marche pas",
            "ça marche pas", "wrong", "fix", "incorrect", "missing bracket",
            "crochet", "crochets", "brackets", "pas le bon", "pas correct",
            "faut corriger", "il faut", "reessaie", "réessaie", "retry",
        ]
        # Garde-fou : vérifier que la question contient VRAIMENT un mot-clé de correction
        # FastText peut confondre certaines questions normales avec correct_sql
        _has_correction_keyword = any(phrase in question.lower() for phrase in _CORRECTION_PHRASES)
        _is_correction = _has_correction_keyword  # On n'utilise QUE les mots-clés, pas l'intent FastText
        # Note : intent FastText désactivé pour correct_sql car trop de faux positifs

        if _is_correction:
            yield await _sse({"type": "thinking", "content": "🔧 Correction SQL en cours…"})
            try:
                # Récupérer le dernier SQL de cette conversation
                _corr_pool = await get_pg_pool()
                _last_sql = None
                _last_question = None
                async with _corr_pool.acquire() as _corr_conn:
                    # conversation_id non disponible dans ChatRequest → fallback source_id
                    _conv_key = source_id_str
                    _last_row = await _corr_conn.fetchrow("""
                        SELECT question, sql_generated
                        FROM   nlu_query_log
                        WHERE  source_id = $1::uuid
                          AND  sql_generated   IS NOT NULL
                          AND  sql_generated   != ''
                        ORDER  BY created_at DESC
                        LIMIT  1
                    """, UUID(_conv_key))
                    if _last_row:
                        _last_sql      = _last_row["sql_generated"]
                        _last_question = _last_row["question"]

                if not _last_sql:
                    _msg = "Je n'ai pas de SQL précédent à corriger dans cette conversation. Posez d'abord une question pour générer un SQL."
                    for _w in _msg.split(" "):
                        yield await _sse({"type": "token", "content": _w + " "})
                        await asyncio.sleep(0.015)
                    yield await _sse({"type": "done", "duration_ms": 0})
                    yield "data: [DONE]\n\n"
                    return

                # Prompt de correction
                import httpx as _hx_corr
                OLLAMA_HOST  = os.environ.get("OLLAMA_HOST",  "http://host.docker.internal:11434")
                OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "qwen2.5-coder:3b")
                _corr_prompt = f"""Tu es un expert SQL Server (MSSQL). 
L'utilisateur veut corriger ce SQL.

SQL actuel:
{_last_sql}

Demande de correction: "{question}"

Règles STRICTES:
- Retourne UNIQUEMENT le SQL corrigé, sans explication, sans markdown
- Conserve la structure du SQL original
- Applique EXACTEMENT la correction demandée
- Pour les noms de colonnes avec espaces: utilise [crochets]
- Pour SXA: utilise uniquement les vues: [Comptes], [Transactions bancaires], [SI_Trésorerie], [FINANCEMENT_BI], [Journal], [Dernière integration bancaire]
- Ne change que ce qui est demandé, garde le reste identique

SQL corrigé:"""

                _corr_resp = _hx_corr.post(
                    f"{OLLAMA_HOST}/api/generate",
                    json={
                        "model":  OLLAMA_MODEL,
                        "prompt": _corr_prompt,
                        "stream": False,
                        "options": {"temperature": 0.05, "num_predict": 400},
                    },
                    timeout=45,
                )
                _corrected_sql = ""
                if _corr_resp.status_code == 200:
                    _raw = _corr_resp.json().get("response", "").strip()
                    # Nettoyer le SQL
                    import re as _re_corr
                    _raw = _re_corr.sub(r"```(?:sql|SQL)?\s*", "", _raw)
                    _raw = _re_corr.sub(r"```", "", _raw)
                    _corrected_sql = _raw.strip()

                if _corrected_sql:
                    # Sauvegarder le SQL corrigé dans nlu_query_log
                    async with _corr_pool.acquire() as _corr_conn2:
                        await _corr_conn2.execute("""
                            INSERT INTO nlu_query_log
                                (conversation_id, question, intent, confidence,
                                 sql_generated, created_at)
                            VALUES ($1, $2, 'correct_sql', 0.99, $3, NOW())
                            ON CONFLICT DO NOTHING
                        """, source_id_str,
                           f"[CORRECTION] {question[:200]}",
                           _corrected_sql)

                    yield await _sse({"type": "sql", "content": _corrected_sql, "method": "correction"})
                    _confirm = "✅ SQL corrigé. C'est mieux comme ça ?"
                    for _w in _confirm.split(" "):
                        yield await _sse({"type": "token", "content": _w + " "})
                        await asyncio.sleep(0.02)
                else:
                    _err = "⚠️ Je n'ai pas pu appliquer la correction. Décrivez plus précisément : par exemple 'ajoute [crochets] autour de montant avec signe'."
                    for _w in _err.split(" "):
                        yield await _sse({"type": "token", "content": _w + " "})
                        await asyncio.sleep(0.015)

            except Exception as _corr_err:
                logger.error(f"[Correction] Erreur: {_corr_err}")
                _msg = f"Erreur lors de la correction : {str(_corr_err)[:100]}"
                yield await _sse({"type": "token", "content": _msg})

            ms = int((_time.time() - t0) * 1000)
            yield await _sse({"type": "done", "duration_ms": ms})
            yield "data: [DONE]\n\n"
            return

        # ── FIN CORRECTION SQL ────────────────────────────────────────────────

        _NON_SQL_INTENTS = {
            "greeting", "help", "list_entities", "list_fields",
            "describe_entity", "get_relations",
            "show_dashboard", "llm_explain",
        }
        # Si intent non-SQL ET pas de mots-clés financiers/complexes → délègue à _build_chat_answer
        _financial_kw_early = [
            # Finance / trésorerie
            "solde", "trésorerie", "tresorerie", "bancaire", "banque",
            "montant", "clôture", "cloture", "devise", "financement",
            "compte", "virement", "flux", "balance", "somme",
            "mois dernier", "mois precedent", "trimestre", "annee derniere",
            "superieur", "inferieur", "moyenne", "chiffre",
            # Requêtes structurées qui nécessitent RAG
            "paramètre", "parametre", "paramètres", "parametres",
            "catégorie", "categorie", "catégories", "categories",
            "champ", "champs", "règle", "regle", "règles",
            "liste", "affiche", "montre", "donne", "quels", "quelles",
            "structure", "données de", "données du", "données des",
            # Agent RAG — forcer bypass list_entities
            "codes pays", "codes iso", "liste les pays",
            "liste les devises", "devises disponibles",
            "utilisateurs avec", "jointure entre",
            # Transactions avec filtres numériques — forcer generate_sql_stream
            "transaction", "supérieur", "superieur", "inférieur", "inferieur",
            "depasse", "dépasse", "plus de", "moins de", "superieure", "inferieure",
            "50000", "100000", "10000", "encaissement", "decaissement",
        ]
        _force_llm_stream = any(kw in question.lower() for kw in _financial_kw_early)
        if slots and slots.intent in _NON_SQL_INTENTS and not _force_llm_stream:
            yield await _sse({"type": "thinking", "content": "preparing"})
            try:
                answer = await _build_chat_answer(source_id_str, question)
            except Exception as _e:
                answer = f"Erreur : {_e}"
            # Stream la réponse mot par mot (délai minimal)
            words = answer.split(" ")
            for i, word in enumerate(words):
                token = word + (" " if i < len(words) - 1 else "")
                yield await _sse({"type": "token", "content": token})
                await asyncio.sleep(0.015)
            ms = int((_time.time() - t0) * 1000)
            yield await _sse({"type": "done", "duration_ms": ms})
            yield "data: [DONE]\n\n"
            return

        # ── Détection requête complexe → LLM streaming ────────────────
        from .llm_engine import is_complex_query, generate_sql_stream

        table_names_hint = slots.table_names if slots else None
        _financial_kw = [
            "solde", "trésorerie", "tresorerie", "bancaire", "banque",
            "montant", "clôture", "cloture", "devise", "financement",
            "compte", "virement", "flux", "balance", "somme",
            "codes pays", "codes iso", "liste les pays",
            "liste les devises", "devises disponibles",
            "utilisateurs avec", "affiche les utilisateurs",
            "jointure entre", "liste les utilisateurs",
            # Sprint 8.5 — synonymes et nouveaux patterns
            "taux de change", "cours de change", "cours marché", "forex",
            "affiche-moi", "montre-moi", "donne-moi",
            "afficher les", "montrer les",
            "quels sont les", "quelles sont les",
            "liste les sociétés", "liste les banques",
            "intégration bancaire", "flux de trésorerie",
            "solde bancaire", "si bancaire", "pays",
            "liste des", "lister les",
            "utilisateurs bloqués", "utilisateurs avec", "utilisateurs ayant",
            "comptes actifs", "comptes par",
            "quelle banque", "quel est le montant", "combien de",
            "nombre de comptes", "comptes par banque",
        "quelle banque", "quel est le montant", "combien de",
        "nombre de comptes", "comptes par banque",
        ]
        _needs_llm = (
            is_complex_query(question, table_names_hint)
            or any(kw in question.lower() for kw in _financial_kw)
        )

        if _needs_llm and schema:
            # ── Sprint 8 : Agentic RAG dans /chat/stream ─────────────
            _q_lower_stream = question.lower()
            _agent_kw_stream = [
                "avec leur", "avec sa", "et leur", "avec les",
                "ont accès", "accès à", "jointure",
                "par compte", "par devise",
                "montant total", "total des virements", "virement", "paiement",
                "employé", "employe", "personnel", "utilisateur", "utilisateurs",
                "liste les comptes", "comptes avec",
                "financement", "financements", "amortissement",
                "solde de trésorerie", "solde trésorerie",
                "total des transactions", "transactions par banque",
                # Devises / pays
                "liste les devises", "devises disponibles", "codes pays",
                "liste les pays", "codes iso", "jointure entre",
                "montre les devises", "affiche les devises", "quelles devises",
                "montre les pays", "affiche les pays", "quels pays",
                # Cours / taux
                "taux de change", "cours de change", "cours marché", "forex",
                # Sociétés / banques
                "liste les sociétés", "toutes les sociétés", "liste les banques",
                "toutes les banques", "groupe de sociétés",
                # Trésorerie / flux
                "intégration bancaire", "flux de trésorerie", "journal des flux",
                "solde bancaire", "si bancaire", "rapprochement bancaire",
                # Variantes synonymes fréquentes
                "montre-moi", "affiche-moi", "donne-moi", "montrer les",
                "afficher les", "quelles sont les", "quels sont les",
                "liste des", "lister les", "lister des",
                "utilisateurs bloqués", "utilisateurs avec", "utilisateurs ayant",
                "comptes actifs", "comptes par",
            "quelle banque", "quel est le montant", "combien de",
            "nombre de comptes", "comptes par banque",
                "quelle banque", "quel est le montant", "combien de",
                "nombre de comptes", "comptes par banque",
        "quelle banque", "quel est le montant", "combien de",
        "nombre de comptes", "comptes par banque",
            # ── Sprint 9 fix : questions composées compare/vs ───────────────────
            "compare ", "vs ", " vs ", "versus",
            # ── Sprint 8.6 fix : questions filtrées dynamiques ────────────────
            "transactions tnd", "transactions eur", "transactions usd",
            "supérieures à", "supérieure à", "inférieures à", "inférieure à",
            "banque postale", "la banque postale",
            "par devise en", "par devise",
            "en 2024", "en 2023", "en 2022", "en 2025",
            "maturité dépasse", "maturité supérieure",
            "accès à plus",
            ]
            # Sprint 8.6 : routing élargi — toute question non-triviale va vers AgentRAG
            # Le matching flou dans agentic_rag.py gère les variantes lexicales
            # ── Exclure les questions dashboard du pipeline AgentRAG ─────────
            # Les questions "dashboard ..." doivent aller vers dashboard_engine
            # via le bloc show_dashboard, pas vers l'Orchestrateur
            _is_dashboard_question = (
                _q_lower_stream.startswith("dashboard ") or
                _q_lower_stream.startswith("génère un dashboard") or
                _q_lower_stream.startswith("genere un dashboard") or
                _q_lower_stream.startswith("crée un dashboard") or
                _q_lower_stream.startswith("create dashboard") or
                " dashboard " in _q_lower_stream
            )
            # ── PRIORITÉ ABSOLUE : Pattern Direct SQL connu → bypass CRAG+LLM ─
            # Intercepte AVANT _is_complex_question et AVANT _use_agent_stream
            # Ex: "encaissements par banque", "flux trésorerie par banque", etc.
            if not _is_dashboard_question:
                try:
                    from .agentic_rag import _find_direct_sql, _tool_execute_sql
                    _matched_pat, _direct_sql_early, _match_score_early = _find_direct_sql(question)
                    if _matched_pat and _direct_sql_early:
                        logger.info(f"[DirectSQL Early] Pattern='{_matched_pat}' score={_match_score_early:.2f}")
                        _src_dict_early = {
                            "id": source_id_str,
                            "connector_type": source.connector_type.value if hasattr(source.connector_type, "value") else str(source.connector_type),
                            "host": source.host, "port": source.port,
                            "database_name": source.database_name,
                            "username": source.username, "base_url": source.base_url, "password": "",
                        }
                        try:
                            async with (await get_pg_pool()).acquire() as _cs_early:
                                _sec_early = await _cs_early.fetchrow(
                                    "SELECT secret_value FROM connection_secrets WHERE source_id=$1 AND secret_key='password'",
                                    source_id
                                )
                                if _sec_early:
                                    _src_dict_early["password"] = _sec_early["secret_value"]
                        except Exception:
                            pass
                        import time as _time_early
                        _exec_early = await _tool_execute_sql(_direct_sql_early, _src_dict_early, dialect)
                        if _exec_early.get("success") or _exec_early.get("row_count", 0) == 0:
                            yield await _sse({"type": "sql", "content": _direct_sql_early, "method": "direct_sql_early"})
                            # ── Sauvegarder pour correction interactive ───────
                            try:
                                _sv_pool = await get_pg_pool()
                                async with _sv_pool.acquire() as _sv_conn:
                                    await _sv_conn.execute("""
                                        UPDATE nlu_query_log
                                        SET    sql_generated = $1
                                        WHERE  source_id     = $2
                                          AND  question      = $3
                                          AND  sql_generated IS NULL
                                    """, _direct_sql_early, UUID(source_id_str), question)
                                    _chk = await _sv_conn.fetchval("""
                                        SELECT COUNT(*) FROM nlu_query_log
                                        WHERE source_id=$1 AND question=$2 AND sql_generated=$3
                                    """, UUID(source_id_str), question, _direct_sql_early)
                                    if not _chk:
                                        await _sv_conn.execute("""
                                            INSERT INTO nlu_query_log
                                                (source_id, question, intent, confidence,
                                                 sql_generated, created_at)
                                            VALUES ($1,$2,$3,$4,$5,NOW())
                                            ON CONFLICT DO NOTHING
                                        """, UUID(source_id_str), question,
                                           "direct_sql", 1.0, _direct_sql_early)
                            except Exception as _sv_err:
                                logger.debug(f"[DirectSQL] save error: {_sv_err}")
                            ms = int((_time.time() - t0) * 1000)
                            yield await _sse({"type": "done", "duration_ms": ms})
                            yield "data: [DONE]\n\n"
                            return
                except Exception as _de:
                    logger.debug(f"[DirectSQL Early] Ignoré: {_de}")

            # ── Sprint 10 : questions complexes → toujours via Orchestrateur ──
            # _is_complex_question() détecte multi-dims, agrégations croisées
            # Ces questions doivent passer par l'Agent Precision, pas par CRAG direct
            from .orchestrator import _is_complex_question as _orch_is_complex
            _is_complex_q = _orch_is_complex(_q_lower_stream)

            _use_agent_stream = (
                not _is_dashboard_question and
                (_is_complex_q or any(kw in _q_lower_stream for kw in _agent_kw_stream))
            )

            if _use_agent_stream:
                # ── Sprint 9 : Orchestrateur Multi-Agent ─────────────────────
                logger.info(f"[Orchestrator] Démarrage dans /chat/stream — question='{question[:60]}'")
                try:
                    from .orchestrator import run_orchestrator, orchestrator_result_to_dict
                    _pool_agent = await get_pg_pool()
                    _src_dict_stream = {
                        "id":             source_id_str,
                        "connector_type": source.connector_type.value if hasattr(source.connector_type, "value") else str(source.connector_type),
                        "host":           source.host,
                        "port":           source.port,
                        "database_name":  source.database_name,
                        "username":       source.username,
                        "base_url":       source.base_url,
                        "password":       "",
                    }
                    try:
                        async with _pool_agent.acquire() as _cs:
                            _sec = await _cs.fetchrow(
                                "SELECT secret_value FROM connection_secrets WHERE source_id=$1 AND secret_key='password'",
                                source_id
                            )
                            if _sec:
                                _src_dict_stream["password"] = _sec["secret_value"]
                    except Exception:
                        pass

                    yield await _sse({"type": "thinking", "content": "🤖 Orchestrateur Multi-Agent — analyse en cours…"})
                    _orch_res = await run_orchestrator(
                        question=question, source_id=source_id,
                        pg_pool=_pool_agent, source_dict=_src_dict_stream, dialect=dialect,
                    )
                    if _orch_res.success and _orch_res.sql:
                        logger.info(
                            f"[Orchestrator] Succès stream — "                            f"{_orch_res.iterations} itérations | "                            f"agent={_orch_res.agent_type.value} | "                            f"méthode={_orch_res.method}"
                        )
                        # Si multi-query : envoyer tous les SQLs
                        if _orch_res.sqls and len(_orch_res.sqls) > 1:
                            yield await _sse({
                                "type": "sql",
                                "content": _orch_res.sql,
                                "method": _orch_res.method,
                                "sqls": _orch_res.sqls,
                                "sub_queries": [
                                    {"text": sq.text, "sql": sq.sql, "success": sq.success}
                                    for sq in _orch_res.sub_queries
                                ],
                            })
                        else:
                            yield await _sse({"type": "sql", "content": _orch_res.sql, "method": _orch_res.method})
                        if _orch_res.warnings:
                            for w in _orch_res.warnings:
                                yield await _sse({"type": "warning", "content": w})
                        # ── Sauvegarde SQL pour correction interactive ────────
                        if _orch_res.sql:
                            try:
                                _sv2 = await get_pg_pool()
                                async with _sv2.acquire() as _sv2c:
                                    await _sv2c.execute("""
                                        UPDATE nlu_query_log
                                        SET    sql_generated = $1
                                        WHERE  source_id=$2 AND question=$3
                                          AND  sql_generated IS NULL
                                    """, _orch_res.sql, UUID(source_id_str), question)
                                    _chk2 = await _sv2c.fetchval(
                                        "SELECT COUNT(*) FROM nlu_query_log WHERE source_id=$1 AND question=$2 AND sql_generated=$3",
                                        UUID(source_id_str), question, _orch_res.sql)
                                    if not _chk2:
                                        await _sv2c.execute("""
                                            INSERT INTO nlu_query_log
                                                (source_id,question,intent,confidence,sql_generated,created_at)
                                            VALUES ($1,$2,$3,$4,$5,NOW()) ON CONFLICT DO NOTHING
                                        """, UUID(source_id_str), question,
                                           _orch_res.method, 1.0, _orch_res.sql)
                            except Exception as _sv2e:
                                logger.debug(f"[Orch] sql save error: {_sv2e}")
                        ms = int((_time.time() - t0) * 1000)
                        yield await _sse({"type": "done", "duration_ms": ms})
                        yield "data: [DONE]\n\n"
                        return
                    else:
                        logger.info(f"[AgentRAG] Pas de résultat — fallback LLM")
                except Exception as _ae:
                    logger.warning(f"[AgentRAG] Échec stream, fallback LLM : {_ae}", exc_info=True)

            # ── Sprint 9 fix : intercepter show_dashboard AVANT generate_sql_stream ──
            if slots and slots.intent == Intent.SHOW_DASHBOARD:
                try:
                    from .dashboard_engine import get_dashboard_generator
                    logger.info(f"[Chat/Dashboard] Génération dashboard — question='{question[:60]}'")
                    generator = get_dashboard_generator()
                    spec = await generator.generate(
                        question          = question,
                        slots             = slots,
                        schema            = schema,
                        source_id         = source_id_str,
                        pg_pool           = await get_pg_pool(),
                        redis             = await get_redis(),
                        connector_factory = ConnectorFactory,
                    )
                    spec_dict = spec.to_dict()
                    n_widgets = len(spec_dict.get("widgets", []))
                    logger.info(f"[Chat/Dashboard] {n_widgets} widgets générés")
                    import json as _json_db
                    dashboard_json = _json_db.dumps(spec_dict, default=str)
                    yield await _sse({"type": "sql", "content": f"```dashboard\n{dashboard_json}\n```", "method": "dashboard_engine"})
                    ms = int((_time.time() - t0) * 1000)
                    yield await _sse({"type": "done", "duration_ms": ms})
                    yield "data: [DONE]\n\n"
                    return
                except Exception as _de:
                    logger.error(f"[Chat/Dashboard] Erreur: {_de}", exc_info=True)
                    yield await _sse({"type": "token", "content": f"Erreur génération dashboard : {_de}"})
                    yield await _sse({"type": "done", "duration_ms": 0})
                    yield "data: [DONE]\n\n"
                    return

            # ── Phase 2 : Streaming LLM ──────────────────────────────
            # Stratégie : on affiche les tokens bruts du LLM en temps réel
            # SAUF les tokens markdown (```sql, ```) qui polluent l'affichage.
            # ── Sprint 13 : A/B Testing — tirage 50/50 Prompt A vs Prompt B ──
            import random as _random
            _ab_variant = "B" if _random.random() < 0.5 else "A"
            _ab_t0 = _time.time()
            _ab_few_shot = []
            if _ab_variant == "B":
                try:
                    from .llm_engine import _get_few_shot_examples
                    _ab_pool = await get_pg_pool()
                    _ab_few_shot = await _get_few_shot_examples(_ab_pool, source_id_str, limit=5)
                    logger.info(f"[AB] Variant B — {len(_ab_few_shot)} few-shot exemples chargés")
                except Exception as _abe:
                    logger.debug(f"[AB] Few-shot fetch error: {_abe}")
                    _ab_variant = "A"  # fallback A si erreur
            else:
                logger.info("[AB] Variant A — prompt baseline")

            # Le SQL final nettoyé (type:sql) remplace tout à la fin.
            pool = await get_pg_pool()
            collected_sql = None
            method = "llm"
            raw_token_buf = []   # accumule pour vérifier si c'est un bloc markdown
            in_code_block = False
            skip_next_lang = False

            async for event in generate_sql_stream(
                question    = question,
                schema      = schema,
                dialect     = dialect,
                table_names = table_names_hint,
                pg_pool     = pool,
                source_id   = source_id_str,
                ab_variant  = _ab_variant,
                few_shot_examples = _ab_few_shot,
            ):
                etype = event.get("type")
                if etype == "thinking":
                    yield await _sse(event)
                elif etype == "conceptual":
                    # ── Sprint 7C : Routing → réponse texte LLM ──────────────
                    try:
                        import httpx as _hx
                        _cp = (
                            f"Tu es un expert ERP et systèmes d'information. "
                            f"Réponds en français à cette question de manière concise.\n"
                            f"Question : {event.get('content', question)}\n\n"
                            f"Réponds directement sans générer de SQL. 2-4 phrases maximum."
                        )
                        _cr = _hx.post(
                            f"{OLLAMA_HOST}/api/generate",
                            json={"model": OLLAMA_MODEL, "prompt": _cp, "stream": False,
                                  "options": {"temperature": 0.3, "num_predict": 300}},
                            timeout=30,
                        )
                        _ans = _cr.json().get("response", "").strip() if _cr.status_code == 200 else "Question conceptuelle non supportée."
                    except Exception:
                        _ans = "Cette question est de nature conceptuelle. Précisez si vous souhaitez une requête SQL."
                    words = _ans.split(" ")
                    for i, word in enumerate(words):
                        yield await _sse({"type": "token", "content": word + (" " if i < len(words)-1 else "")})
                        await asyncio.sleep(0.015)
                    ms = int((_time.time() - t0) * 1000)
                    yield await _sse({"type": "done", "duration_ms": ms})
                    yield "data: [DONE]\n\n"
                    return
                elif etype == "token":
                    tok = event.get("content", "")
                    raw_token_buf.append(tok)
                    # Filtre les tokens markdown parasites
                    tok_stripped = tok.strip()
                    if tok_stripped == "```":
                        in_code_block = not in_code_block
                        if in_code_block:
                            skip_next_lang = True  # prochaine ligne = "sql"
                        continue
                    if skip_next_lang and tok_stripped.lower() in ("sql", ""):
                        skip_next_lang = False
                        continue
                    skip_next_lang = False
                    # Envoie le token au client
                    yield await _sse({"type": "token", "content": tok})
                elif etype == "sql":
                    collected_sql = event.get("content", "")
                    method  = event.get("method", "llm")
                elif etype == "error":
                    # LLM indisponible → fallback vers _build_chat_answer
                    yield await _sse({"type": "thinking", "content": "llm_offline"})
                    answer = await _build_chat_answer(source_id_str, question)
                    words = answer.split(" ")
                    for i, word in enumerate(words):
                        yield await _sse({"type": "token", "content": word + (" " if i < len(words)-1 else "")})
                        await asyncio.sleep(0.015)
                    ms = int((_time.time() - t0) * 1000)
                    yield await _sse({"type": "done", "duration_ms": ms})
                    yield "data: [DONE]\n\n"
                    return
                elif etype == "done":
                    pass  # On envoie notre propre done après

            # SQL final — on remplace tout ce qui a été affiché par le SQL propre
            ms = int((_time.time() - t0) * 1000)
            if collected_sql:
                badge = "LLM+RAG" if "rag" in (method or "") else "LLM"
                # Efface les tokens bruts et affiche le SQL nettoyé final
                yield await _sse({"type": "replace", "content": f"```sql\n{collected_sql}\n```"})

            # ── Sprint 13 : Enregistrement résultat A/B ──────────────────────
            try:
                _ab_valid = bool(collected_sql and len(collected_sql) > 10)
                _ab_score = float(_ab_valid)
                _ab_duration = int((_time.time() - _ab_t0) * 1000)
                _ab_pg = await get_pg_pool()
                async with _ab_pg.acquire() as _ab_conn:
                    await _ab_conn.execute("""
                        INSERT INTO ab_test_results
                            (question, source_id, prompt_variant, sql_generated,
                             sql_valid, score, duration_ms, few_shot_count, created_at)
                        VALUES ($1,$2,$3,$4,$5,$6,$7,$8,NOW())
                    """,
                        question[:500],
                        UUID(source_id_str) if source_id_str else None,
                        _ab_variant,
                        collected_sql,
                        _ab_valid,
                        _ab_score,
                        _ab_duration,
                        len(_ab_few_shot),
                    )
                logger.info(f"[AB] Variant={_ab_variant} valid={_ab_valid} score={_ab_score} duration={_ab_duration}ms")
            except Exception as _ab_err:
                logger.debug(f"[AB] Enregistrement ignoré: {_ab_err}")

            # ── Sauvegarde sql_generated dans nlu_query_log pour la correction ─
            if collected_sql:
                try:
                    _save_pool = await get_pg_pool()
                    async with _save_pool.acquire() as _save_conn:
                        # Mise à jour de la ligne la plus récente pour cette question
                        _rows_updated = await _save_conn.execute("""
                            UPDATE nlu_query_log
                            SET    sql_generated = $1
                            WHERE  id = (
                                SELECT id FROM nlu_query_log
                                WHERE  source_id = $2
                                  AND  question  = $3
                                ORDER  BY created_at DESC
                                LIMIT  1
                            )
                        """, collected_sql, UUID(source_id_str), question)
                        # Si aucune ligne → INSERT direct
                        if _rows_updated == "UPDATE 0":
                            await _save_conn.execute("""
                                INSERT INTO nlu_query_log
                                    (source_id, question, intent, confidence,
                                     sql_generated, created_at)
                                VALUES ($1, $2, $3, $4, $5, NOW())
                            """, UUID(source_id_str), question,
                               slots.intent if slots else "generate_aggregate",
                               float(slots.confidence) if slots and slots.confidence else 0.9,
                               collected_sql)
                    logger.info(f"[Stream] sql_generated sauvegardé — {len(collected_sql)} chars")
                except Exception as _save_err:
                    logger.warning(f"[Stream] sql_generated save error: {_save_err}")

            yield await _sse({"type": "done", "duration_ms": ms})
            yield "data: [DONE]\n\n"
            return

        # ── Phase 2 alt : Template SQLGenerator (requête simple) ──────
        yield await _sse({"type": "thinking", "content": "sql_template"})
        try:
            answer = await _build_chat_answer(source_id_str, question)
        except Exception as _e:
            answer = f"Erreur : {_e}"

        # ── Extraire et sauvegarder le SQL depuis la réponse ─────────────
        try:
            import re as _re_sql_save
            _sql_match = _re_sql_save.search(r'```sql\s*([\s\S]+?)\s*```', answer)
            if _sql_match:
                _extracted_sql = _sql_match.group(1).strip()
                _sv3_pool = await get_pg_pool()
                async with _sv3_pool.acquire() as _sv3_conn:
                    _upd3 = await _sv3_conn.execute("""
                        UPDATE nlu_query_log
                        SET    sql_generated = $1
                        WHERE  id = (
                            SELECT id FROM nlu_query_log
                            WHERE  source_id = $2 AND question = $3
                            ORDER  BY created_at DESC LIMIT 1
                        )
                    """, _extracted_sql, UUID(source_id_str), question)
                    if _upd3 == "UPDATE 0":
                        await _sv3_conn.execute("""
                            INSERT INTO nlu_query_log
                                (source_id, question, intent, confidence,
                                 sql_generated, created_at)
                            VALUES ($1,$2,$3,$4,$5,NOW())
                        """, UUID(source_id_str), question,
                           slots.intent if slots else "generate_aggregate",
                           float(slots.confidence) if slots and slots.confidence else 0.9,
                           _extracted_sql)
                logger.info(f"[Phase2] sql_generated sauvegardé — {len(_extracted_sql)} chars")
        except Exception as _sv3e:
            logger.debug(f"[Phase2] sql save error: {_sv3e}")

        # Stream la réponse template mot par mot
        words = answer.split(" ")
        for i, word in enumerate(words):
            yield await _sse({"type": "token", "content": word + (" " if i < len(words)-1 else "")})
            await asyncio.sleep(0.012)

        ms = int((_time.time() - t0) * 1000)
        yield await _sse({"type": "done", "duration_ms": ms})
        yield "data: [DONE]\n\n"

    return StreamingResponse(
        _stream_pipeline(),
        media_type="text/event-stream",
        headers={
            "Cache-Control":  "no-cache",
            "X-Accel-Buffering": "no",  # Désactive le buffering Nginx
        },
    )

# ══════════════════════════════════════════════════════════════
# VALIDATION EXPERTE NLU — Sprint 13
# ══════════════════════════════════════════════════════════════

@app.get("/admin/nlu/low-confidence", tags=["Admin"])
async def get_low_confidence_queries(
    min_conf: float = Query(0.0, ge=0.0, le=1.0),
    max_conf: float = Query(0.7, ge=0.0, le=1.0),
    limit:    int   = Query(50, ge=1, le=200),
):
    """
    Retourne les questions avec une faible confiance NLU.
    L'expert admin peut les corriger pour améliorer FastText.
    """
    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT
                    id, question, intent, confidence,
                    sql_generated IS NOT NULL AS has_sql,
                    created_at
                FROM   nlu_query_log
                WHERE  confidence >= $1
                  AND  confidence <  $2
                  AND  question   != ''
                  AND  question   NOT LIKE '[CORRECTION]%'
                ORDER  BY created_at DESC
                LIMIT  $3
            """, min_conf, max_conf, limit)

        return {
            "total": len(rows),
            "queries": [
                {
                    "id":         r["id"],
                    "question":   r["question"],
                    "intent":     r["intent"],
                    "confidence": round(float(r["confidence"] or 0), 3),
                    "has_sql":    r["has_sql"],
                    "created_at": r["created_at"].isoformat() if r["created_at"] else None,
                }
                for r in rows
            ]
        }
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/admin/nlu/correct-intent", tags=["Admin"])
async def correct_intent(body: dict):
    """
    Permet à l'expert admin de corriger l'intent d'une question.
    La correction est ajoutée au dataset FastText et déclenche un retrain.

    Body: { question: str, correct_intent: str, question_id?: int }
    """
    question       = body.get("question", "").strip()
    correct_intent = body.get("correct_intent", "").strip()
    question_id    = body.get("question_id")

    if not question or not correct_intent:
        raise HTTPException(400, "question et correct_intent requis")

    # Vérifier que l'intent est valide
    VALID_INTENTS = [
        "generate_aggregate", "list_entities", "count_entities",
        "show_dashboard", "list_fields", "describe_entity",
        "get_relations", "generate_join", "generate_sql",
        "profile_entity", "search_entity", "find_path",
        "greeting", "help", "correct_sql", "forecast",
    ]
    if correct_intent not in VALID_INTENTS:
        raise HTTPException(400, f"Intent invalide. Valides: {VALID_INTENTS}")

    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            # 1. Mettre à jour nlu_query_log si question_id fourni (UUID string)
            if question_id:
                try:
                    from uuid import UUID as _UUID
                    _qid = _UUID(str(question_id))
                    await conn.execute("""
                        UPDATE nlu_query_log
                        SET    intent     = $1,
                               confidence = 1.0
                        WHERE  id         = $2
                    """, correct_intent, _qid)
                except Exception:
                    pass  # question_id invalide → on continue quand même

            # 2. Ajouter dans chat_feedback comme exemple validé
            await conn.execute("""
                INSERT INTO chat_feedback
                    (conversation_id, source_id, message_id,
                     feedback_type, question, intent,
                     used_for_training, created_at)
                VALUES
                    ('admin_correction', NULL, 'expert_' || $1::text,
                     'like', $2, $3,
                     FALSE, NOW())
                ON CONFLICT DO NOTHING
            """, question_id or 0, question, correct_intent)

        # 3. Retrain FastText immédiat avec cet exemple
        result = retrain_fasttext_with_feedback([(question, correct_intent)])

        logger.info(
            f"[ExpertValidation] Correction: '{question[:40]}' → {correct_intent} "
            f"| retrain: +{result.get('examples_added',0)} exemples"
        )

        return {
            "status":          "corrected",
            "question":        question,
            "correct_intent":  correct_intent,
            "retrain_result":  result,
            "message":         f"Intent corrigé et FastText réentraîné avec cet exemple."
        }
    except Exception as e:
        raise HTTPException(500, str(e))


@app.get("/admin/nlu/stats", tags=["Admin"])
async def get_nlu_validation_stats():
    """Statistiques globales de validation NLU pour le dashboard admin."""
    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            # Distribution des confiances
            dist = await conn.fetch("""
                SELECT
                    CASE
                        WHEN confidence >= 0.9  THEN 'high (>=0.9)'
                        WHEN confidence >= 0.7  THEN 'medium (0.7-0.9)'
                        WHEN confidence >= 0.5  THEN 'low (0.5-0.7)'
                        ELSE                         'very_low (<0.5)'
                    END AS bucket,
                    COUNT(*) AS cnt
                FROM   nlu_query_log
                WHERE  confidence IS NOT NULL
                GROUP  BY bucket
                ORDER  BY bucket
            """)

            # Top intents
            intents = await conn.fetch("""
                SELECT intent, COUNT(*) AS cnt,
                       ROUND(AVG(confidence)::NUMERIC, 3) AS avg_conf
                FROM   nlu_query_log
                WHERE  intent IS NOT NULL AND intent != ''
                GROUP  BY intent
                ORDER  BY cnt DESC
                LIMIT  15
            """)

            # Total corrections faites
            corrections = await conn.fetchval("""
                SELECT COUNT(*) FROM chat_feedback
                WHERE  conversation_id = 'admin_correction'
            """)

        return {
            "confidence_distribution": [{"bucket": r["bucket"], "count": r["cnt"]} for r in dist],
            "top_intents":             [{"intent": r["intent"], "count": r["cnt"], "avg_conf": float(r["avg_conf"] or 0)} for r in intents],
            "total_corrections":       corrections or 0,
        }
    except Exception as e:
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════════════════════
# A/B TESTING — Sprint 13
# ══════════════════════════════════════════════════════════════

@app.get("/ab-testing/stats", tags=["AB Testing"])
async def ab_testing_stats(source_id: Optional[str] = None, days: int = 30):
    """
    Statistiques A/B testing des prompts LLM.
    Compare Prompt A (baseline) vs Prompt B (few-shot enrichi).
    """
    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            # Stats globales par variant
            rows = await conn.fetch("""
                SELECT
                    prompt_variant,
                    COUNT(*)                                            AS total_tests,
                    COUNT(*) FILTER (WHERE sql_valid = TRUE)           AS valid_sql,
                    COUNT(*) FILTER (WHERE has_results = TRUE)         AS has_results,
                    ROUND(AVG(score)::NUMERIC, 3)                      AS avg_score,
                    ROUND(AVG(duration_ms)::NUMERIC, 0)                AS avg_duration_ms,
                    ROUND((COUNT(*) FILTER (WHERE sql_valid = TRUE)::FLOAT
                          / NULLIF(COUNT(*),0)*100)::NUMERIC, 1)       AS valid_pct,
                    AVG(few_shot_count)                                AS avg_few_shot
                FROM ab_test_results
                WHERE created_at >= NOW() - ($1 || ' days')::INTERVAL
                  AND ($2::UUID IS NULL OR source_id = $2::UUID)
                GROUP BY prompt_variant
                ORDER BY prompt_variant
            """, str(days), source_id)

            # Total général
            total = await conn.fetchval(
                "SELECT COUNT(*) FROM ab_test_results WHERE created_at >= NOW() - ($1 || ' days')::INTERVAL",
                str(days)
            )

            # Évolution quotidienne
            daily = await conn.fetch("""
                SELECT
                    DATE(created_at)  AS day,
                    prompt_variant,
                    COUNT(*)          AS tests,
                    ROUND(AVG(score)::NUMERIC, 2) AS avg_score
                FROM ab_test_results
                WHERE created_at >= NOW() - ($1 || ' days')::INTERVAL
                GROUP BY DATE(created_at), prompt_variant
                ORDER BY day DESC
                LIMIT 60
            """, str(days))

        stats = {}
        for r in rows:
            stats[r["prompt_variant"]] = {
                "total_tests":    r["total_tests"],
                "valid_sql":      r["valid_sql"],
                "has_results":    r["has_results"],
                "avg_score":      float(r["avg_score"] or 0),
                "avg_duration_ms": float(r["avg_duration_ms"] or 0),
                "valid_pct":      float(r["valid_pct"] or 0),
                "avg_few_shot":   float(r["avg_few_shot"] or 0),
            }

        # Déterminer le gagnant actuel
        winner = None
        if "A" in stats and "B" in stats:
            winner = "B" if stats["B"]["avg_score"] > stats["A"]["avg_score"] else "A"
        elif stats:
            winner = list(stats.keys())[0]

        return {
            "period_days":  days,
            "total_tests":  total,
            "by_variant":   stats,
            "current_winner": winner,
            "daily_evolution": [
                {
                    "day":          str(r["day"]),
                    "variant":      r["prompt_variant"],
                    "tests":        r["tests"],
                    "avg_score":    float(r["avg_score"] or 0),
                }
                for r in daily
            ],
            "recommendation": (
                f"Prompt {winner} est plus performant — continuer à l'utiliser."
                if winner else "Pas assez de données pour recommander un variant."
            ),
        }
    except Exception as e:
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════════════════════
# LLM ENGINE — §2.3.3
# ══════════════════════════════════════════════════════════════

class DictationRequest(BaseModel):
    text:      str
    source_id: str = ""

@app.post("/llm/extract-entities", tags=["LLM"])
async def extract_entities_from_dictation(req: DictationRequest):
    """
    Extrait les entités structurées depuis une dictée vocale.
    Ex: "Crée une commande pour le client Dupont, 50 unités du produit ABC-123"
    → {action: "create", entity: "order", client: "Dupont", qty: 50, product: "ABC-123"}
    """
    try:
        from .llm_engine import generate_sql_with_llm
        prompt = f"""Extract structured entities from this voice dictation in JSON format.
Text: "{req.text}"
Return ONLY a JSON object with these fields if present:
- action: create/update/delete/search
- entity: order/product/customer/invoice
- client/customer name
- quantity (number)
- product code or name
- amount (number)
- date
- any other relevant fields
Return only valid JSON, no explanation."""

        import httpx, json as _json, os
        OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://host.docker.internal:11434")
        response = httpx.post(
            f"{OLLAMA_HOST}/api/generate",
            json={"model": "qwen2.5-coder:3b", "prompt": prompt,
                  "stream": False, "options": {"temperature": 0.1}},
            timeout=30,
        )
        raw = response.json().get("response", "").strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        entities = _json.loads(raw)
        return {"entities": entities, "text": req.text}
    except Exception as e:
        return {"entities": None, "error": str(e), "text": req.text}

@app.get("/llm/status", tags=["LLM"])
async def llm_status():
    """Vérifie si Ollama est disponible et retourne les modèles."""
    try:
        from .llm_engine import check_ollama_available
        return check_ollama_available()
    except Exception as e:
        return {"available": False, "error": str(e)}

class LLMSQLRequest(BaseModel):
    question:  str
    source_id: str

@app.post("/llm/generate-sql", tags=["LLM"])
async def llm_generate_sql(req: LLMSQLRequest):
    """Génère du SQL directement via LLM (pour tests)."""
    try:
        from .llm_engine import generate_sql_with_llm
        source_id = UUID(req.source_id)
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT se.name, ef.name AS field
                FROM source_entities se
                JOIN entity_fields ef ON ef.entity_id = se.id
                WHERE se.source_id = $1 AND se.is_visible = TRUE
                LIMIT 2000
            """, source_id)
        schema = {}
        for r in rows:
            schema.setdefault(r["name"], []).append(r["field"])
        # Sprint 7A : passage de pg_pool et source_id pour le RAG
        result = await generate_sql_with_llm(
            question    = req.question,
            schema      = schema,
            dialect     = "mssql",
            table_names = None,
            pg_pool     = pool,
            source_id   = req.source_id,
        )
        return result
    except Exception as e:
        raise HTTPException(500, str(e))

# ══════════════════════════════════════════════════════════════
# TEXT-TO-SPEECH & STT ENDPOINTS
# ══════════════════════════════════════════════════════════════

class TTSRequest(BaseModel):
    text:   str
    voice:  str = "fr_FR-upmc-medium"
    speed:  float = 1.0
    format: str = "wav"

@app.get("/tts/status", tags=["Voice"])
async def tts_status():
    """Vérifie si Piper TTS est disponible."""
    try:
        from .voice_engine import check_piper_available
        return check_piper_available()
    except Exception as e:
        return {"available": False, "error": str(e)}

@app.post("/tts", tags=["Voice"])
async def text_to_speech(req: TTSRequest):
    """
    Synthétise du texte en audio WAV via Piper TTS (on-premise).
    Supporte SSML basique, vitesse ajustable, voix multiples.
    """
    try:
        from .voice_engine import get_tts_engine
        from fastapi.responses import Response

        if not req.text:
            raise HTTPException(400, "Texte vide")

        tts = get_tts_engine()

        # Préparation vocale dans le endpoint
        text_for_tts = req.text
        original_len = len(text_for_tts)

        # Simplifie Markdown, emojis, backticks
        text_for_tts = tts.simplify_for_voice(text_for_tts)

        # Si le simplificateur a beaucoup coupé → texte original était long → ajoute conclusion
        if original_len > 400 and len(text_for_tts) < original_len * 0.5:
            text_for_tts = text_for_tts.rstrip(".:") + ". La réponse complète est affichée à l'écran."

        # Résumé vocal si encore trop long
        if len(text_for_tts) > 600:
            text_for_tts = tts.vocal_summary(text_for_tts)

        # Pronunciation hints
        text_for_tts = tts.pronunciation_hints(text_for_tts)

        logger.info(f"[TTS] Préparé: {original_len} → {len(text_for_tts)} chars : {repr(text_for_tts[:80])}")

        wav_bytes = tts.synthesize(
            text  = text_for_tts,
            voice = req.voice,
            speed = req.speed,
        )

        if not wav_bytes:
            raise HTTPException(500, "Synthèse audio échouée")

        return Response(
            content    = wav_bytes,
            media_type = "audio/wav",
            headers    = {"Content-Disposition": "inline; filename=speech.wav"},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"TTS error: {e}")
        raise HTTPException(500, str(e))

@app.post("/tts/ssml", tags=["Voice"])
async def text_to_speech_ssml(req: TTSRequest):
    """Synthétise du texte SSML en audio WAV."""
    return await text_to_speech(req)

@app.post("/tts/stream", tags=["Voice"])
async def text_to_speech_stream(req: TTSRequest):
    """
    Streaming TTS — génère et envoie l'audio phrase par phrase.
    Latence réduite : le navigateur joue dès que la première phrase est prête.
    """
    import re
    from fastapi.responses import StreamingResponse

    try:
        from .voice_engine import get_tts_engine

        tts = get_tts_engine()
        text = req.text[:800]

        # Découpe en phrases
        sentences = re.split(r'(?<=[.!?]) +', text)
        sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 2]
        if not sentences:
            sentences = [text]

        async def generate():
            for sentence in sentences:
                try:
                    wav = tts.synthesize(
                        text  = sentence,
                        voice = req.voice,
                        speed = req.speed,
                    )
                    if wav and len(wav) > 44:  # > header WAV vide
                        yield wav
                except Exception as e:
                    logger.warning(f"[TTS/stream] Phrase skip: {e}")
                    continue

        return StreamingResponse(
            generate(),
            media_type="audio/wav",
            headers={"X-Accel-Buffering": "no"},
        )

    except Exception as e:
        logger.error(f"TTS stream error: {e}")
        raise HTTPException(500, str(e))

@app.get("/stt/status", tags=["Voice"])
async def stt_status():
    """Vérifie si Whisper STT et Vosk sont disponibles."""
    try:
        from .voice_engine import check_whisper_available, check_vosk_available
        whisper = check_whisper_available()
        vosk    = check_vosk_available()
        return {**whisper, "vosk": vosk}
    except Exception as e:
        return {"available": False, "error": str(e)}

@app.post("/stt/vosk", tags=["Voice"])
async def speech_to_text_vosk(file: UploadFile = File(...)):
    """
    Transcrit un fichier audio via Vosk (alternative légère à Whisper).
    Plus rapide mais moins précis que Whisper.
    """
    try:
        from .voice_engine import get_vosk_engine
        audio_bytes = await file.read()
        if not audio_bytes:
            raise HTTPException(400, "Fichier audio vide")
        vosk = get_vosk_engine()
        result = vosk.transcribe_audio_file(audio_bytes, filename=file.filename or "audio.webm")
        if result.get("error"):
            raise HTTPException(500, result["error"])
        return result
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Vosk STT error: {e}")
        raise HTTPException(500, str(e))

@app.post("/stt", tags=["Voice"])
async def speech_to_text(file: UploadFile = File(...)):
    """
    Transcrit un fichier audio en texte via Whisper.
    Supporte : webm, mp3, wav, ogg, m4a
    """
    try:
        from .voice_engine import get_stt_engine
        audio_bytes = await file.read()
        if not audio_bytes:
            raise HTTPException(400, "Fichier audio vide")

        stt    = get_stt_engine()

        # VAD — détection automatique silence/parole
        # Note: webm nécessite conversion PCM pour VAD
        # On applique VAD seulement sur WAV/PCM
        fname = file.filename or "audio.webm"
        if fname.endswith(".wav") or fname.endswith(".pcm"):
            try:
                from .voice_engine import get_vad
                vad = get_vad()
                audio_bytes = vad.filter_silence(audio_bytes)
                logger.info(f"[STT] VAD appliqué — {len(audio_bytes)} bytes après filtrage")
            except Exception as e:
                logger.warning(f"[STT] VAD skip: {e}")

        result = stt.transcribe(audio_bytes, filename=fname)

        if result.get("error"):
            raise HTTPException(500, result["error"])

        return {
            "text":        result["text"],
            "raw_text":    result.get("raw_text", ""),
            "language":    result.get("language", "fr"),
            "duration_ms": result.get("duration_ms", 0),
            "command":     result.get("command"),
            "model":       result.get("model", "base"),
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"STT error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ══════════════════════════════════════════════════════════════
# VOICE EXTENSIONS — §2.3.2 (custom vocab, voice messages)
# ══════════════════════════════════════════════════════════════

class VocabularyRequest(BaseModel):
    terms:     List[str]
    source_id: Optional[str] = ""

class VoiceMessageIn(BaseModel):
    text:      str
    source_id: Optional[str] = ""
    user_id:   Optional[str] = "default"

@app.post("/stt/vocabulary", tags=["Voice"])
async def add_custom_vocabulary(req: VocabularyRequest):
    """
    Ajoute des termes métier custom au vocabulaire Whisper.
    Ces termes sont injectés dans initial_prompt pour améliorer la précision STT.
    §2.3.2A — Custom vocabulary persistant à runtime.
    """
    try:
        from .voice_engine import extend_vocabulary, get_full_vocabulary
        if not req.terms:
            raise HTTPException(400, "Liste de termes vide")
        cleaned = [t.strip() for t in req.terms if t.strip() and len(t.strip()) <= 60]
        if not cleaned:
            raise HTTPException(400, "Aucun terme valide (max 60 car. par terme)")
        total = extend_vocabulary(cleaned)
        return {
            "added":       len(cleaned),
            "total":       total,
            "vocabulary":  get_full_vocabulary(),
            "message":     f"{len(cleaned)} terme(s) ajouté(s) au vocabulaire Whisper",
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Vocabulary error: {e}")
        raise HTTPException(500, str(e))

@app.get("/stt/vocabulary", tags=["Voice"])
async def get_vocabulary():
    """Retourne le vocabulaire STT complet (base + custom)."""
    try:
        from .voice_engine import get_full_vocabulary, BUSINESS_VOCABULARY, _CUSTOM_VOCABULARY_EXTRA
        return {
            "base":   BUSINESS_VOCABULARY,
            "custom": _CUSTOM_VOCABULARY_EXTRA,
            "total":  len(BUSINESS_VOCABULARY) + len(_CUSTOM_VOCABULARY_EXTRA),
        }
    except Exception as e:
        return {"base": [], "custom": [], "total": 0, "error": str(e)}

@app.get("/tts/voices", tags=["Voice"])
async def list_voices():
    """Liste les voix TTS disponibles."""
    try:
        from .voice_engine import get_tts_engine, check_piper_available
        info = check_piper_available()
        return {
            "available": info.get("available", False),
            "voices":    info.get("voices", []),
            "default":   info.get("default", "fr_FR-upmc-medium"),
            "aliases": {
                "femme":  "fr_FR-upmc-medium",
                "homme":  "fr_FR-gilles-low",
                "female": "fr_FR-upmc-medium",
                "male":   "fr_FR-gilles-low",
            },
        }
    except Exception as e:
        return {"available": False, "voices": [], "error": str(e)}

@app.websocket("/stt/stream")
async def stt_stream_websocket(ws: WebSocket):
    """
    Streaming STT via WebSocket + Vosk — transcription en temps réel.
    §2.3.2A — Streaming recognition : retourne les mots au fur et à mesure.

    Protocole :
      Client → serveur : chunks PCM 16-bit mono 16kHz (bytes bruts)
      Client → serveur : string "END" pour signaler fin de parole
      Serveur → client : JSON {"partial": "...", "final": false}
      Serveur → client : JSON {"text": "...", "final": true, "command": ...}
    """
    await ws.accept()
    logger.info("[WS/STT] Connexion streaming STT ouverte")

    try:
        from vosk import KaldiRecognizer
        from .voice_engine import get_vosk_engine, normalize_voice_text, detect_voice_command
        import json as _json

        vosk = get_vosk_engine()
        model = vosk._load_model()
        rec = KaldiRecognizer(model, 16000)
        rec.SetWords(True)

        partial_text = ""
        full_results = []

        while True:
            try:
                data = await ws.receive()
            except WebSocketDisconnect:
                break

            # "END" signal — envoie résultat final
            if data.get("type") == "websocket.receive" and data.get("text") == "END":
                final_raw = _json.loads(rec.FinalResult()).get("text", "")
                if final_raw:
                    full_results.append(final_raw)
                full = " ".join(full_results).strip()
                normalized = normalize_voice_text(full)
                command = detect_voice_command(normalized)
                await ws.send_json({
                    "text":    normalized,
                    "raw":     full,
                    "final":   True,
                    "command": command,
                })
                break

            # Chunk audio PCM
            if data.get("type") == "websocket.receive" and data.get("bytes"):
                chunk = data["bytes"]
                if rec.AcceptWaveform(chunk):
                    result = _json.loads(rec.Result())
                    word = result.get("text", "")
                    if word:
                        full_results.append(word)
                        partial_text = " ".join(full_results)
                        await ws.send_json({"partial": partial_text, "final": False})
                else:
                    partial = _json.loads(rec.PartialResult()).get("partial", "")
                    if partial:
                        await ws.send_json({
                            "partial": (" ".join(full_results) + " " + partial).strip(),
                            "final": False,
                        })

    except ImportError:
        await ws.send_json({"error": "Vosk non disponible — streaming STT désactivé", "final": True})
    except Exception as e:
        logger.error(f"[WS/STT] Erreur: {e}")
        try:
            await ws.send_json({"error": str(e), "final": True})
        except Exception:
            pass
    finally:
        logger.info("[WS/STT] Connexion fermée")
        try:
            await ws.close()
        except Exception:
            pass


class DictationEnhancedRequest(BaseModel):
    model_config = {"arbitrary_types_allowed": True}
    text:        str
    source_id:   str = ""
    schema_data: Optional[Any] = None  # tables/champs connus (Dict[str, List[str]])

DictationEnhancedRequest.model_rebuild()


@app.post("/llm/extract-entities-v2", tags=["LLM"])
async def extract_entities_enhanced(req: DictationEnhancedRequest):
    """
    Extraction d'entités structurées depuis une dictée vocale — version améliorée.
    §2.3.2D — Dictée structurée avec contexte schéma.

    Retourne:
    - entities: {action, entity_type, fields: {field: value}, confidence}
    - suggested_sql: SQL généré si pertinent
    - confirmation_message: phrase à lire à l'utilisateur pour confirmation
    """
    try:
        import httpx, json as _json, os, re as _re
        OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://host.docker.internal:11434")
        OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "qwen2.5-coder:3b")

        schema_hint = ""
        if req.schema_data:
            schema_hint = "\nKnown tables and fields:\n"
            for tbl, fields in list(req.schema_data.items())[:4]:
                schema_hint += f"  [{tbl}]: {', '.join(fields[:8])}\n"

        prompt = f"""Extract structured entities from this French/English voice dictation.
Text: "{req.text}"{schema_hint}

Return ONLY valid JSON (no markdown, no explanation):
{{
  "action": "create|update|delete|search|filter",
  "entity_type": "order|product|customer|invoice|payment|other",
  "fields": {{"field_name": "value"}},
  "confidence": 0.0-1.0,
  "confirmation_message": "short French sentence to read back to user for confirmation"
}}

Examples:
- "Crée une commande pour Dupont, 50 unités ABC-123" → action=create, entity=order, fields={{customer:"Dupont",qty:50,product:"ABC-123"}}
- "Filtre les clients dont le montant dépasse 1000" → action=filter, entity=customer, fields={{amount_filter:">1000"}}
"""
        resp = httpx.post(
            f"{OLLAMA_HOST}/api/generate",
            json={"model": OLLAMA_MODEL, "prompt": prompt, "stream": False,
                  "options": {"temperature": 0.05, "num_predict": 300}},
            timeout=30,
        )
        raw = resp.json().get("response", "").strip()
        raw = _re.sub(r"```(?:json)?\s*", "", raw).replace("```", "").strip()
        match = _re.search(r"\{.*\}", raw, _re.DOTALL)
        entities = _json.loads(match.group(0)) if match else {}

        # Génère un SQL simple si action=search/filter
        suggested_sql = None
        if entities.get("action") in ("search", "filter") and req.schema_data:
            tbl = list(req.schema_data.keys())[0]
            filters = entities.get("fields", {})
            where_parts = []
            for k, v in filters.items():
                if isinstance(v, str) and v.startswith(">"):
                    where_parts.append(f"[{k}] > {v[1:].strip()}")
                elif isinstance(v, str) and v.startswith("<"):
                    where_parts.append(f"[{k}] < {v[1:].strip()}")
                else:
                    where_parts.append(f"[{k}] = '{v}'")
            if where_parts:
                suggested_sql = f"SELECT TOP 100 * FROM [{tbl}] WHERE {' AND '.join(where_parts)}"

        return {
            "entities":             entities,
            "suggested_sql":        suggested_sql,
            "confirmation_message": entities.get("confirmation_message", ""),
            "text":                 req.text,
        }

    except Exception as e:
        logger.error(f"[Dictation] Erreur: {e}")
        return {"entities": None, "error": str(e), "text": req.text}


# ══════════════════════════════════════════════════════════════
# SEMANTIC ENRICHMENT — §2.2.3
# ══════════════════════════════════════════════════════════════

try:
    from .semantic_enricher import enrich_source, semantic_search as _semantic_search
    _SEMANTIC_AVAILABLE = True
except ImportError:
    _SEMANTIC_AVAILABLE = False
    logger.warning("semantic_enricher non disponible")


@app.post("/sources/{source_id}/enrich", tags=["Semantic"])
async def enrich_source_endpoint(
    source_id: UUID,
    background_tasks: BackgroundTasks,
):
    """
    Lance l'enrichissement sémantique d'une source :
    - Classification domaine métier (Finance, RH, Ventes...)
    - Extraction concept métier (Customer, Order, Invoice...)
    - Détection dimensions analytiques (temps, géo, produit)
    - Indexation MeiliSearch (full-text search)
    - Calcul embeddings TF-IDF + pgvector (recherche sémantique)
    """
    if not _SEMANTIC_AVAILABLE:
        raise HTTPException(503, "Module semantic_enricher non disponible")

    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    async def _run_enrich():
        try:
            result = await enrich_source(
                source_id=source_id,
                source_name=source.name,
                source_type=source.connector_type.value,
            )
            logger.info(f"[Enrich] ✅ {source.name}: {result.get('enriched', 0)} entités enrichies")
        except Exception as e:
            logger.error(f"[Enrich] ❌ {source.name}: {e}", exc_info=True)

    background_tasks.add_task(_run_enrich)
    return {
        "source_id": str(source_id),
        "source_name": source.name,
        "message": "Enrichissement sémantique lancé en arrière-plan",
        "status": "running",
    }


@app.get("/sources/{source_id}/enrich/status", tags=["Semantic"])
async def get_enrich_status(source_id: UUID):
    """Retourne le statut de l'enrichissement sémantique d'une source."""
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        stats = await conn.fetchrow("""
            SELECT
                COUNT(*)                                            AS total,
                COUNT(*) FILTER (WHERE business_domain IS NOT NULL) AS enriched,
                COUNT(DISTINCT business_domain)                     AS domains_count,
                COUNT(DISTINCT business_concept)                    AS concepts_count,
                COUNT(*) FILTER (WHERE embedding IS NOT NULL)       AS with_embedding
            FROM source_entities
            WHERE source_id = $1
              AND entity_type IN ('table','view','materialized_view',
                                  'odata_entity','rest_resource','excel_sheet')
        """, source_id)

        domain_stats = await conn.fetch("""
            SELECT business_domain, COUNT(*) AS cnt
            FROM source_entities
            WHERE source_id = $1
              AND business_domain IS NOT NULL
            GROUP BY business_domain
            ORDER BY cnt DESC
        """, source_id)

        concept_stats = await conn.fetch("""
            SELECT business_concept, COUNT(*) AS cnt
            FROM source_entities
            WHERE source_id = $1
              AND business_concept IS NOT NULL
            GROUP BY business_concept
            ORDER BY cnt DESC
        """, source_id)

    enriched = stats["enriched"] or 0
    total    = stats["total"] or 1
    return {
        "source_id":       str(source_id),
        "source_name":     source.name,
        "total_entities":  stats["total"],
        "enriched":        enriched,
        "progress_pct":    round(enriched / total * 100, 1),
        "with_embedding":  stats["with_embedding"],
        "domains_count":   stats["domains_count"],
        "concepts_count":  stats["concepts_count"],
        "domain_stats":    {r["business_domain"]: r["cnt"] for r in domain_stats},
        "concept_stats":   {r["business_concept"]: r["cnt"] for r in concept_stats},
    }


@app.get("/search", tags=["Semantic"])
async def search_entities(
    q:          str   = Query(..., description="Requête en langage naturel"),
    source_ids: str   = Query(None, description="IDs sources séparés par virgule"),
    limit:      int   = Query(10, ge=1, le=50),
    use_vector: bool  = Query(True, description="Activer la recherche vectorielle"),
):
    """
    Recherche sémantique hybride sur toutes les entités indexées.

    Exemples :
    - /search?q=chiffre affaires client
    - /search?q=customer invoice total amount
    - /search?q=commande fournisseur&source_ids=uuid1,uuid2
    """
    if not _SEMANTIC_AVAILABLE:
        raise HTTPException(503, "Module semantic_enricher non disponible")

    sids = [s.strip() for s in source_ids.split(",")] if source_ids else None

    try:
        results = await _semantic_search(
            query=q,
            source_ids=sids,
            limit=limit,
            use_vector=use_vector,
        )
        return {
            "query":   q,
            "total":   len(results),
            "results": results,
        }
    except Exception as e:
        logger.error(f"[Search] {q}: {e}", exc_info=True)
        raise HTTPException(500, f"Erreur recherche: {str(e)}")


@app.get("/sources/{source_id}/entities/by-domain/{domain}", tags=["Semantic"])
async def get_entities_by_domain(source_id: UUID, domain: str):
    """Retourne les entités d'un domaine métier spécifique."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT name, entity_type, business_concept, entity_class,
                   semantic_tags, row_count, description
            FROM source_entities
            WHERE source_id = $1
              AND LOWER(business_domain) = LOWER($2)
            ORDER BY entity_class, name
        """, source_id, domain)

    return {
        "source_id": str(source_id),
        "domain":    domain,
        "total":     len(rows),
        "entities":  [
            {
                "name":        r["name"],
                "type":        r["entity_type"],
                "concept":     r["business_concept"],
                "class":       r["entity_class"],
                "tags":        json_module.loads(r["semantic_tags"] or "[]"),
                "row_count":   r["row_count"],
                "description": r["description"],
            }
            for r in rows
        ],
    }


@app.get("/sources/{source_id}/dimensions", tags=["Semantic"])
async def get_source_dimensions(source_id: UUID):
    """
    Retourne toutes les dimensions analytiques détectées pour une source,
    avec les hiérarchies SQL (Year→Quarter→Month→Week→Day, etc.)
    Utile pour construire des dashboards et des analyses OLAP.
    """
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT se.name AS entity_name, ef.name AS field_name,
                   ef.dimension_type, ef.data_type,
                   ef.dimension_hierarchy
            FROM entity_fields ef
            JOIN source_entities se ON se.id = ef.entity_id
            WHERE se.source_id = $1
              AND ef.dimension_type IS NOT NULL
            ORDER BY ef.dimension_type, se.name, ef.name
        """, source_id)

    # Grouper par type de dimension
    dimensions: dict = {}
    for r in rows:
        dim = r["dimension_type"]
        if dim not in dimensions:
            dimensions[dim] = []

        # Parser la hiérarchie JSON
        hierarchy = None
        if r["dimension_hierarchy"]:
            try:
                import json as json_module
                raw = r["dimension_hierarchy"]
                hierarchy = json_module.loads(raw) if isinstance(raw, str) else dict(raw)
            except Exception:
                hierarchy = None

        dimensions[dim].append({
            "entity":    r["entity_name"],
            "field":     r["field_name"],
            "type":      r["data_type"],
            "hierarchy": hierarchy,
        })

    return {
        "source_id":  str(source_id),
        "dimensions": dimensions,
        "summary":    {dim: len(cols) for dim, cols in dimensions.items()},
    }


# ══════════════════════════════════════════════════════════════
# SYNONYMES CLIENT — §2.2.3.A
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/synonyms", tags=["Semantic"])
async def add_source_synonym(source_id: UUID, body: dict):
    """
    Ajoute ou met à jour un synonyme métier spécifique à une source.
    Ex: {"term": "BALI", "synonyms": ["balance", "solde"], "description": "Balance SXA"}
    """
    term        = body.get("term", "").strip().upper()
    synonyms    = body.get("synonyms", [])
    description = body.get("description", "")
    created_by  = body.get("created_by", "user")

    if not term or not synonyms:
        raise HTTPException(400, "term et synonyms sont requis")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        row = await conn.fetchrow("""
            INSERT INTO source_synonyms (source_id, term, synonyms, description, created_by)
            VALUES ($1, $2, $3, $4, $5)
            ON CONFLICT (source_id, term)
            DO UPDATE SET synonyms=$3, description=$4, updated_at=NOW()
            RETURNING *
        """, source_id, term, synonyms, description, created_by)

    # Synchroniser dans MeiliSearch
    try:
        from .semantic_enricher import _get_meili_client, MEILI_INDEX
        meili = _get_meili_client()
        if meili:
            # Charger tous les synonymes de la source pour MeiliSearch
            async with pool.acquire() as conn:
                all_rows = await conn.fetch(
                    "SELECT term, synonyms FROM source_synonyms WHERE source_id=$1",
                    source_id
                )
            meili_synonyms = {}
            for r in all_rows:
                t = r["term"].lower()
                syns = [s.lower() for s in r["synonyms"]]
                meili_synonyms[t] = syns
                for s in syns:
                    meili_synonyms[s] = meili_synonyms.get(s, []) + [t]

            meili.index(MEILI_INDEX).update_synonyms(meili_synonyms)
            logger.info(f"[Synonyms] MeiliSearch mis à jour: {len(meili_synonyms)} synonymes")
    except Exception as e:
        logger.warning(f"[Synonyms] MeiliSearch sync failed: {e}")

    return {
        "success":     True,
        "term":        term,
        "synonyms":    synonyms,
        "description": description,
        "message":     f"Synonyme '{term}' ajouté et synchronisé dans MeiliSearch",
    }


@app.get("/sources/{source_id}/synonyms", tags=["Semantic"])
async def list_source_synonyms(source_id: UUID):
    """Liste tous les synonymes métier d'une source."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT term, synonyms, description, created_by, created_at, updated_at
            FROM source_synonyms
            WHERE source_id = $1
            ORDER BY term
        """, source_id)

    return {
        "source_id": str(source_id),
        "count":     len(rows),
        "synonyms":  [
            {
                "term":        r["term"],
                "synonyms":    list(r["synonyms"]),
                "description": r["description"],
                "created_by":  r["created_by"],
                "updated_at":  r["updated_at"].isoformat() if r["updated_at"] else None,
            }
            for r in rows
        ],
    }


@app.delete("/sources/{source_id}/synonyms/{term}", tags=["Semantic"])
async def delete_source_synonym(source_id: UUID, term: str):
    """Supprime un synonyme métier et met à jour MeiliSearch."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        result = await conn.execute(
            "DELETE FROM source_synonyms WHERE source_id=$1 AND term=$2",
            source_id, term.upper()
        )

    if result == "DELETE 0":
        raise HTTPException(404, f"Synonyme '{term}' introuvable")

    # Re-synchroniser MeiliSearch sans ce synonyme
    try:
        from .semantic_enricher import _get_meili_client, MEILI_INDEX
        meili = _get_meili_client()
        if meili:
            async with pool.acquire() as conn:
                all_rows = await conn.fetch(
                    "SELECT term, synonyms FROM source_synonyms WHERE source_id=$1",
                    source_id
                )
            meili_synonyms = {}
            for r in all_rows:
                t = r["term"].lower()
                syns = [s.lower() for s in r["synonyms"]]
                meili_synonyms[t] = syns
            meili.index(MEILI_INDEX).update_synonyms(meili_synonyms)
    except Exception as e:
        logger.warning(f"[Synonyms] MeiliSearch sync failed: {e}")

    return {"success": True, "message": f"Synonyme '{term}' supprimé"}


# ══════════════════════════════════════════════════════════════
# CALENDRIERS FISCAUX — §2.2.3.B
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/fiscal-calendar", tags=["Semantic"])
async def set_fiscal_calendar(source_id: UUID, body: dict):
    """
    Configure l'année fiscale d'une source.
    Ex: {"fiscal_year_start": 7, "description": "Exercice juillet-juin"}
    fiscal_year_start : 1=Janvier (défaut), 4=Avril, 7=Juillet, 10=Octobre
    """
    fiscal_year_start = int(body.get("fiscal_year_start", 1))
    description       = body.get("description", "")
    fiscal_year_label = body.get("fiscal_year_label", "FY")

    if not 1 <= fiscal_year_start <= 12:
        raise HTTPException(400, "fiscal_year_start doit être entre 1 et 12")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        await conn.execute("""
            INSERT INTO source_fiscal_calendars
                (source_id, fiscal_year_start, fiscal_year_label, description)
            VALUES ($1, $2, $3, $4)
            ON CONFLICT (source_id)
            DO UPDATE SET fiscal_year_start=$2, fiscal_year_label=$3,
                          description=$4, updated_at=NOW()
        """, source_id, fiscal_year_start, fiscal_year_label, description)

    month_names = {1:"Janvier",2:"Février",3:"Mars",4:"Avril",5:"Mai",6:"Juin",
                   7:"Juillet",8:"Août",9:"Septembre",10:"Octobre",11:"Novembre",12:"Décembre"}

    return {
        "success":           True,
        "fiscal_year_start": fiscal_year_start,
        "fiscal_year_label": fiscal_year_label,
        "description":       description,
        "message":           f"Année fiscale configurée: début en {month_names.get(fiscal_year_start, str(fiscal_year_start))}. Lancez un sync pour recalculer les hiérarchies.",
    }


@app.get("/sources/{source_id}/fiscal-calendar", tags=["Semantic"])
async def get_fiscal_calendar(source_id: UUID):
    """Retourne la configuration du calendrier fiscal d'une source."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        row = await conn.fetchrow(
            "SELECT * FROM source_fiscal_calendars WHERE source_id=$1",
            source_id
        )

    if not row:
        return {
            "source_id":         str(source_id),
            "fiscal_year_start": 1,
            "fiscal_year_label": "FY",
            "description":       "Calendrier standard (janvier)",
            "configured":        False,
        }

    return {
        "source_id":         str(source_id),
        "fiscal_year_start": row["fiscal_year_start"],
        "fiscal_year_label": row["fiscal_year_label"],
        "description":       row["description"],
        "configured":        True,
        "updated_at":        row["updated_at"].isoformat() if row["updated_at"] else None,
    }


# ══════════════════════════════════════════════════════════════
# TAXONOMIE PERSONNALISABLE — §2.2.3.A
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/domains", tags=["Semantic"])
async def add_custom_domain(source_id: UUID, body: dict):
    """
    Ajoute un domaine métier personnalisé pour une source.
    Ex: {"domain_name": "Trésorerie", "patterns": ["tresor","cash","bnk"], "color": "#f59e0b"}
    """
    domain_name = body.get("domain_name", "").strip()
    patterns    = body.get("patterns", [])
    color       = body.get("color", "#6366f1")
    icon        = body.get("icon", "database")
    description = body.get("description", "")
    priority    = int(body.get("priority", 10))

    if not domain_name or not patterns:
        raise HTTPException(400, "domain_name et patterns sont requis")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        await conn.execute("""
            INSERT INTO source_domains
                (source_id, domain_name, patterns, color, icon, description, priority)
            VALUES ($1, $2, $3, $4, $5, $6, $7)
            ON CONFLICT (source_id, domain_name)
            DO UPDATE SET patterns=$3, color=$4, icon=$5,
                          description=$6, priority=$7, updated_at=NOW()
        """, source_id, domain_name, patterns, color, icon, description, priority)

    return {
        "success":     True,
        "domain_name": domain_name,
        "patterns":    patterns,
        "message":     f"Domaine '{domain_name}' ajouté. Lancez un sync pour reclassifier les entités.",
    }


@app.get("/sources/{source_id}/domains", tags=["Semantic"])
async def list_custom_domains(source_id: UUID):
    """Liste tous les domaines personnalisés d'une source."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT domain_name, patterns, color, icon, description, priority, updated_at
            FROM source_domains
            WHERE source_id = $1
            ORDER BY priority DESC, domain_name
        """, source_id)

    return {
        "source_id": str(source_id),
        "count":     len(rows),
        "domains":   [
            {
                "domain_name": r["domain_name"],
                "patterns":    list(r["patterns"]),
                "color":       r["color"],
                "icon":        r["icon"],
                "description": r["description"],
                "priority":    r["priority"],
                "updated_at":  r["updated_at"].isoformat() if r["updated_at"] else None,
            }
            for r in rows
        ],
    }


@app.delete("/sources/{source_id}/domains/{domain_name}", tags=["Semantic"])
async def delete_custom_domain(source_id: UUID, domain_name: str):
    """Supprime un domaine personnalisé."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        result = await conn.execute(
            "DELETE FROM source_domains WHERE source_id=$1 AND domain_name=$2",
            source_id, domain_name
        )
    if result == "DELETE 0":
        raise HTTPException(404, f"Domaine '{domain_name}' introuvable")
    return {"success": True, "message": f"Domaine '{domain_name}' supprimé. Relancez un sync."}


# ══════════════════════════════════════════════════════════════
# HISTORIQUE RECHERCHES — §2.2.3.C
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/search/click", tags=["Semantic"])
async def record_search_click(source_id: UUID, body: dict):
    """
    Enregistre le clic d'un utilisateur sur un résultat de recherche.
    Ex: {"query": "facture client", "entity_id": "uuid", "entity_name": "GS_GLACC"}
    Utilisé pour boosting des résultats futurs.
    """
    query       = body.get("query", "").strip()
    entity_id   = body.get("entity_id")
    entity_name = body.get("entity_name", "")
    user_id     = body.get("user_id")

    if not query or not entity_id:
        raise HTTPException(400, "query et entity_id sont requis")

    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        await conn.execute("""
            UPDATE search_history
            SET clicked_id=$1, clicked_name=$2
            WHERE source_id=$3
              AND query_norm=$4
              AND clicked_id IS NULL
              AND search_at = (
                SELECT MAX(search_at) FROM search_history
                WHERE source_id=$3 AND query_norm=$4
              )
        """, UUID(entity_id), entity_name, source_id, query.lower().strip())

    return {"success": True, "message": f"Clic enregistré pour '{entity_name}'"}


@app.get("/sources/{source_id}/search/history", tags=["Semantic"])
async def get_search_history(source_id: UUID, limit: int = 20):
    """Retourne les statistiques des recherches passées pour une source."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT query_norm, search_count, click_count,
                   click_rate_pct, last_searched_at, most_clicked
            FROM search_history_stats
            WHERE source_id = $1
            LIMIT $2
        """, source_id, limit)

    return {
        "source_id": str(source_id),
        "count":     len(rows),
        "history":   [
            {
                "query":           r["query_norm"],
                "search_count":    r["search_count"],
                "click_count":     r["click_count"],
                "click_rate_pct":  float(r["click_rate_pct"] or 0),
                "last_searched":   r["last_searched_at"].isoformat() if r["last_searched_at"] else None,
                "most_clicked":    r["most_clicked"],
            }
            for r in rows
        ],
    }


# ══════════════════════════════════════════════════════════════
# CDC / VERSIONING — §2.2.5
# ══════════════════════════════════════════════════════════════

def _get_cdc() -> CDCEngine:
    """Fabrique un CDCEngine (pool + redis déjà initialisés au lifespan)."""
    from .database import _pg_pool, _redis  # singletons module-level
    return CDCEngine(_pg_pool, _redis)


def _get_versioner() -> SchemaVersioner:
    from .database import _pg_pool, _redis
    return SchemaVersioner(_pg_pool, _redis)


# ── 1. Détecter les changements (CDC) ────────────────────────

@app.post("/sources/{source_id}/cdc/detect", tags=["CDC"])
async def cdc_detect(source_id: UUID):
    """
    Compare le schéma courant avec la dernière version enregistrée.
    - Crée une nouvelle version si des changements sont détectés.
    - Publie les breaking changes dans Redis.
    - Invalide le cache de la source.
    """
    cdc = _get_cdc()
    try:
        result = await cdc.detect_and_record(source_id)
        return result
    except Exception as e:
        logger.error(f"[CDC] detect error source {source_id}: {e}", exc_info=True)
        raise HTTPException(500, str(e))


# ── 1b. Reset baseline CDC ────────────────────────────────────

@app.post("/sources/{source_id}/cdc/reset-baseline", tags=["CDC"])
async def cdc_reset_baseline(source_id: UUID):
    """
    Réinitialise la baseline CDC pour une source.
    - Supprime tout l'historique des versions.
    - Crée un nouveau v1 propre à partir du schéma actuel.
    - Utile quand la source a été restaurée depuis un .bak ou réimportée.
    """
    pool  = await get_pg_pool()
    redis = await get_redis()
    cdc   = CDCEngine(pool, redis)

    try:
        current_schema = await cdc.snapshot_schema(source_id)
        if not current_schema:
            raise HTTPException(400, "Aucune entité trouvée — faites une synchro d'abord.")

        from .cdc_engine import compute_schema_fingerprint
        import json
        fp      = compute_schema_fingerprint(current_schema)
        summary = {"added": 0, "dropped": 0, "modified": 0, "total": 0}

        # Supprime l'historique (tags en premier pour respecter les FK)
        await pool.execute("DELETE FROM schema_version_tags WHERE source_id = $1", source_id)
        await pool.execute("DELETE FROM schema_versions      WHERE source_id = $1", source_id)

        # Nettoie Redis
        try:
            keys = await redis.keys(f"cdc:*{source_id}*")
            if keys:
                await redis.delete(*keys)
        except Exception:
            pass

        # Insère le nouveau v1 propre
        await pool.execute("""
            INSERT INTO schema_versions (
                source_id, version_number, fingerprint,
                schema_snapshot, changes_delta, has_breaking_changes,
                change_summary, created_at
            ) VALUES ($1, 1, $2, $3, $4, FALSE, $5, NOW())
        """,
            source_id,
            fp,
            json.dumps(current_schema, default=str),
            json.dumps([]),
            json.dumps(summary),
        )

        logger.info(f"[CDC] reset-baseline source {source_id} → v1 propre ({len(current_schema)} entités)")
        return {
            "status":      "reset",
            "new_version": 1,
            "fingerprint": fp,
            "entities":    len(current_schema),
            "message":     "Baseline réinitialisée. v1 = état actuel de la source.",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[CDC] reset-baseline error: {e}", exc_info=True)
        raise HTTPException(500, str(e))

@app.get("/sources/{source_id}/versions", tags=["CDC"])
async def schema_log(
    source_id: UUID,
    limit: int = Query(50, ge=1, le=200),
):
    """
    Retourne l'historique des versions de schéma pour une source.
    Enrichi avec les tags et le résumé des changements.
    """
    versioner = _get_versioner()
    return {
        "source_id": str(source_id),
        "versions":  await versioner.get_log(source_id, limit),
    }


# ── 3. Détail d'une version (git show) ───────────────────────

@app.get("/sources/{source_id}/versions/{version}", tags=["CDC"])
async def schema_version_detail(source_id: UUID, version: int):
    """
    Retourne le détail complet d'une version :
    snapshot complet du schéma + liste des deltas + tags.
    """
    versioner = _get_versioner()
    v = await versioner.get_version(source_id, version)
    if not v:
        raise HTTPException(404, f"Version {version} introuvable pour source {source_id}")
    return v


# ── 4. Diff entre deux versions (git diff) ───────────────────

@app.get("/sources/{source_id}/versions/diff", tags=["CDC"])
async def schema_diff(
    source_id: UUID,
    v1: int = Query(..., description="Version de base"),
    v2: int = Query(..., description="Version cible"),
):
    """
    Diff entre deux versions arbitraires.
    Retourne la liste complète des changements avec flag breaking.
    """
    versioner = _get_versioner()
    try:
        return await versioner.diff_versions(source_id, v1, v2)
    except ValueError as e:
        raise HTTPException(404, str(e))


# ── 5. Rollback vers une version passée ──────────────────────

@app.post("/sources/{source_id}/versions/{version}/rollback", tags=["CDC"])
async def schema_rollback(source_id: UUID, version: int):
    """
    Rollback logique : crée une nouvelle version HEAD
    qui restaure le snapshot d'une version passée.
    Ne modifie PAS la source de données réelle —
    sert à tracer le retour arrière dans l'historique OnePilot.
    """
    cdc = _get_cdc()
    try:
        return await cdc.rollback_to_version(source_id, version)
    except ValueError as e:
        raise HTTPException(404, str(e))
    except Exception as e:
        logger.error(f"[CDC] rollback error: {e}", exc_info=True)
        raise HTTPException(500, str(e))


# ── 6. Tagger une version (git tag) ──────────────────────────

@app.post("/sources/{source_id}/versions/{version}/tag", tags=["CDC"])
async def schema_tag(source_id: UUID, version: int, body: dict):
    """
    Pose un tag nommé sur une version.
    Body: { "tag": "v1.2-stable", "note": "Avant migration SAGE" }
    """
    tag  = body.get("tag", "").strip()
    note = body.get("note", "")
    if not tag:
        raise HTTPException(400, "Le champ 'tag' est requis")

    versioner = _get_versioner()
    try:
        return await versioner.tag_version(source_id, version, tag, note)
    except ValueError as e:
        raise HTTPException(404, str(e))


@app.delete("/sources/{source_id}/tags/{tag}", tags=["CDC"])
async def schema_delete_tag(source_id: UUID, tag: str):
    """Supprime un tag."""
    versioner = _get_versioner()
    try:
        return await versioner.delete_tag(source_id, tag)
    except ValueError as e:
        raise HTTPException(404, str(e))


@app.get("/sources/{source_id}/tags", tags=["CDC"])
async def schema_list_tags(source_id: UUID):
    """Liste tous les tags d'une source."""
    versioner = _get_versioner()
    tags = await versioner.list_tags(source_id)
    return {"source_id": str(source_id), "tags": tags}


# ── 7. Notifications breaking changes ────────────────────────

@app.get("/sources/{source_id}/cdc/notifications", tags=["CDC"])
async def cdc_notifications(
    source_id: UUID,
    limit: int = Query(50, ge=1, le=100),
):
    """
    Retourne les dernières notifications de breaking changes.
    Lecture depuis Redis (TTL 7j), fallback PostgreSQL.
    """
    versioner = _get_versioner()
    notifs = await versioner.get_notifications(source_id, limit)
    return {
        "source_id": str(source_id),
        "count":     len(notifs),
        "notifications": notifs,
    }


# ── 7b. Subscriber status ─────────────────────────────────────

@app.get("/sources/{source_id}/cdc/subscriber-log", tags=["CDC"])
async def cdc_subscriber_log(
    source_id: UUID,
    limit: int = Query(20, ge=1, le=100),
):
    """
    Retourne l'historique des traitements CDC automatiques :
    - Invalidations de cache déclenchées
    - Réindexations MeiliSearch effectuées
    - Version et nombre de changements traités
    Utile pour vérifier que le subscriber fonctionne correctement.
    """
    pool = await get_pg_pool()
    try:
        rows = await pool.fetch("""
            SELECT
                source_id, version, change_count,
                cache_invalidated, reindex_triggered, processed_at
            FROM cdc_subscriber_log
            WHERE source_id = $1
            ORDER BY processed_at DESC
            LIMIT $2
        """, source_id, limit)
        return {
            "source_id": str(source_id),
            "count": len(rows),
            "events": [
                {
                    "version":           r["version"],
                    "change_count":      r["change_count"],
                    "cache_invalidated": r["cache_invalidated"],
                    "reindex_triggered": r["reindex_triggered"],
                    "processed_at":      r["processed_at"].isoformat(),
                }
                for r in rows
            ],
        }
    except Exception as e:
        return {"source_id": str(source_id), "count": 0, "events": [], "note": str(e)}


@app.post("/sources/{source_id}/cdc/trigger-invalidation", tags=["CDC"])
async def cdc_trigger_invalidation(source_id: UUID):
    """
    Déclenche manuellement l'invalidation du cache et la réindexation
    pour une source donnée. Utile pour tester le pipeline CDC sans attendre
    un vrai breaking change.
    """
    redis = await get_redis()
    pool  = await get_pg_pool()

    # Récupère la dernière version
    last_version = await pool.fetchval("""
        SELECT COALESCE(MAX(version_number), 0)
        FROM schema_versions WHERE source_id = $1
    """, source_id)

    # Invalide le cache Redis
    cache_count = 0
    try:
        dash_keys = await redis.keys(f"onepilot:dashboard:{source_id}:*")
        step_keys = await redis.keys("onepilot:step:*")
        prof_keys = await redis.keys(f"onepilot:profile:{source_id}:*")
        all_keys  = dash_keys + step_keys + prof_keys
        if all_keys:
            await redis.delete(*all_keys)
            cache_count = len(all_keys)
    except Exception as e:
        logger.warning(f"[CDC Trigger] Cache error: {e}")

    # Déclenche réindexation
    reindex_result = {}
    try:
        from .cdc_reindexer import CDCReindexer
        reindexer = CDCReindexer(pool, redis)
        reindex_result = await reindexer.full_reindex(source_id)
    except Exception as e:
        reindex_result = {"error": str(e)}

    return {
        "source_id":        str(source_id),
        "version":          last_version,
        "cache_invalidated": cache_count,
        "reindex":           reindex_result,
        "status":            "ok",
    }



# ══════════════════════════════════════════════════════════════
# FEEDBACK LOOP — Sprint 12
# ══════════════════════════════════════════════════════════════

async def _auto_retrain_fasttext_bg():
    """
    Réentraîne FastText automatiquement en arrière-plan.
    Acquiert son propre pool — évite le bug pool fermé avant exécution.
    """
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            feedback_rows = await conn.fetch("""
                SELECT question, intent
                FROM   chat_feedback
                WHERE  feedback_type     = 'like'
                  AND  used_for_training = FALSE
                  AND  question          != ''
                ORDER  BY created_at ASC
                LIMIT  500
            """)

        if not feedback_rows:
            logger.info("[AutoRetrain] Aucun exemple disponible — abandon")
            return

        examples = [(r["question"], r["intent"]) for r in feedback_rows]
        result   = retrain_fasttext_with_feedback(examples)

        if result["status"] != "error":
            questions = [r["question"] for r in feedback_rows]
            async with pool.acquire() as conn:
                await conn.execute("""
                    UPDATE chat_feedback
                    SET    used_for_training = TRUE
                    WHERE  feedback_type     = 'like'
                      AND  used_for_training = FALSE
                      AND  question          = ANY()
                """, questions)
            logger.info(
                f"[AutoRetrain] ✅ FastText réentraîné automatiquement — "
                f"+{result['examples_added']} exemples, total={result['total_examples']}"
            )
        else:
            logger.error(f"[AutoRetrain] ❌ Erreur retrain: {result['message']}")

    except Exception as e:
        logger.error(f"[AutoRetrain] Exception non gérée: {e}")


@app.post("/feedback", tags=["Feedback"])
async def submit_chat_feedback(req: dict):
    """
    Enregistre le feedback utilisateur (👍/👎) depuis le chat.
    Appelé automatiquement par chat.html quand l'utilisateur clique
    sur like/dislike.

    Body: {
        conversation_id, source_id, message_id,
        feedback_type: "like"|"dislike",
        answer: str (texte de la réponse)
    }
    """
    pool = await get_pg_pool()

    conversation_id = req.get("conversation_id", "")
    source_id_str   = req.get("source_id", "")
    message_id      = req.get("message_id", "")
    feedback_type   = req.get("feedback_type", "")
    answer          = req.get("answer", "")[:500]
    # Question envoyée directement depuis chat.html
    question_direct = req.get("question", "")

    if feedback_type not in ("like", "dislike"):
        raise HTTPException(400, "feedback_type doit être 'like' ou 'dislike'")

    # Récupère question + intent depuis nlu_query_log
    # Si question déjà fournie par le client, on l'utilise directement
    question   = question_direct
    intent     = ""
    confidence = 0.0
    try:
        async with pool.acquire() as conn:
            row = await conn.fetchrow("""
                SELECT question, intent, confidence
                FROM   nlu_query_log
                WHERE  conversation_id = $1
                ORDER  BY created_at DESC
                LIMIT  1
            """, conversation_id)
            if row:
                if not question:
                    question = row["question"] or ""
                intent     = str(row["intent"] or "")
                confidence = float(row["confidence"] or 0)
    except Exception as e:
        logger.warning(f"[Feedback] nlu_query_log lookup error: {e}")

    # Enregistre le feedback
    try:
        async with pool.acquire() as conn:
            await conn.execute("""
                INSERT INTO chat_feedback
                    (conversation_id, source_id, message_id,
                     question, answer, intent, confidence,
                     feedback_type, created_at)
                VALUES ($1,$2,$3,$4,$5,$6,$7,$8,NOW())
            """,
                conversation_id,
                UUID(source_id_str) if source_id_str else None,
                message_id,
                question,
                answer,
                intent,
                confidence,
                feedback_type,
            )
    except Exception as e:
        logger.error(f"[Feedback] Insert error: {e}")
        raise HTTPException(500, f"Erreur enregistrement feedback: {e}")

    # Comptage des likes non encore utilisés pour training
    like_count = 0
    try:
        async with pool.acquire() as conn:
            like_count = await conn.fetchval("""
                SELECT COUNT(*) FROM chat_feedback
                WHERE  feedback_type     = 'like'
                  AND  used_for_training = FALSE
                  AND  question          != ''
            """)
    except Exception:
        pass

    logger.info(
        f"[Feedback] {feedback_type} enregistré — "
        f"question='{question[:40]}' intent={intent} "
        f"(likes non traités: {like_count})"
    )

    # ── Retrain automatique — seuil 10 likes non traités ─────────────
    auto_retrained = False
    if feedback_type == "like" and like_count >= 10:
        logger.info("[Feedback] Seuil 10 likes atteint → retrain FastText automatique")
        asyncio.create_task(_auto_retrain_fasttext_bg())
        auto_retrained = True

    return {
        "status":                 "ok",
        "feedback_type":          feedback_type,
        "likes_pending":          like_count,
        "retrain_suggested":      like_count >= 10,
        "auto_retrain_triggered": auto_retrained,
    }


@app.get("/feedback/stats", tags=["Feedback"])
async def feedback_stats():
    """
    Statistiques du feedback collecté.
    Retourne le nombre de likes/dislikes, les intents les plus corrigés,
    et si un réentraînement est recommandé.
    """
    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            total = await conn.fetchrow("""
                SELECT
                    COUNT(*) FILTER (WHERE feedback_type='like')    AS likes,
                    COUNT(*) FILTER (WHERE feedback_type='dislike') AS dislikes,
                    COUNT(*) FILTER (WHERE feedback_type='like'
                                      AND used_for_training=FALSE
                                      AND question != ''
                                      )           AS likes_pending
                FROM chat_feedback
            """)

            by_intent = await conn.fetch("""
                SELECT intent,
                       COUNT(*) FILTER (WHERE feedback_type='like')    AS likes,
                       COUNT(*) FILTER (WHERE feedback_type='dislike') AS dislikes
                FROM   chat_feedback
                WHERE  intent != ''
                GROUP  BY intent
                ORDER  BY dislikes DESC
                LIMIT  10
            """)

        return {
            "total_likes":        total["likes"],
            "total_dislikes":     total["dislikes"],
            "likes_pending":      total["likes_pending"],
            "retrain_suggested":  total["likes_pending"] >= 10,
            "by_intent": [
                {
                    "intent":   r["intent"],
                    "likes":    r["likes"],
                    "dislikes": r["dislikes"],
                }
                for r in by_intent
            ],
        }
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/feedback/retrain-fasttext", tags=["Feedback"])
async def retrain_fasttext_endpoint():
    """
    Réentraîne FastText avec les questions validées (👍).
    Ajoute les questions likées au dataset d'entraînement
    et recharge le modèle en mémoire immédiatement (hot reload).

    Déclencher manuellement ou automatiquement quand likes_pending >= 10.
    """
    pool = await get_pg_pool()

    # Récupère les questions likées non encore utilisées
    try:
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT question, intent
                FROM   chat_feedback
                WHERE  feedback_type     = 'like'
                  AND  used_for_training = FALSE
                  AND  question          != ''
                ORDER  BY created_at ASC
                LIMIT  500
            """)
    except Exception as e:
        raise HTTPException(500, f"Erreur lecture feedback: {e}")

    if not rows:
        return {
            "status":  "nothing_to_train",
            "message": "Aucun feedback 👍 disponible pour l'entraînement",
        }

    examples = [(r["question"], r["intent"]) for r in rows]

    # Réentraîne FastText
    result = retrain_fasttext_with_feedback(examples)

    if result["status"] == "error":
        raise HTTPException(500, result["message"])

    # Marque les exemples comme utilisés
    try:
        async with pool.acquire() as conn:
            questions = [r["question"] for r in rows]
            await conn.execute("""
                UPDATE chat_feedback
                SET    used_for_training = TRUE
                WHERE  feedback_type     = 'like'
                  AND  used_for_training = FALSE
                  AND  question          = ANY($1)
            """, questions)
    except Exception as e:
        logger.warning(f"[Feedback Retrain] Mark used error: {e}")

    logger.info(
        f"[Feedback Retrain] ✅ FastText réentraîné — "
        f"+{result['examples_added']} exemples, "
        f"total={result['total_examples']}"
    )

    return {
        "status":          "retrained",
        "examples_added":  result["examples_added"],
        "total_examples":  result["total_examples"],
        "model_reloaded":  result["model_reloaded"],
        "message":         f"FastText réentraîné avec {result['examples_added']} nouveaux exemples",
    }


@app.get("/feedback/history", tags=["Feedback"])
async def feedback_history(
    limit: int = Query(50, ge=1, le=200),
    feedback_type: Optional[str] = Query(None),
):
    """Historique des feedbacks collectés."""
    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            where = "WHERE 1=1"
            params = [limit]
            if feedback_type in ("like", "dislike"):
                where += " AND feedback_type = $2"
                params.append(feedback_type)

            rows = await conn.fetch(f"""
                SELECT conversation_id, source_id, message_id,
                       question, intent, confidence, feedback_type,
                       used_for_training, created_at
                FROM   chat_feedback
                {where}
                ORDER  BY created_at DESC
                LIMIT  $1
            """, *params)

        return {
            "count": len(rows),
            "feedbacks": [
                {
                    "question":          r["question"],
                    "intent":            r["intent"],
                    "confidence":        r["confidence"],
                    "feedback_type":     r["feedback_type"],
                    "used_for_training": r["used_for_training"],
                    "created_at":        r["created_at"].isoformat(),
                }
                for r in rows
            ],
        }
    except Exception as e:
        raise HTTPException(500, str(e))


# ── 8. Impact analysis ────────────────────────────────────────

@app.get("/sources/{source_id}/versions/{version}/impact", tags=["CDC"])
async def schema_impact(source_id: UUID, version: int):
    """
    Analyse l'impact des breaking changes d'une version :
    - Tables et colonnes affectées
    - Jointures / relations potentiellement cassées
    """
    versioner = _get_versioner()
    try:
        return await versioner.impact_analysis(source_id, version)
    except ValueError as e:
        raise HTTPException(404, str(e))


# ── 9. CDC intégré au pipeline de sync ───────────────────────
# Le sync background (_run_sync_background) appelle detect_and_record
# automatiquement après chaque sync complet.
# Ce endpoint permet un CDC on-demand sans sync complet.

@app.post("/sources/{source_id}/cdc/watch", tags=["CDC"])
async def cdc_watch(
    source_id: UUID,
    background_tasks: BackgroundTasks,
    interval_seconds: int = Query(300, ge=30, le=3600),
):
    """
    Lance une détection CDC différée (une seule fois, non-récurrente).
    Pour un polling continu, utilise un cron externe ou Celery Beat.
    """
    async def _delayed_detect():
        import asyncio
        await asyncio.sleep(interval_seconds)
        cdc = _get_cdc()
        try:
            result = await cdc.detect_and_record(source_id)
            logger.info(f"[CDC Watch] source {source_id}: {result.get('status')}")
        except Exception as e:
            logger.error(f"[CDC Watch] error: {e}")

    background_tasks.add_task(_delayed_detect)
    return {
        "status":           "scheduled",
        "source_id":        str(source_id),
        "detect_in_seconds": interval_seconds,
        "message":          f"Détection CDC programmée dans {interval_seconds}s",
    }


def _get_file_cdc():
    from .database import _pg_pool, _redis
    from .cdc_file_watcher import FileCDCEngine
    return FileCDCEngine(_pg_pool, _redis)

def _get_reindexer():
    from .database import _pg_pool, _redis
    from .cdc_reindexer import CDCReindexer
    return CDCReindexer(_pg_pool, _redis)


# ── 10. CDC Fichiers — détecter changement ────────────────────

@app.post("/sources/{source_id}/cdc/file/detect", tags=["CDC"])
async def cdc_file_detect(source_id: UUID):
    """
    Détecte si le fichier source a changé (checksum MD5 + timestamp).
    Pour sources : file_csv, file_excel, file_json.
    Si changement détecté → publie notification Redis + retourne needs_resync=True.
    """
    engine = _get_file_cdc()
    try:
        return await engine.detect_file_change(source_id)
    except Exception as e:
        logger.error(f"[CDC File] detect error source {source_id}: {e}")
        raise HTTPException(500, str(e))


# ── 11. CDC Fichiers — historique ────────────────────────────

@app.get("/sources/{source_id}/cdc/file/history", tags=["CDC"])
async def cdc_file_history(
    source_id: UUID,
    limit: int = Query(20, ge=1, le=100),
):
    """
    Retourne l'historique des changements de fichier détectés
    (checksums, tailles, deltas).
    """
    engine = _get_file_cdc()
    history = await engine.get_file_history(source_id, limit)
    return {
        "source_id": str(source_id),
        "count":     len(history),
        "history":   history,
    }


# ── 12. CDC Fichiers — notifications Redis ───────────────────

@app.get("/sources/{source_id}/cdc/file/notifications", tags=["CDC"])
async def cdc_file_notifications(
    source_id: UUID,
    limit: int = Query(20, ge=1, le=50),
):
    """Retourne les dernières notifications de changement fichier."""
    engine = _get_file_cdc()
    notifs = await engine.get_file_notifications(source_id, limit)
    return {
        "source_id":     str(source_id),
        "count":         len(notifs),
        "notifications": notifs,
    }


# ── 13. Réindexation MeiliSearch forcée ──────────────────────

@app.post("/sources/{source_id}/cdc/reindex", tags=["CDC"])
async def cdc_force_reindex(
    source_id: UUID,
    background_tasks: BackgroundTasks,
):
    """
    Force une réindexation complète de la source dans MeiliSearch.
    Utile après un rollback ou une resync complète.
    Lance en background — ne bloque pas la réponse.
    """
    async def _do_reindex():
        reindexer = _get_reindexer()
        try:
            result = await reindexer.full_reindex(source_id)
            logger.info(f"[CDC Reindex] {result}")
        except Exception as e:
            logger.error(f"[CDC Reindex] error: {e}")

    background_tasks.add_task(_do_reindex)
    return {
        "status":    "reindex_scheduled",
        "source_id": str(source_id),
        "message":   "Réindexation MeiliSearch lancée en arrière-plan",
    }


# ── 14. Relations nécessitant une révision ───────────────────

@app.get("/sources/{source_id}/cdc/relations/review", tags=["CDC"])
async def cdc_relations_to_review(source_id: UUID):
    """
    Retourne les relations validées manuellement qui nécessitent
    une révision suite à un breaking change CDC.
    """
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT
                id, source_entity, source_field,
                target_entity, target_field,
                relation_type, confidence, review_reason, created_at
            FROM   entity_relations_to_review
            WHERE  source_id = $1
            ORDER  BY confidence DESC
        """, source_id)

    return {
        "source_id": str(source_id),
        "count":     len(rows),
        "relations": [
            {
                "id":            str(r["id"]),
                "from":          f"{r['source_entity']}.{r['source_field']}",
                "to":            f"{r['target_entity']}.{r['target_field']}",
                "relation_type": r["relation_type"],
                "confidence":    r["confidence"],
                "reason":        r["review_reason"],
                "created_at":    r["created_at"].isoformat(),
            }
            for r in rows
        ],
    }


@app.post("/sources/{source_id}/cdc/relations/{relation_id}/reviewed", tags=["CDC"])
async def cdc_mark_relation_reviewed(source_id: UUID, relation_id: UUID):
    """
    Marque une relation comme révisée après un breaking change.
    Confirme que la relation est toujours valide.
    """
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        result = await conn.execute("""
            UPDATE entity_relations
            SET    needs_review = FALSE,
                   reviewed_at  = NOW()
            WHERE  id        = $1
              AND  source_id = $2
        """, relation_id, source_id)

    if result == "UPDATE 0":
        raise HTTPException(404, "Relation introuvable")

    return {
        "success":     True,
        "relation_id": str(relation_id),
        "message":     "Relation marquée comme révisée et toujours valide",
    }


# ── 15. Log réindexation ──────────────────────────────────────

@app.get("/sources/{source_id}/cdc/reindex/log", tags=["CDC"])
async def cdc_reindex_log(
    source_id: UUID,
    limit: int = Query(20, ge=1, le=100),
):
    """Retourne l'historique des réindexations automatiques."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT
                schema_version, deleted_count, reindexed_count,
                relations_preserved, errors, created_at
            FROM   cdc_reindex_log
            WHERE  source_id = $1
            ORDER  BY created_at DESC
            LIMIT  $2
        """, source_id, limit)

    return {
        "source_id": str(source_id),
        "count":     len(rows),
        "log": [
            {
                "schema_version":      r["schema_version"],
                "deleted":             r["deleted_count"],
                "reindexed":           r["reindexed_count"],
                "relations_preserved": r["relations_preserved"],
                "errors":              r["errors"],
                "created_at":          r["created_at"].isoformat(),
            }
            for r in rows
        ],
    }


@app.get("/sources/{source_id}/enrichment/summary", tags=["Metadata"])
async def get_enrichment_summary(source_id: UUID):
    """
    Résumé de l'enrichissement des métadonnées pour une source.
    Indique combien d'entités ont des descriptions, dépendances, stats.
    """
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        row = await conn.fetchrow("""
            SELECT * FROM source_enrichment_summary WHERE source_id = $1
        """, source_id)

    if not row:
        raise HTTPException(404, f"Source {source_id} introuvable")

    return {
        "source_id":                  str(source_id),
        "source_name":                row["source_name"],
        "total_entities":             row["total_entities"],
        "entities_with_description":  row["entities_with_description"],
        "entities_with_dependencies": row["entities_with_dependencies"],
        "entities_with_rowcount":     row["entities_with_rowcount"],
        "hateoas_graphql_relations":  row["hateoas_graphql_relations"],
        "description_coverage_pct":   round(
            (row["entities_with_description"] or 0) /
            max(row["total_entities"] or 1, 1) * 100, 1
        ),
    }


@app.post("/sources/{source_id}/enrichment/descriptions", tags=["Metadata"])
async def save_descriptions(source_id: UUID, body: dict):
    """
    Sauvegarde les descriptions COMMENT ON extraites.
    Body: { "descriptions": { "table_name": { "_table": "desc", "col_name": "desc" } } }
    """
    descriptions = body.get("descriptions", {})
    if not descriptions:
        raise HTTPException(400, "descriptions requis")

    pool = await get_pg_pool()
    enricher = MetadataEnricher(pool)
    count = await enricher.save_descriptions(source_id, descriptions)
    return {"success": True, "saved": count}


@app.post("/sources/{source_id}/generate-descriptions", status_code=202, tags=["Metadata"])
async def generate_llm_descriptions(
    source_id: UUID,
    background_tasks: BackgroundTasks,
    limit: int = 500,
):
    """Génère automatiquement des descriptions métier LLM pour les tables de la source."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        src = await conn.fetchrow(
            "SELECT id, name, connector_type FROM data_sources WHERE id = $1", source_id
        )
    if not src:
        raise HTTPException(404, f"Source {source_id} introuvable")

    from .llm_description_generator import run_description_generation_bg
    src_type = src["connector_type"] or "unknown"

    background_tasks.add_task(
        run_description_generation_bg,
        source_id   = source_id,
        source_name = src["name"],
        source_type = src_type,
        pg_pool     = pool,
        limit       = limit,
    )
    logger.info(f"[DescGen] Tâche lancée pour {src['name']} ({limit} tables max)")
    return {
        "source_id":  str(source_id),
        "source":     src["name"],
        "status":     "started",
        "message":    f"Génération descriptions LLM lancée — jusqu\'à {limit} tables",
        "check_progress": f"GET /sources/{source_id}/description-stats",
    }


@app.get("/sources/{source_id}/description-stats", tags=["Metadata"])
async def get_description_stats(source_id: UUID):
    """Retourne les statistiques d\'enrichissement des descriptions LLM."""
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        stats = await conn.fetchrow("""
            SELECT
                COUNT(*) AS total,
                COUNT(*) FILTER (WHERE description LIKE '[LLM]%') AS with_llm_desc,
                COUNT(*) FILTER (WHERE description IS NULL OR description = '') AS without_desc,
                COUNT(*) FILTER (WHERE description IS NOT NULL AND description != ''
                                  AND description NOT LIKE '[LLM]%') AS with_manual_desc
            FROM source_entities
            WHERE source_id = $1 AND is_visible = TRUE
        """, source_id)
    return {
        "source_id":        str(source_id),
        "total_entities":   stats["total"],
        "llm_described":    stats["with_llm_desc"],
        "manual_described": stats["with_manual_desc"],
        "undescribed":      stats["without_desc"],
        "coverage_pct":     round(
            (stats["with_llm_desc"] + stats["with_manual_desc"]) / max(stats["total"], 1) * 100, 1
        ),
    }


@app.post("/sources/{source_id}/enrichment/cardinality", tags=["Metadata"])
async def save_cardinality(source_id: UUID, body: dict):
    """
    Sauvegarde les statistiques de cardinalité.
    Body: { "stats": { "table_name": { "row_count": 123, "size_kb": 456 } } }
    """
    stats = body.get("stats", {})
    if not stats:
        raise HTTPException(400, "stats requis")

    pool = await get_pg_pool()
    enricher = MetadataEnricher(pool)
    count = await enricher.save_cardinality_stats(source_id, stats)
    return {"success": True, "updated": count}


@app.get("/sources/{source_id}/entities/{entity_name}/metadata", tags=["Metadata"])
async def get_entity_metadata(source_id: UUID, entity_name: str):
    """
    Retourne les métadonnées enrichies d'une entité :
    descriptions, dépendances, statistiques, formules Excel.
    """
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        row = await conn.fetchrow("""
            SELECT id, name, entity_type, description, row_count, metadata
            FROM source_entities
            WHERE source_id = $1 AND name = $2
        """, source_id, entity_name)

    if not row:
        raise HTTPException(404, f"Entité {entity_name} introuvable")

    meta = {}
    try:
        raw = row["metadata"]
        meta = json_module.loads(raw) if isinstance(raw, str) else (raw or {})
    except Exception:
        pass

    return {
        "source_id":    str(source_id),
        "name":         row["name"],
        "entity_type":  row["entity_type"],
        "description":  row["description"],
        "row_count":    row["row_count"],
        "size_kb":      meta.get("size_kb"),
        "dependencies": meta.get("dependencies", []),
        "col_stats":    meta.get("col_stats", {}),
        "excel_formulas": meta.get("excel_formulas", []),
        "formula_count":  meta.get("formula_count", 0),
        "partitions":   meta.get("partitions"),
    }


# ══════════════════════════════════════════════════════════════
# §2.2.2 — GRAPHQL RELATIONS + HATEOAS
# ══════════════════════════════════════════════════════════════

@app.post("/sources/{source_id}/relations/graphql", tags=["Relations"])
async def save_graphql_relations(source_id: UUID, body: dict):
    """
    Sauvegarde les relations GraphQL nested types dans entity_relations.
    Body: { "relations": [ { "source_type": "Order", "source_field": "customer",
                              "target_type": "Customer", "relation_type": "many_to_one" } ] }
    """
    relations = body.get("relations", [])
    if not relations:
        raise HTTPException(400, "relations requis")

    pool = await get_pg_pool()
    saver = GraphQLRelationSaver(pool)
    count = await saver.save_graphql_relations(source_id, relations)
    return {"success": True, "saved": count}


@app.post("/sources/{source_id}/relations/hateoas", tags=["Relations"])
async def save_hateoas_relations(source_id: UUID, body: dict):
    """
    Sauvegarde les liens HATEOAS détectés dans entity_relations.
    Body: { "spec": { ... openapi spec ... } }
    ou   { "relations": [ { ... } ] }
    """
    pool = await get_pg_pool()
    extractor = HATEOASLinkExtractor(pool)

    # Mode 1 : spec OpenAPI fournie → extraction automatique
    spec = body.get("spec")
    if spec:
        relations = extractor.extract_from_openapi(spec)
    else:
        # Mode 2 : relations déjà extraites
        relations = body.get("relations", [])

    if not relations:
        return {"success": True, "saved": 0, "message": "Aucun lien HATEOAS détecté"}

    count = await extractor.save_hateoas_relations(source_id, relations)
    return {"success": True, "saved": count, "detected": len(relations)}


@app.get("/sources/{source_id}/relations/by-method", tags=["Relations"])
async def get_relations_by_method(source_id: UUID):
    """
    Retourne le compte des relations par méthode de détection.
    Permet de voir combien de relations viennent de GraphQL, HATEOAS, etc.
    """
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT detection_method, COUNT(*) AS count,
                   AVG(confidence) AS avg_confidence
            FROM entity_relations
            WHERE source_id = $1
            GROUP BY detection_method
            ORDER BY count DESC
        """, source_id)

    return {
        "source_id": str(source_id),
        "by_method": [
            {
                "method":         r["detection_method"],
                "count":          r["count"],
                "avg_confidence": round(float(r["avg_confidence"] or 0), 3),
            }
            for r in rows
        ],
    }




# ══════════════════════════════════════════════════════════════
# §2.2.2A — IMPORT VUES SQL (view_parser.py)
# ══════════════════════════════════════════════════════════════

class ViewImportRequest(BaseModel):
    views: List[Dict] = Field(
        ...,
        description="Liste de vues : [{name: str, definition: str}]",
        example=[{"name": "Comptes", "definition": "CREATE VIEW Comptes AS SELECT ..."}]
    )

@app.post("/sources/{source_id}/views/import", tags=["Relations"])
async def import_views_paste(source_id: UUID, req: ViewImportRequest):
    """
    Importe les jointures extraites depuis les définitions SQL des vues.
    Les relations sont sauvegardées dans entity_relations (method=view_join).
    
    Utilisation : coller les résultats de :
        SELECT v.name, m.definition FROM sys.views v
        JOIN sys.sql_modules m ON v.object_id = m.object_id
    """
    try:
        from .view_parser import import_views_from_paste
        result = await import_views_from_paste(source_id, req.views)
        return result
    except Exception as e:
        logger.error(f"[Views] Import error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/sources/{source_id}/views/import-from-db", tags=["Relations"])
async def import_views_from_db_endpoint(source_id: UUID):
    """
    Connecte directement à la base source et importe automatiquement
    toutes les vues SQL et leurs jointures.
    """
    try:
        from .view_parser import import_views_from_db
        result = await import_views_from_db(source_id)
        return result
    except Exception as e:
        logger.error(f"[Views] Auto-import error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/sources/{source_id}/views/stats", tags=["Relations"])
async def get_views_stats(source_id: UUID):
    """
    Retourne les statistiques des jointures extraites depuis les vues SQL.
    """
    try:
        from .view_parser import get_view_relations_stats
        return await get_view_relations_stats(source_id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



# ══════════════════════════════════════════════════════════════
# §Sprint 8 — AGENTIC RAG ENDPOINT
# ══════════════════════════════════════════════════════════════

class AgentQueryRequest(BaseModel):
    question:  str      = Field(..., description="Question en langage naturel")
    source_id: str      = Field(..., description="UUID de la source de données")
    dialect:   str      = Field("mssql", description="Dialecte SQL")
    verbose:   bool     = Field(False, description="Inclure les étapes de raisonnement")


@app.post("/orchestrator/query", tags=["Agent"])
async def orchestrator_query(req: AgentQueryRequest):
    """
    Sprint 9 — Orchestrateur Multi-Agent RAG.
    Route vers DirectSQL, MultiQuery ou ReAct+ selon la question.
    Supporte les questions composées (ET / VS) avec exécution parallèle.
    """
    try:
        from uuid import UUID as _UUID
        from .orchestrator import run_orchestrator, orchestrator_result_to_dict

        source_id = _UUID(req.source_id)
        pool      = await get_pg_pool()

        async with pool.acquire() as conn:
            src_row = await conn.fetchrow(
                """SELECT id, name, connector_type, host, port,
                          database_name, username, base_url
                   FROM data_sources WHERE id = $1""",
                source_id
            )
            if not src_row:
                raise HTTPException(status_code=404, detail="Source introuvable")
            sec = await conn.fetchrow(
                "SELECT secret_value FROM connection_secrets WHERE source_id=$1 AND secret_key='password'",
                source_id
            )

        source_dict = dict(src_row)
        source_dict["password"] = sec["secret_value"] if sec else ""

        result = await run_orchestrator(
            question    = req.question,
            source_id   = source_id,
            pg_pool     = pool,
            source_dict = source_dict,
            dialect     = req.dialect,
        )

        return orchestrator_result_to_dict(result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[OrchestratorQuery] Erreur : {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/agent/query", tags=["Agent"])
async def agent_query(req: AgentQueryRequest):
    """
    Sprint 8 — Agentic RAG (legacy).
    Utilisé en fallback ou pour les tests unitaires.
    Pour le benchmark Sprint 9, utiliser /orchestrator/query.
    """
    try:
        from uuid import UUID as _UUID
        from .agentic_rag import run_agentic_rag, agent_result_to_dict

        source_id = _UUID(req.source_id)
        pool      = await get_pg_pool()

        # Récupérer la source
        async with pool.acquire() as conn:
            src_row = await conn.fetchrow(
                """SELECT id, name, connector_type, host, port,
                          database_name, username, base_url
                   FROM data_sources WHERE id = $1""",
                source_id
            )
            if not src_row:
                raise HTTPException(status_code=404, detail="Source introuvable")

            sec = await conn.fetchrow(
                "SELECT secret_value FROM connection_secrets WHERE source_id=$1 AND secret_key='password'",
                source_id
            )

        source_dict = dict(src_row)
        source_dict["password"] = sec["secret_value"] if sec else ""

        result = await run_agentic_rag(
            question    = req.question,
            source_id   = source_id,
            pg_pool     = pool,
            source_dict = source_dict,
            dialect     = req.dialect,
        )

        data = agent_result_to_dict(result)
        if not req.verbose:
            data.pop("steps", None)

        return data

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[AgentRAG] Endpoint error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/agent/status", tags=["Agent"])
async def agent_status():
    """Statut et configuration du Sprint 8 Agentic RAG."""
    return {
        "sprint":          "Sprint 8 — Agentic RAG",
        "enabled":         True,
        "max_iterations":  5,
        "max_sql_retries": 3,
        "tools": [
            "search_schema (RAG 7C)",
            "search_views (vues SXA)",
            "get_table_columns",
            "execute_sql (ConnectorFactory)",
            "validate_result",
        ],
        "supported_sources": ["mssql", "postgresql", "odata", "file_csv", "rest", "graphql"],
    }

# ══════════════════════════════════════════════════════════════
# §2.2.5 — CDC DB TRIGGERS (WAL + MSSQL CDC)
# ══════════════════════════════════════════════════════════════

@app.get("/sources/{source_id}/cdc/wal/status", tags=["CDC"])
async def cdc_wal_status(source_id: UUID):
    """
    Vérifie si PostgreSQL WAL logical replication est configuré
    pour cette source.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    if source.connector_type.value not in ("postgresql", "pg"):
        return {
            "supported": False,
            "message": "WAL CDC disponible uniquement pour PostgreSQL",
        }

    from .database import _redis
    pool = await get_pg_pool()

    # Vérifie l'état réel du slot WAL
    async with pool.acquire() as conn:
        slot_row = await conn.fetchrow("""
            SELECT slot_name, active, restart_lsn
            FROM pg_replication_slots
            WHERE slot_name = 'onepilot_cdc'
        """)
        wal_level_row = await conn.fetchrow(
            "SELECT setting FROM pg_settings WHERE name = 'wal_level'"
        )

    slot_exists = slot_row is not None
    slot_active = slot_row["active"] if slot_row else False
    wal_level   = wal_level_row["setting"] if wal_level_row else "minimal"
    wal_ok      = wal_level == "logical" and slot_exists

    return {
        "supported":    True,
        "source_id":    str(source_id),
        "slot_name":    PostgreSQLWALCDC.SLOT_NAME,
        "slot_exists":  slot_exists,
        "slot_active":  slot_active,
        "wal_level":    wal_level,
        "wal_ok":       wal_ok,
        "message":      "WAL Logical Replication actif" if wal_ok else "Configurez wal_level=logical",
        "setup_sql": (
            "ALTER SYSTEM SET wal_level = logical;\n"
            "SELECT pg_reload_conf();\n"
            "SELECT pg_create_logical_replication_slot('onepilot_cdc', 'wal2json');"
        ),
    }


@app.get("/sources/{source_id}/cdc/mssql/status", tags=["CDC"])
async def cdc_mssql_status(source_id: UUID):
    """
    Vérifie si SQL Server CDC natif est activé pour cette source.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    if source.connector_type.value not in ("mssql", "sage_100"):
        return {
            "supported": False,
            "message": "MSSQL CDC disponible uniquement pour SQL Server",
        }

    from .database import _redis
    cdc = SQLServerCDC(await get_pg_pool(), _redis)
    result = await cdc.check_cdc_enabled(source_id)
    return {
        "source_id": str(source_id),
        **result,
        "setup_sql": (
            "EXEC sys.sp_cdc_enable_db;\n"
            "EXEC sys.sp_cdc_enable_table "
            "@source_schema='dbo', @source_name='Products', "
            "@role_name=NULL;"
        ),
    }


@app.post("/sources/{source_id}/cdc/mssql/poll", tags=["CDC"])
async def cdc_mssql_poll(source_id: UUID):
    """
    Lit les changements DDL récents depuis le default trace SQL Server.
    Ne nécessite pas CDC natif activé — utilise sys.fn_trace_gettable.
    """
    source = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {source_id} introuvable")

    from .database import _redis
    cdc = SQLServerCDC(await get_pg_pool(), _redis)
    result = await cdc.poll_ddl_changes(source_id)
    return result


@app.get("/sources/{source_id}/cdc/wal/notifications", tags=["CDC"])
async def cdc_wal_notifications(source_id: UUID, limit: int = 20):
    """Retourne les dernières notifications WAL/DDL."""
    redis = await get_redis()
    try:
        key   = f"cdc:wal:notifications:{source_id}"
        items = await redis.lrange(key, 0, limit - 1)
        notifs = [json_module.loads(i) for i in items]
    except Exception:
        notifs = []

    # Fallback PostgreSQL
    if not notifs:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT event_type, object_name, payload, detected_at
                FROM cdc_wal_log
                WHERE source_id = $1
                ORDER BY detected_at DESC
                LIMIT $2
            """, source_id, limit)
            notifs = [
                {
                    "event_type":  r["event_type"],
                    "object_name": r["object_name"],
                    "payload":     r["payload"],
                    "detected_at": r["detected_at"].isoformat(),
                }
                for r in rows
            ]

    return {
        "source_id":     str(source_id),
        "count":         len(notifs),
        "notifications": notifs,
    }


class NLURequest(BaseModel):
    question:        str
    source_id:       Optional[str] = None
    conversation_id: Optional[str] = None

@app.post("/nlu/analyze", tags=["NLU"])
async def nlu_analyze(req: NLURequest):
    """
    Analyse une question en langage naturel.
    Retourne intent, entités extraites, slots remplis.
    """
    import time
    t0 = time.time()

    nlu     = get_nlu_pipeline()
    context = get_context(req.conversation_id or "default")

    # Récupère les entités connues de la source
    known_entities = []
    known_fields   = {}
    if req.source_id:
        try:
            source_id = UUID(req.source_id)
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                rows = await conn.fetch(
                    "SELECT name FROM source_entities WHERE source_id=$1 AND is_visible=TRUE ORDER BY name LIMIT 2000",
                    source_id
                )
                known_entities = [r["name"] for r in rows]

                # Récupère les champs pour les tables détectées
                field_rows = await conn.fetch("""
                    SELECT se.name AS table_name, ef.name AS field_name
                    FROM source_entities se
                    JOIN entity_fields ef ON ef.entity_id = se.id
                    WHERE se.source_id = $1 AND se.is_visible = TRUE
                    LIMIT 1000
                """, source_id)
                for r in field_rows:
                    known_fields.setdefault(r["table_name"], []).append(r["field_name"])
        except Exception as e:
            logger.warning(f"[NLU] Entities fetch error: {e}")

    # Pipeline NLU
    slots = nlu.process(req.question, context, known_entities)

    # Détection ambiguïtés
    resolver   = AmbiguityResolver()
    questions  = resolver.analyze(slots, known_entities, known_fields)
    clarif     = resolver.build_clarification_response(questions)

    # Log dans DB
    ms = int((time.time() - t0) * 1000)
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.execute("""
                INSERT INTO nlu_query_log
                    (conversation_id, source_id, question, intent, confidence,
                     tables_detected, slots_json, needs_clarification, clarification_json, response_ms)
                VALUES ($1,$2,$3,$4,$5,$6,$7,$8,$9,$10)
            """,
                req.conversation_id or "default",
                UUID(req.source_id) if req.source_id else None,
                req.question,
                slots.intent,
                slots.confidence,
                slots.table_names,
                json_module.dumps({
                    "metric":       slots.metric,
                    "group_by":     slots.group_by,
                    "date_filter":  slots.date_filter,
                    "amount_filter":slots.amount_filter,
                    "top_n":        slots.top_n,
                }),
                clarif.get("needs_clarification", False),
                json_module.dumps(clarif) if clarif.get("needs_clarification") else None,
                ms,
            )
    except Exception as e:
        logger.warning(f"[NLU] Log error: {e}")

    return {
        "intent":      slots.intent,
        "confidence":  slots.confidence,
        "nlu_method":  getattr(slots, "nlu_method", "hybrid"),
        "tables":      slots.table_names,
        "metric":      slots.metric,
        "group_by":    slots.group_by,
        "date_filter": slots.date_filter,
        "top_n":       slots.top_n,
        "ambiguities": slots.ambiguities,
        "needs_clarification": clarif.get("needs_clarification", False),
        "clarification":       clarif if clarif.get("needs_clarification") else None,
        "response_ms": ms,
    }


# ══════════════════════════════════════════════════════════════
# §2.3.3 — SQL GENERATOR
# ══════════════════════════════════════════════════════════════

class SQLGenRequest(BaseModel):
    question:        str
    source_id:       str
    conversation_id: Optional[str] = None

@app.post("/nlu/generate-sql", tags=["NLU"])
async def generate_sql(req: SQLGenRequest):
    """
    Génère du SQL depuis une question en langage naturel.
    Pipeline complet : NLU → slots → SQL Generator.
    """
    import time
    t0 = time.time()

    source_id = UUID(req.source_id)
    source    = await get_source(source_id)
    if not source:
        raise HTTPException(404, f"Source {req.source_id} introuvable")

    # Récupère le schéma
    pool = await get_pg_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT se.name AS table_name, ef.name AS field_name
            FROM source_entities se
            JOIN entity_fields ef ON ef.entity_id = se.id
            WHERE se.source_id = $1 AND se.is_visible = TRUE
            ORDER BY se.name, ef.position
            LIMIT 2000
        """, source_id)

    schema: dict = {}
    for r in rows:
        schema.setdefault(r["table_name"], []).append(r["field_name"])

    known_entities = list(schema.keys())

    # NLU
    nlu     = get_nlu_pipeline()
    context = get_context(req.conversation_id or req.source_id)
    slots   = nlu.process(req.question, context, known_entities)

    # Corrige intents mal classés
    slots = await _fix_slots(slots, req.question, schema, known_entities, source_id_str=req.source_id, pg_pool=pool)

    # Vérifie ambiguïtés bloquantes
    resolver  = AmbiguityResolver()
    questions = resolver.analyze(slots, known_entities, schema)
    if questions and questions[0].required:
        clarif = resolver.build_clarification_response(questions)
        return {
            "sql":         None,
            "explanation": "Clarification nécessaire avant génération SQL",
            "needs_clarification": True,
            "clarification":       clarif,
            "intent":      slots.intent,
            "response_ms": int((time.time() - t0) * 1000),
        }

    # Dialecte SQL selon le type de connecteur
    dialect_map = {"mssql": "mssql", "sage_100": "mssql", "mysql": "mysql",
                   "postgresql": "postgresql", "sqlite": "sqlite"}
    dialect = dialect_map.get(source.connector_type.value, "sql")

    # Génère le SQL
    sql_gen = SQLGenerator()
    result  = sql_gen.generate(slots, schema, dialect)

    # Sauvegarde dans le log
    try:
        async with pool.acquire() as conn:
            await conn.execute("""
                UPDATE nlu_query_log SET sql_generated = $1
                WHERE conversation_id = $2 AND question = $3
                ORDER BY created_at DESC LIMIT 1
            """, result["sql"], req.conversation_id or "default", req.question)
    except Exception:
        pass

    ms = int((time.time() - t0) * 1000)
    return {
        "sql":         result["sql"],
        "explanation": result["explanation"],
        "warnings":    result.get("warnings", []),
        "intent":      slots.intent,
        "tables":      slots.table_names,
        "dialect":     dialect,
        "needs_clarification": False,
        "response_ms": ms,
    }


# ══════════════════════════════════════════════════════════════
# §2.3.3.C — UNIVERSAL QUERY PLANNER
# ══════════════════════════════════════════════════════════════

class QueryPlanRequest(BaseModel):
    question:   str
    source_ids: List[str]

@app.post("/nlu/query-plan", tags=["NLU"])
async def create_query_plan(req: QueryPlanRequest):
    """
    Crée un plan d'exécution cross-source pour une question complexe.
    Décompose la question en sous-requêtes atomiques.
    """
    pool = await get_pg_pool()

    # Charge les sources
    sources = []
    schemas = {}
    for sid in req.source_ids[:5]:
        try:
            src = await get_source(UUID(sid))
            if src:
                sources.append(src.model_dump())
                # Charge le schéma
                async with pool.acquire() as conn:
                    rows = await conn.fetch("""
                        SELECT se.name AS t, ef.name AS f
                        FROM source_entities se
                        JOIN entity_fields ef ON ef.entity_id = se.id
                        WHERE se.source_id = $1 AND se.is_visible = TRUE
                        LIMIT 500
                    """, UUID(sid))
                schema = {}
                for r in rows:
                    schema.setdefault(r["t"], []).append(r["f"])
                schemas[sid] = schema
        except Exception as e:
            logger.warning(f"[QueryPlanner] Source {sid}: {e}")

    if not sources:
        raise HTTPException(400, "Aucune source valide")

    # NLU
    nlu     = get_nlu_pipeline()
    known   = [t for schema in schemas.values() for t in schema.keys()]
    slots   = nlu.process(req.question, None, known)

    # Query Planner
    planner = UniversalQueryPlanner()
    plan    = planner.plan(slots, sources, schemas)

    # Sauvegarde le plan
    plan_id = None
    try:
        async with pool.acquire() as conn:
            row = await conn.fetchrow("""
                INSERT INTO query_plans (question, source_ids, plan_json, status)
                VALUES ($1, $2, $3, 'ready')
                RETURNING id
            """,
                req.question,
                [UUID(s) for s in req.source_ids],
                json_module.dumps([
                    {
                        "step_id":     s.step_id,
                        "source_id":   s.source_id,
                        "source_type": s.source_type,
                        "action":      s.action,
                        "query":       s.query,
                        "depends_on":  s.depends_on,
                        "parallel":    s.parallel,
                    }
                    for s in plan.steps
                ]),
            )
            plan_id = str(row["id"])
    except Exception as e:
        logger.warning(f"[QueryPlanner] Save error: {e}")

    return {
        "plan_id":        plan_id,
        "question":       req.question,
        "intent":         slots.intent,
        "steps":          [
            {
                "step_id":     s.step_id,
                "source_id":   s.source_id,
                "source_type": s.source_type,
                "action":      s.action,
                "query":       s.query,
                "depends_on":  s.depends_on,
                "parallel":    s.parallel,
            }
            for s in plan.steps
        ],
        "merge_strategy": plan.merge_strategy,
        "explanation":    plan.explanation,
        "estimated_ms":   plan.estimated_ms,
    }


# ══════════════════════════════════════════════════════════════
# §2.4.3 — DASHBOARD GENERATOR
# ══════════════════════════════════════════════════════════════

class DashboardRequest(BaseModel):
    question:  str
    source_id: str

class DashboardDataRequest(BaseModel):
    sql:       str
    source_id: str
    limit:     int = 200

# ─────────────────────────────────────────────────────────────────────────────
# EXPORT POWERPOINT
# ─────────────────────────────────────────────────────────────────────────────

class PptxMultiRequest(BaseModel):
    specs_json: str   # JSON array de specs
    title:      str = "Tableau de bord consolidé"

# ─── PARTAGE PAR LIEN ────────────────────────────────────────────────────────

class ShareRequest(BaseModel):
    spec_json:  str
    title:      str
    source_id:  Optional[str] = None
    expires_in: int = 24   # heures

@app.post("/api/share")
async def create_share_link(req: ShareRequest):
    """Crée un lien de partage temporaire pour un dashboard."""
    import uuid, hashlib
    from datetime import timedelta
    try:
        # Token stable basé sur le contenu — même dashboard = même lien
        content_hash = hashlib.md5(req.spec_json.encode()).hexdigest()[:16]
        token = content_hash
        expires_at = datetime.utcnow() + timedelta(hours=req.expires_in)
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            # Upsert — si le token existe déjà, on met à jour l'expiration
            await conn.execute("""
                INSERT INTO shared_dashboards (token, title, spec_json, source_id, expires_at)
                VALUES ($1,$2,$3,$4,$5)
                ON CONFLICT (token) DO UPDATE
                SET expires_at = EXCLUDED.expires_at,
                    title      = EXCLUDED.title,
                    view_count = shared_dashboards.view_count
            """, token, req.title, req.spec_json, req.source_id, expires_at)
        base_url = os.environ.get("BASE_URL", "http://localhost:8000")
        share_url = f"{base_url}/shared/{token}"
        return {"token": token, "url": share_url, "expires_at": expires_at.isoformat(), "expires_in_hours": req.expires_in}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/shared/{token}", response_class=HTMLResponse)
async def view_shared_dashboard(token: str):
    """Page HTML lecture seule pour un dashboard partagé."""
    import json
    from datetime import timezone
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            row = await conn.fetchrow("SELECT * FROM shared_dashboards WHERE token=$1", token)
            if not row:
                return HTMLResponse("<h2>Lien invalide ou expiré.</h2>", status_code=404)
            if row["expires_at"].replace(tzinfo=timezone.utc) < datetime.now(timezone.utc):
                return HTMLResponse("<h2>Ce lien a expiré.</h2>", status_code=410)
            await conn.execute("UPDATE shared_dashboards SET view_count=view_count+1 WHERE token=$1", token)
        spec = json.loads(row["spec_json"])
        title = row["title"]
        expires = row["expires_at"].strftime("%d/%m/%Y %H:%M")
        spec_escaped = json.dumps(spec).replace("</", "<\/")
        html = f"""<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title} — OnePilot</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<style>
  *{{box-sizing:border-box;margin:0;padding:0;}}
  body{{font-family:'Inter',system-ui,sans-serif;background:#f4f6fa;color:#1a2b3c;min-height:100vh;}}
  .header{{background:linear-gradient(135deg,#006eb8,#2d72ff);padding:16px 24px;border-bottom:1px solid #d0e4f4;display:flex;align-items:center;justify-content:space-between;box-shadow:0 2px 12px rgba(0,110,184,.15);}}
  .logo{{display:flex;align-items:center;gap:10px;font-size:18px;font-weight:700;color:#fff;}}
  .badge{{background:rgba(255,255,255,.2);border:1px solid rgba(255,255,255,.3);color:#fff;padding:3px 10px;border-radius:20px;font-size:11px;}}
  .meta{{font-size:11px;color:rgba(255,255,255,.7);}}
  .main{{padding:24px;max-width:1400px;margin:0 auto;}}
  .dashboard-title{{font-size:22px;font-weight:700;margin-bottom:16px;color:#1a2b3c;}}
  .widgets-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(400px,1fr));gap:16px;}}
  .widget{{background:#fff;border:1px solid #d0e4f4;border-radius:12px;padding:16px;box-shadow:0 2px 8px rgba(0,110,184,.06);}}
  .widget-title{{font-size:13px;font-weight:600;color:#4a6a82;margin-bottom:12px;}}
  .chart-container{{position:relative;height:260px;}}
  .kpi-card{{text-align:center;padding:20px;}}
  .kpi-value{{font-size:42px;font-weight:800;color:#006eb8;}}
  .kpi-label{{font-size:13px;color:#4a6a82;margin-top:8px;}}
  .kpi-delta{{font-size:13px;margin-top:6px;}}
  .table-wrap{{overflow-x:auto;max-height:280px;}}
  table{{width:100%;border-collapse:collapse;font-size:11px;}}
  th{{background:#006eb8;color:#fff;padding:6px 8px;text-align:left;position:sticky;top:0;}}
  td{{padding:5px 8px;border-bottom:1px solid #e8f0f8;color:#1a2b3c;}}
  tr:nth-child(even) td{{background:#f0f6ff;}}
  .expired-banner{{background:rgba(192,48,48,.1);border:1px solid rgba(192,48,48,.3);color:#c03030;padding:8px 16px;text-align:center;font-size:12px;}}
  footer{{text-align:center;padding:20px;color:#9ab4c8;font-size:11px;border-top:1px solid #d0e4f4;margin-top:30px;}}
</style>
</head>
<body>
<div class="header">
  <div class="logo">
    <img src="data:image/png;base64,/9j/4AAQSkZJRgABAQAAAQABAAD/4gHYSUNDX1BST0ZJTEUAAQEAAAHIAAAAAAQwAABtbnRyUkdCIFhZWiAH4AABAAEAAAAAAABhY3NwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAA9tYAAQAAAADTLQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAlkZXNjAAAA8AAAACRyWFlaAAABFAAAABRnWFlaAAABKAAAABRiWFlaAAABPAAAABR3dHB0AAABUAAAABRyVFJDAAABZAAAAChnVFJDAAABZAAAAChiVFJDAAABZAAAAChjcHJ0AAABjAAAADxtbHVjAAAAAAAAAAEAAAAMZW5VUwAAAAgAAAAcAHMAUgBHAEJYWVogAAAAAAAAb6IAADj1AAADkFhZWiAAAAAAAABimQAAt4UAABjaWFlaIAAAAAAAACSgAAAPhAAAts9YWVogAAAAAAAA9tYAAQAAAADTLXBhcmEAAAAAAAQAAAACZmYAAPKnAAANWQAAE9AAAApbAAAAAAAAAABtbHVjAAAAAAAAAAEAAAAMZW5VUwAAACAAAAAcAEcAbwBvAGcAbABlACAASQBuAGMALgAgADIAMAAxADb/2wBDAAUDBAQEAwUEBAQFBQUGBwwIBwcHBw8LCwkMEQ8SEhEPERETFhwXExQaFRERGCEYGh0dHx8fExciJCIeJBweHx7/2wBDAQUFBQcGBw4ICA4eFBEUHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh7/wAARCAH0AfQDASIAAhEBAxEB/8QAHQABAAEFAQEBAAAAAAAAAAAAAAUBAgMEBgcICf/EAEgQAAIBAwIEAwQIBAMECAcAAAABAgMEEQUhBhIxQVFhcRMigZEHFDKhscHR8BUjQlJy4fEIM2LSFiRDU4KTosIXNGODkpTi/8QAGgEBAAMBAQEAAAAAAAAAAAAAAAMEBQIBBv/EADQRAQACAQIEAgkDBQEAAwAAAAABAgMEERIhMUEFEyIyUWFxgZGh8BRCUiOxwdHhMwZi8f/aAAwDAQACEQMRAD8A+MgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAy2tVUa8ZuMZx6Si1nK7gYgTtXTrW5iqtJqnldYdH3zh/ht26EdW065ptcsVVz2hu/THw7EtsF6xvtyR1y1tO27TBVpp4aw0UIkgAAAAAAAAAAAAAAAADZtbK4uN4R5Y4zzS2X72Ja20i3i4ubnUa652Xbt5b9ybHgyZPVhFkzUx+tKDpUqlWXLSpzqS8Ixyzco6Te1En7NQysrmf4+HxOlt6VKCUKcIxWcpJbZ77eJtwgupfx+G7+vb6KGXxHblWHPU+HJuKcrlc2d0o7Y8n1+5G5R4atMp1K9eSx0i0t/kyVqV6FCLlVqRgl1cmlg1rjW9OoRlmsqkotYjD3ub0a2+bJp02kxetP3V/wBTq8nqR9llDh7To/bpup/ik1+DMz4f0pra2S8+ef8AzGhU4ot/Zp07ep7TupJY/Exx4rXs8SsW556qthY9HF/iczk0MdnsY9fPf+zclwvYSTaqV4vtyyWF89zBU4SjzR9neNR7uUFnHks7mKPFbT3sm14e2X/KZY8WwdX3rOVOn4KfNL57HM20Fu231dcPiFekxP0aN1wxqFLmlSlTqxXTDxL5f5kbd6Zf2ufb2lWKX9SjmPzWx1NLirT5VYxlSrxg+spRW3wTyStvrWkV019co8rlhe0fLn4Pt8jz9Jpcn/nk2+P5Dr9Xq8X/AKY9/h+S82B6ZeaJpV/zudvBVJbucdpZaW/+uxzuo8G3MFzWFaNZf2T2fbo+hBl8Nz0jiiN49ybD4ngyTwzPDPvcqDJcUK1vUdOvSnSmu0o4MZQaIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAJLR7rkl7CbSi902/u/fmS+xy8W4yUotpp5TXVE/p9yrqhzNr2kcKa/P479P0zoaPPt/Tt8lLU4v3wy1qFCqmqlKMsrGe68Md/2smhW0mL3o1uV4zia2eF4rf0WPiSLzkFrJpsd+3NXpqL1c/cWtxbpSq0pRi3hS6xb8M+PkYDqG0atxYW1aUpuHJJ7+5iKXwwUsmivHq81umqrPrckCCQr6XXi/5L9sn0SWJP4fozRnCdOXLOMoyXZrDKlq2rO0ws1tFo3haADl6AAAASOn6e6mKtxlQ6qHeX6I6pSbztVza0VjeWpbW1W4linH3U8OT6L9+HUmbTT6FBxm17Wa6uazH4L9cmxShCEFCEVGK8EXmnh0da8785UMuqtPKq9Nyee7+8yJpYXjjHmyPvr6lbJxzzVcbRXb1/H5EPd31xc5jKfLTbyoLov1Osuspj5V5y4x6W9+duUJ6vrFrb5jHNWaeMRxj5/wCpFXWtX1bMYT9jF5WIdcepGgoZNVkydZXcemx4+kLqlSdSXNUnKcvGTyy0ArpwAAAAAAAGzaX15aSUre5qU8dEpbfLodBpXGN1QxG8oqvHbEovEvP1OWBNi1GXDO9LbIc2nxZo2yV3ep22ocP8RUPq9XkqNb+zqe7Jbbtfnj/M5ziDgm5oN1tKcrinu/ZP7UfTxOQTaeU8NHSaDxhqOnyVO5bu6Gy5ZS95LZbP08S9Grw6jlqK7T/KP8s79Fn0s8WmtvH8Z/xLnJwlTnKE4uMovEotYafgy09NuLLh/i+1de1qeyu1lcyWJxx/dHuv3k4LXNIvtHu3b3lLH9k47xmvFMrajSWwxxRPFWekx0/4t6bW0zzNJjhvHWJ6/wDYR4AKi4AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABmtLidtXVWG+OsX0kvBmEDoOno1KdamqtKWYS+a8n5/tdS5nPWV3UtZvlbcJNOUc7PH+r+Z0FKcKtGNWnOMoyXZ9H3Xjn99DX02ojJHDPVm58E0neOijWAGC2rKxKVI05wcZwhJN53inh+Kz++gBxasWjaYdRaazvEtC40ulJN0ZOEnulJ5j6dM/Hc0bixuaDWYc8Wm1KG626+nx9ehOlVJprGdmmvIq30VLeryWaaq0etzcuDoK9nbVYtOnGLbbyljfx29OnQ0rnS2pN28+aPZS6+mf1wUsmlyU7brVM9LM2n2CopVbiKc2vdi9+Xo848fmb6jt4Ig43F7aNRm24pcqU/eWE+ifbfwN+21ShNKNaMqU+jl1i/PxXyZY02fFSNpjaUGfDkvO8TybspKPdJfJEXf6lJ5pW7ws4c/H08PUwajfSuJOFPMaSfT+71NIj1Gqm/o16JMOninO3UABTWQAAAAAAAAAAAAAAAAAAZbS5r2lxC4tqsqVWDzGUWeg6HxJpnEdktG4jhGFaTxCt0jJ9nn+mXXyfxPOQWMGpvh325xPWO0q2o0tM8RvymOkx1hO8X8N3fD95y1E6ltUb9lV/J+DII7DhbielUt1oXESVfTqqcI1ZbypZ6bvsvHqiI4r0Gtol8oqXtrSsue3rReVOL6b9MnubHSY8zF09ns/57Jc4Mt4nys3rdp7T/wB9sIYAFZbAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAzW1xVt5uVOWM9U+jMIHQStLVYtJVaDTzvKEtkseD7/E2qN5b1XGMa0OZptqXu4+L2z8yEoUalaTVNJtLLy0tviKlGtSSdSlOCfRyi1ks01WWvdBbT47dnSNYeGmn4dPMoc7QuK9DPsqsoKWOZJ7PHTK6M26OqVY4VaEaiy22vdb8vD7izTXR+6EFtJPaUuUTNSnqNrJRzKdN43UllJ+q6r4I2qU41E3TlGWFl8u/bPb9/Es0z479JQWw3r1hcAUbJUOy7ZLomtsprKePFdGaVxp1vU3jmk293FZXnt/obe5QiyYqX9aEtMtqdJQ1XTbmOXCHtEv7evft16Lsab2eGdOi2tSo1lGNWnGaT2z19M+HkVMmhnrSVqmrifWhzQJa40qDWaFVp5+zNbY8cr9CPubWvbtqrTaS/qW6+a2KV8dqetCzW9bdJYQAcOwAAAAAAAAAAACsIuc4xXVvCAoDflpVdY5alKeVnbm/NBaVXcXJ1KSx23y/uwSeTk/jLjzae1oAlaOkxkv5lxOHpST/wDcZ1otpje/r/8A6y/5zr9Pk/i587H7UGdLw7rNvVs3oOuS59Pq7Uqz3lbS7SXkaVXR6SX8q7qTfnRS/wDczBU0m4jFSVSlLK6LOV934ZOqY8uOd4hze2LJG0yx63ptfStRqWdwsuO8JrpOD6SXk0aRIXsdRdrRoXMHOFFNwksScU98NrovJ9NyPIbxtKWkzMcwAHLoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAX0ak6VWNSDxKLyjobK6hWh7SlJxljEknvFeH3fvtzZns7idtWVSPvL+qOWkybBl8u3Poiy4+OPenZW1rKWZ29Oa7rpnbbdb/fuak9Lt3J+/OCb7b4814/P4m5Rq061NTpzUk0u3R+HwKvBpzhxZI4ojqoebkpO0yiqmlVFOSp1YSils5Jpvw2WTDWsLy3qqLp+/lY9nNT37Y5WybyXIinRUnpKWurtHWEFSvrqhPE/fw23Gpnr59H95s0dTj7NqrTkpJbOPR/DsStSXPBU5tuKTST7LwSNapZ2s171CGdveTa2XkjiNPmp6lnXn4r+tClK6t6kVJVaaztiUuVp+D/UzpfNGhPS6Uk3CpKD8Mcy/X8Sisr6jh2ty5Y2UXLl269Htjyz8DuM2anr13czjxW9WyRKZNFXd7S964tIzhjMnHZpZ67ZS+K9C6lqNrKHvOpGeejisP45/I7rq8czz5S4nTXjnHNuphSazu9016prDLISjOmqkJwlF/wBs02vVLp8S7bxJ4mt45Tuhms1nm1rmwt6qXLBQl4x2z+X3GlX0uvBSlSlGpFeai8eOP0ZK9C5MgvpMd+nKfcmpqb1683NVITpzcKkZRkuqaw0WnT1IwqpRqxVSKWMS32fh+8kfcaZSeZUpODztHGU/0Kd9Hkr05rVNTS3XkiAbFzZ17fLnFOOftRec+fjg1yrMTHKU8Tv0AAePQAADa0um6l/SfK5RhJTn6J/tfE1ToNPtvqtHleHOSTm0/u+H77E2DFOS+yPLkild2ytlhLYowDbZCkpKKzJ49Sz29JPHtI/NF7TaxnbwMbt6XMpeyi5J5zg5ni7O68PdnoxnVpupTpznBPDlGLazjPb1KOcYvEm0/BowulTdRVHCMpL+5Jr5PY1L3T4VFzW8Iwl4ZwpfPp9yIL3y0jeIiUtKY7TtvMJNSxLbKcX8U0atzZ29eMuaGJvfnX2s+fj26kdaw1Km1GEJ8vRRm8J+mfyJRyrqnTnKnScpZ5owk/d89/yb7nFctc8bXr9nc4rYudbIG6tqlvPE0+Vv3ZY6mE6avThVi4TjGcX6/Mj6+lRe9Ko4+TWV8/Aq5dJes715wnx6mto58pRIMlelUo1OSpHEvXJjKiyAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAM9ncztqmY7xf2o5J6jONajGrT3g9s+D7p+fTb0OaM1pc1LapzQw1lc0X0l6ljBqJxT7kOXDGSPe6BoIxWd5RuUlGXLNvHJLr8PH4fEzrDWU+ZeKNWmSl43rLOvjtSdpha2UwXPGSmTtwp0KplH1AFcvKaymnlPui2UYyTjKKkm28PdZx59yoObVieUuotMdGu7G1zJxhyt9MSfu+hV0LpTWLucoL/vFzSSxss+HxM+SuSKdPTffp8EkZ7xy6sdSVaNVKVKDTaTdFtrt2lhr4svhLmqOCTim/dlPEY4/xdPky7JSSi4uMkmn1T7nsY716W+p5lbda/QjLKTXf4DJVTnGUpRk1KUeVy7tYxjftj4FkViWZPmWMdcPrnPhnGUvwOotfvDma17Sv2fYwXFpQrtSqQ3Wd08N/r9/bsXx9quqjJb9zK47LDTydWpTJG1oeVtak7xKHr6ZUjHmpT5/FNY/fxwaValUoy5akJRfmuvp4nSbhqLTUoxlFvdNZT69V0ZTyaGOtJWaauf3Q5gE3X0y3qbwbpPHbdfI0K1hWoTi5x9pTzu6Ty8d9uq+KKWTDfH60LVMtb9JZ9Eteef1ma92LxDf+rx+H76EtjBgsbijOnGNJw2X2Y7PZdceS7mdvJqaWla09Gd/az9Ra1rc42UABZQBSXQuSbeEZKVCVSaiurERM9HkzEc5YRsSlS1s7eGa1Rylj7KNCtKk3inBRXm8s9mk1c0yRboxLboVyY6s4whKc5KMY9W33/eemTQrarCMl7Ck5pPOamyfwX6lbJqKY+U9VmmC9+cJPOwbI3TZ6pqd27SxhCVSab5OVPZebyzVq3d9Rqzo1J8s6cuWScFlNbY6eRDOtrtvwpY0s77bpDV6Sq2rkk26eWl4dM/gQZKUdSTtqsaqxUcJRXhLMcfPfPh6Y3iylqL1vbiqt4a2rXawACBKAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAS+n2dvUs4uquZzTbcZYcXlpL7s9Oj+K7x47ZJ2q5taKxvKIM0bu6jTdONzWUHs4qbw/gZr2wnb01UU1OLeHhPKNM8tWaztL2LRaN4bVC/u6LzGrzeHtIqaX/5JmShqdal1pUqv+Pm/JoyaZCyrQ9nUpJ1c7Zct1jyfqbUtMt+yUfi/wBSfHjy3jek/dDkvjrO1o+zTp6pUT/mUoSWf6W0/TuvuM38UpSqf7qpTp+GVN/Pb8DI9LodpNfMp/CqP98iaKamO6Kb6eR6na+0xFVuTxcVn5ZL1e2vsnN14c3aniXN+GPvMf8AC6P98i7+FUP+8mdR+phzM6eWancUZxc3WpRS7Sms/IyQnCpjknGbfSMJKT+S6GtDSrXni5zqyj3UWk/nh/gZnpmldqd7/wCfH/kO4tqI6xEuZjBPdlqNU6jp1HyTXWMtmu/Qq1hb5SfiYY6dYR+x9dj/APej/wApd9VpxkpQuL2LXR/WN1/6TuL5u9Pu4mmLtb7L8LHQo2vAs9him6catXEurbjzfPlz95bC3nDpVk/8UmySL371n7POCn8mVMqmi2EakebMaM8rCcubZ+WGjFGlcxzmtTl5cmDqtp/i5mkfyZ0VKUoS/wC0mn/hj+rMzVHGzn8kSI5Yimxc0k9nleaLHsHjFVtrarLM4YfXmg+WWevX9cllJXdJ7p3Sk+2faL8pZb836GzsNlFSeyfd9CtfBXfij0Z9yxXNbbaecK21SNem503lLaS7xfn3X+XxMjia8o0as1KLXtIr3ZQwmn5+K8mn3MlGclJQmpNZUVOMXhvtnwb/AB7ntM01na/17f8AHN8UTzp9GSDSwvFk3OjTtbH2md2t34vwRDuKe6K1K1WdKNKU24Rey7IvY7xHZRyUnJMc9mKpNyk231MbTk+Vde26W/q9i5o2aEbSWkalGvh13Cn7FZayuZKWPh92SC8ztus0iN4hy9/czuK2ObNOGVBLOMePxNYyXNKVGtKm98dH4oxmBMzM82zG23JLcI3bsuJLGuo8ydVQkvKW355+Bv8A0j2H1LiSpVhDlp3K9qmlhOX9X6/E5ynOVOpGpB4lFpp+DR6L9JtrG60Cw1SgnKEFHMn15JLZ/Nr5l3DXzNNkj+O0/wCJZ+e/lavHP8t4/wAw84ABRaIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHSWqirSk4rH8uLfrhZObOo9m6KdGT/3b5M48Hj8EXdDHpzPuVdXPoQqo5yn0aa9U01+BBaraq1rrkfuTy4rwJ2l9oiuIc+2pZ/tZY1tImnF3V9LeYvw9kWtnlEnYai0lTuZbLpN9vXG7IwGbTJak71aFqRaNpdRFppNNNdmuj80UZAWd5XtZZptOPeMllMlLXUKFZtVJKg8L7e6b77pfj9/fSxaytuV+UqGTS2jnXm3AUX2IzX2ZfZl2fp4jOC3FonorTExylUZK04yq1FTpxlObWVFLLePL7zHVnGlJxqSjCS7SePxweTesdZexSZ6QyIo2Ylc2mcVLulTXK2nvLPl7qbT9TFK5bS9lRq1VJtQlGD5ZPvjOCOdRjjlxO4wX67Nr4FMinze45U5S5otOLlycj23ziWV1L6dNLeWZZ7S2SWN+nV53z9x1XJxdIl5NOHrMLUVw+xmqRjJwfKk4RUU1tnHd9s46leVttktd+6G1ohgw12K7+BnVNeBVUn2R1s4nJDCotroUVLLy+nfBsun5FHHboOF55m6xxoLDVLLW65pdN87YwWT9m5uf1e1k315reDy/iiW0Xh/V9YqNafZVasE8SqNctOHnKTwkvVk7/AeGNFlJa/rMr6uln6tpWJrPhKrL3V54TwczirfnNfqjtrq454YmZn2Q4ynbTvasbejY0KlSbxGNK1hzP5RMtzp1zZVY0ru3nTnSf2KkcNNdE++3gejcL/SVS4Z1ChX0DhjSbGNKW9SUHVrSXdOcntt4JEX9J3FkeLuIa+r/AFaFCVd80oxH6fHtPKPohp4hq7ZYrakxX47uKoUoUoKMXLGHs8dfLGNvmUksoysRjzNJd9j2tYrG0LdskzPFLWeUWM9s+jv6BeIOMOF563Qq06FJwcqMZptzx6Lo/i/I8r4t0G94d1u50u+punXt5uE1v1T6ry7oTH5/tDg12LNfgrPP+/wcxrNHntlVjHLpvd435X4+j/Ehzorul7W2qrb3acpfKLf5HOmNq6RXJO3dv6a3FQPUb6hK4+jGkptTlC0jUj5YWfw/A8uPXLeapcAW1SdNTjG0ptxfSWy2fqW/C68U5K+2ss7xeeGMVo7Xh5GCslhtPsUMpsAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAdff/APz9wv8A6svxZyB1VWoqterVWcTnKS9G8l7QetKnrPVgpLM0iN4kjh0HjqpLPyJKj9tGtxBSUrHmUW5Qkpc3ZR3X4tFzVxvhlV007ZoQ+mKlK9hGtFTg1JYfd8rx9+Dau9KqQxOhJTT6xbw18/8AX8TRs+b63RUVmXtI4Xi8nSuWUmUNNgrlrO/Vdz5bY5iY6OYqU6lNpVISjlZWVjK8Sw6qL5ZKUXhp5TXVeZc61dvLrTz/AIiSdBbtLiNZXvDl6M6sZJUZzjJvblbzk3IUdVzJRVanKP2lzckvlsyblKU5uc25Saw5N5bQzsjuuhmOtvo5trI7Qio6Zc1XCNzdYj6uXL+/Iy0dJoRS9pKdR533wnt4fnn4EhuVjlk1dFir15obau8+5ip21vSUVTowXLutt+nd9fyMq656v7yqizJCn5FmuKtPVjZWvlm3Wd2Nptl8IMzwpI3L76k/ZfU6NWklSiqntKik3US95rCWE30W+PFkm3dBORoqPgVUGy/BkpxTayew4mzJa2zqvlS+J1ugcEatqtNytbOtViurjBv8DU4OsoXOpUac1lSml08WkfoP9GOg6ZovCVhTsrenFzoRlOSisybW+TjPmrgpx2jfdnx52r1H6fDaK7RvM/8AH5/alwXrFG9hZwsa87iTwqag8v4fmXw03QeGakpary6xqEcONrRn/Ig/CclvJrwjt4s+xf8AaJ0GlPQqlzpsJQvK8eRxorEppPptv4/M+NOJNE1Kwqy+s2FzSTfWdNr72ialq5MUZKx1Z9c2aNTfSZ7x6M7cu+8b/L85tHiDifVdXg7epVjbWSfuWdsvZ0Y+GIrZ+ryzn5R7s2qsJ5zyvPoYpRx1RzPOW3hrTHXhpGzA1gt38TPGm5NRSbfkZVYXM6kKcKNSU5yUYRUW8tvCwu7bPOqTjrHWWnjPRZOm0nhlUbSGqcQ3cdLsniUIyTdesv8Agh1afi8L5G/TpafwZRVW8p0b/X5RThQklKlaJ/3rvUX9vRdzktWv73VLypeX9xUuK83mU5vOfLyXktj2YivOUEZL6neMfKvt7z8P9vof6N/9oWy4S4fno1vptepb0Y8ls6lROUo9ubol8M7eh4Z9I/E1bivia71m4UVUuajm4x6Jdl8Fg5uT7NmPJHMxvM7bTPVLpvD6YJiYmZ2323npv1WV5KFCs33o1F84tHNE9qjxYzlnC2j6t9vxfwIEyNbO+R9DpY2oHq7pzo/R5Qpzbcnbw6vs2mjyy2pSr3FOjD7VSaivVvB6/wAVyjQ4djDlUM8iUV0WF0Re8Hryy2/+u31ZfjV/Sw09tt/o8guFi4qLwm/xMZkuXm4qvxm/xMZit0AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADpLba1pLmz/Lj98Vsc2Tuk1ITs4RSfNHaTb6vO3pthfAt6Kdsmytqo3pu3IPDT8GXaxSVbT6kf6ksx8Mp5+/oWpLBvUoqdKLeWvJ77GvanHS1fbDL4+C8W9kuJi3FqSbTW6aOlpVHUpRqSh7NyipcqWEsrOy8NznrqjK3uKlGXWDxnx8yY0irKrZxUnF+zfLst/FZ/BemOiMnRW4cnDPdqamvFTf2NsqguoeDWZqpWKztg3tR0m90yrSp39CVGVWjCvTjLG8JLMXs+jRihHwSR7tKOcldt4YoU/EzUqSlJJ4Se2S9JFejWDqIQ2tMvauL/oZ0bR/oiseL7PiCN1ezjCdelzRcGpbe5jfbK65yeNypKDxjJMaDxNqGnTpUqtWd3YJOM7OtNulOL6rlzhPumt00n2N3iXRaFGFDVtJqTr6TeN+xlLedKS60542Ul8msNHVa8uu7OrkyYb8GWevSfz6uaUHgo47G7Ght0Erd+DPdpTedG6vD+k3Gs6vbabbLNa4qKnTTeFlvCPW+LfoG1bhO1sb/AFrV7KlZ15YrTi96OFnLzhY7dc5aWGcbwZpq0ZU+LdVqzoWtvUTtqcXipc1V0UcraKeMv4dS76SvpI4n47r0f4vdtW1usUbem2oR/wCJpveXm/uPOG0bT27wp3y3z5LVx2mIjvy237pOlxRougN2XCtrlp4nf3EVKrJ9MwXSK7rv6Hsn0X/ThdWtnQ0m6Ua021CE2/xPlNTcejZLcNXlSlq9vUi23Tbn1/tTf5EkzXJtW8RMKuTwzyt82G8xf2783s30u/SzfcS1JUozdClB+7GMnnK9DyyHG/E1pN+w1m5cM/YqS9pH4xllfcc3WvKtVuTk38TBKbkeTbaIivKI6bJ8HhmOsTOWOKZ6zPPd2X/xF1t7V7LR63i56fTy/uMF1xhaXsVDUeFdEqp9ZUaUqE16Sg0vmmcjKW5a5Hk3meq1Tw7T19Wu3w5PXvol0PgrifjCxtqdvf20s81ShWqRnTkl4SwmseGHnfdH0p9J/wBHfC2k8FXurabbW1nd2ttOdKrGEU4yUXhrbZrs0fH/AATKrw/af9LLipOlyNwsaSk4u4qYw3t1jHO/i2kSHE30scV65pH8LvtTnO36Situbbo8dvLoeWi29bRbbbszL6WcmTJSK8ccoiZn1Z7/AB/IcJrGXeVZSm5ycm3JvLeXvv3ZG1HjJsXFV1JOTecmrVawcTzmZfR4azWsRLE9yyRcyypJRi3JpJLdvsRTy5ytxG87QjtbqrFOkpZz77Xh1S+PX7iLMt3V9tcTqL7Lfu7Y27GIwsl+O02bFK8NYhM8F2f13iS0g4uUKc1Vljso7/LODu+PKqqU7a2dTkUm8yf9OcJN+mSE+iuyUql1qEotuKVKDzt4v/2/tl/Hdzm7rRbSUKTgs77tf5m7pK+T4fa/e0/aPyXz+qt5/iVKdqx95/IcJJ5bb7lAD599EAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABKaJWSjKg2855opeaw/wAERZt6VUdO9gl0qe41nGc9M+WcEmK/BeLOMleKswnkblhLMJQfZ5RprdGS2qKFZPs+p9BSdpYmSu8TCM4ptuS5hcqOFUWJY8Utvu/A09GrRpXEoTzyzWFh4w/9M/M6TVrb63ZTpqK58e62+jyvh2/fU45OUJppuMovZrZpmRq8c4c3FHfm0tLeMuHae3J1DWCnVpFtpWVzQhW7y+1tsmuq/D4YL1jJp47ResWjpKjas1tMS6Ktc6HccLUI8l5/HKdZRlUlPmpSo4eEk9008LHTBFpGCmbEXzNJEinwRSJ2VRdFGx9QuXR9t7KXJ/djYwpNPDR1t7UfHWd9pXQhlrbc7PgK4hSnV0vUIyq6XepRrwXWElnlqR8Gm87dVlHKWseaawj0D6OtOd1q1CnOlNxlJZ917ImxV3tEMnxbNWmC3F+S3o/R1q312rQpUXWhB7VIL3ZJrKabfRrczWfAsrW7nc6zCVvY2y56zaw546Qjv1fTK6H2xwLodjpWgW1OhQp88qacpdcvr1OO+nfhnTrzhqteKlGnVWctd3h4ePFEOLXYb5vK4fhKjqvDvEsOgjV2yR0iZjbpHx9r4r4y1Cvq+o+0nBUqFKPs7ehH7NKmukV+b7s5qrDlbR03ENJUrmpFY2bOcr/aZPflMr+hmPLiK8o2alRYZNcEwtnql1WvJwhSoWNxP3njMnTcYpeLbkiGn9oxSystMh3aNqcdJr7U9pN/w9S4J1bT77S6tXWq1ejOyvI1Hy0oJvni1nun55bXTBBJrBjKrY57JIpEb8yeEyY4M0GfEGsq3lUVC0owdW6ryWY0qcVlt+fZLxaIbDnOMYptyeEl1bO24ikuFOFqfC1Fpaje8tfVWusejp0k+myeX5s7rETzlBqstqxGPH61uXwjvPy/vshOM9cjq2oRjaUnb6bbRVGyt859nBeL7tvLb8WQLllFJsxyfgczMzMys4cNcVIpXsrN7GCb3Lpy2MTZxMrNarZEfrFw4QVKMmpTTzt/Tlr8vxNy4qwpU5TnLCSb9X2S8/8AU5+tUlVqyqTeW/u8jP1mbaOCGjpcX7pWFUm2kllvZIodDwJpa1HWY1KkU6NvicsrKb7L9+BQxY7ZbxSvWVrNlrhxzkt0h6Dw5aLSOHaNKo0pQg51N8rPV7o8+4ku51KdaTn79ea5k+6zn8UjuuL7v2GnexT96s8fBbv8l8TzTW6spXEaXPmMI5x2Te/4YN/xa9cOGuCvaIh894NS2XJbUX6zO6PAB84+lAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAdHaXCuKEZ7uTXvZaznv96ePIzdGRGi3KhKVCb2l70fJ9/Tb8CWzk2dLk8zH7+jL1FOC/ulJW0/a0k+r6M5ziSx+r3Cr04KNKps0n0l/n+vQmLOr7Kss55ZbP9TcvKEbm2qUJr3Zxx06efzwT6nD5+LaOsdFXBl/T5efSXJ6PdRoVpU6rSp1O7/pl2f5fHyJvDzhrDXY5u8t6tpczoVo4lH711TJfR7tV6XsZY9rBbbfaXj6r/PxM7R5uGeCzS1OLjjjhJQexs2coqtFy6J7mqtjJF4Zqxylk3rvvD634d4x+iWn9DdPT72nR+uKhyVaLjmpKectr9TwKtq/Bca0pR4Wup77KWoPD+UUcXGrPpzMujLxO4tMMivhdKTvvPTtO3126z73b0uLNJtt9K4R0y3qLpUuJSuGvNKTx9zJrhzifijWNUoWdtcy55y2pUYqnHbd5SSWEll52web2kZ1KkaVKLlOTSSSy8vokd/cyXBmmvTadSL1u6ppXs4vLtoPDVNNbZaeW0/Imx2nffoo6/SYaxFK13tPTfn8+fZ9R8GfTNotnZ/w2/rKr9XXJGrT6Swt0vLOVnyIHjz6Urfiaq9N010M8rdKjX+xVlunHOVu1su258u0tTnBe7Nr0ZhuNSqt86nJNbpp7o4riw0yeZWvNVvpdbm08aS+WZpHbl29vuddrseG9UvK1O4r1uH72MnGVOtF1KTlndZS5o7+Kfqc7q3Buu21u7ylbRvbPqri0mq0GvHMc4Xrg3aValxpQhYXdaFDXKcVG2uJvEbtJbU5vopeEn16PscxQvtZ4f1KpCjXubC6oycJqM3GSw908Hd7RPZa0eHLjjgpbaY7Tz+k9dvq3eG+Ftc4i1mjpWmafWrXNaXLFcjSXm29kvUxcdcK6xwdr1TRdct/YXcIqfKpKSaecNNeOGS2n/Sbxtp91SubfXq8atKSlFyjGXTp1W5l1v6QJcTahO/4v0W01e7nFQlcxlKjUaSwt4+7svIgnbi9zRpbU1iZvWJ90T/vZwy8EZ4W1WcW1Bv0R09vY8FaknUpatfaRVb/AN1cUFWgv/HHDx/4T6V+hL6EuH9S4WhfarOF1OqswlTb5WuzWUn08ULTWlZtadoc31trZIxYqTN57Ty+/R81cHWNLSLarxVqVJVIWklGyoTW1evjb1Udm/gu5zGpXlxfXta9uqsqterNznOTy228tnvH+0nwTDhm+tLS2q/9To0nGjSS2jvlv4t98nz/AHXu1GsndtorHDziXPh+Sc17ZLxtffaY9m3Zhm8sxtoSkY5y2wiOZ7NmtVJvLwjFJxj9qcYru5PC2/fqJyUU230/IhtRu/bS9nTeaafX+79F+/SnqM8Y45dV3Bgm8+5ivrl3FV8rl7NfZT6+rNcAx5mbTvLTiNuULqUJ1akadOEpzk8RjFZbZ67wvpcNH0qFvmLqv3qs1/U/8vyOX+j3QeZx1e6hsv8AcRfycmvwOl4mvnaWLpwf8yqmk090u7/I+i8J0vlUnU5I+D5rxjVTnyRpcc/H89zm+KtQVxdVaik3Tppxhjy9Tias3UqyqNJOUm2ksLcltbr8tCNFNc1R8zw+kU/za77rHmQ5k6/POXLPubeiwRhxRWAAFJbAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAXU5ypzjOEnGUWmmuzOhtKyuKMaiSWeqXZ+H77HOG5plz9XrOL5eSezb7Ps8/v44LGmy+XfeeiHNj8yvvTq8yRsayqQ5H9qK280Rqy8ZTXlgvpzcJqUXhrc26X2ndj5KcUbd2bX9M+u0FOmv50F7uZYyvA5H+ZRrdJQqQl0aw00egW1VV6alHbGzXgyI4j0dV4O5t4pVYr3kl9pfqVNbpOKPNx/NNodXwz5WT5NXTbuNzRSz/MgveXV9ve9P34Z3Iy7M5SMqlCrlZhOL8N0TmnX9O4hGE5ctfOMNbS9P0/HtFptXvtW6xqNN+6qSUty+En6GBPfdnr3+zfwVwXxpq1/Q4t1dWMLekp0aft4Uvatt53fZeXiX5ttEzLMy24I3QPC8aXCmiw4pvIwnqNwnHSqElnGNpVpLwXRJ9Xv2OXrXlW5r1LivUlOtUk5SlJ5bbeW2TP0p3FKfHGo0LW8hdWdtWdC1nTxyKlDaKjhtYSWNuvXucwp7bslm3LZQ0+Diic1+drfaO0fndIKtt1KSq5XU0VU8yvP5ocSbyY3bLquMsxeGvA6uzvbTjKhT03V60bfWqcVC0vpvEbjGypVOyfZSfozinPxZY54fMmc8e3weZdNGSN45THSUjqthe6Ve1bC/oTo16UnGcZLDWO/n5Y2ZpSnhPHU9v/ANn3hyy+lu8loHFN5yQ0ykp0bmEkrmcc4VLLTTgt3lrKeEuprfSvwDwX9GvF9e2uL+tq8YwjUtbLplv/AL2SxhZXRbsb1tfhj4of1N8dP6ld5325d5/x83m/C/Dt3qNGepXVVWOlUGvbXlVPC8orrKXkvLOD2f6Pfp1q8J270yxozubClFQpSr1P5ksLGXs0vJLosLc8S4n4kv8AXKsFcTjTtaOVQtqS5aVFeEYrp0W/V92QqrYezZ1M124dt496O+iyam0ZMk8Mx02nbb5/kPTfpZ+ke/4z1Kd3eVFhLEILGIr5I8vuJuU856sVKrcst5Necs+hza2/KOy9pNJXBXavf6kpeBhrTUKcpzajGO7b6IxXd3St4ZnJOWV/LT3ff4Lf9M4IS5uKlxPmm9l0XgUc+rinKvOWvg002526Ml/du4m1DMafZPq/P/I1QDJtabTvLRiIiNoDoODeH5axdurXjONnT+00vtvwTLOFeHq+sXCnNSp2kH79TH2vKPj+R6haUKFnaQoUYqnSpxxFZ2SNbwzw6c9vMvHox92N4r4nGmr5eOfTn7FWVGztXOeIUqce3ZLscBq9+7qvUua0lGHXxwl0X77kpxPqyu6n1ei/5FN9f7n4+ngcbq9yqklRp5xH7e3V+Hj/AJ+hoeJ6yK14afn/AOKfhGhtH9TJ1lp3NWVevKrJJcz2S7LsvkYwD5h9MAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAJfSrxSpq3qP3o7RfivD4fvoSK6ZOYjKUZKUW4yTymnhpk3p14riHJN/wA1deiz5/v/AE0NJqP2W+SlqcO/pVSVrXlQqqcenRrxRM0ZwrU1ODyn27ryIAz2lzO3nzR3XdPozWx5OHlPRlZsXHG8dV2t6FG8/nW3JTqpbrGFL1/U5K4o1berKlWg4Tj1TPSLWrTuKanTefFPqmYNU0u21CkoV44kuk49V1+7f06FXVeHxk9PEk0viM454MrjbDUuTELjmkuiku3qu/76kvTrYip0p5jLaMovZ46/kROr6Jd6fJyx7Wj2qRX4rsaFtcVrefPSnyvwayn6p7MoU1GTDPDeGpbDjzRxVl0/tG3lvI5/PBFW2qU5csa0PZvvJNtdsbdV67m9RqRq0vaU5KccZfK8437+G/iXsepx36SpX01qdmx7TzKqo84MCe3U3dEpW1xqtpQvKvsrapWhGrPP2INpN/BZZYiUFoisTM9mNOUmoxTbbxt1On0/hG4harUOILmGj2PVKsn7asv+Cn1fhl4XmegfSb/0C4D/AIauAnaajf3FvzXFzVn7Z28ljDinsm3l9O3gzyTVNTvdTu5XV9dVbmtPrOpJtv5/gdcoZ1MubUxvjjhr7e/y7fnR1FPi6noeKHB1KppsIPLvJS/6zV/xSW0V35UsdMtnP6zq99q11O71G7q3NxP7U6k3KT+LIrnx03Mcptnk2lPi0eOk8URz9vdlnPL6lkpGJvxZjuq8aFLnk0s9n19V4422RFky1pG8ruPDNp2hmnOMIuU5qMV1beEv34dSJu9VzmNvHfdOcln5L9fuNG6uKtxU5qk+bG0V2SMJl5tXfJyjlDTxaatOc85Vk3JtybbfVsoDd0rS77U6vs7O3lUWUpTxiMfVlWIm07QnmYrG8tI6nhXhStfyjc30XStuqi9nPy8kdFw7whaWGK9443NwsNbNRg/L49zpJ8tOLlJqKiuucJI3tF4PM7Xz9PZ/t87r/G4jfHp+c+3/AEWtClbUI0qUI06cVhRSwl+hzPE2te0crO0n7i2qTX9XkvL8SziDXpVlK1tJONPdTn0cvJeX4nLXt1ChT7Ob2S8/PfOC/rNZTHThpyiPzaFTQeH3tfzMvOZ/N5W6rcxp0OSMs1J9vBbp5+JCF9apOtVlUqS5pSeWyw+Wy5Zy24pfV46RSu0AAInYAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAF0JyhNThJxlF5TXYtAE5p97C4iozajVit039rzX6fl03Tl03FpptNbpolbHU3KXJdS954SqePm/Pz+ZoafV7ejf6qebTb86pi3r1LeanTeH38H5Mm7G+pXKSfuVPDOzOehJSinF5WPH4lU2nlPD7GrjyzXp0ZWbTxfryl1bT646fcQuocNWV1zSpZoVW/tR6fFem/Z9clbLVqlJKFxH2sF0f9S+PcmratQuYc1CpGWFlruvVE18eHUxtaFKL59LO9Z/04DUeH9Rs1zezVaH91Ld/Lr+RFLMXlZTX3Hq0o74aNa60jT7181zawnJ4zJbNpea3RnZfCJ645+UtHD41Xplj6PO4ahdKSlKp7TfL51ly9X1+82I6pLdyo7t5SjPCS8PH7zpbngu3nJu3u6lJdfeipJfg/0Iy44N1Wnl05UKqSysSab8sY6/cVJ02rxdIn5c16ut0mX90fPk1qWrW8pZre3prwjFT+/MTN/FdNx/vLz/AMiP/Oalxw3rdB/zLCb/AMEoz/Bswy0TV4pOWm3KT6Zps487U167/RJFNNbpMfVt1tWt0/5UatRf8UVD82a89XlzJxoLbrGcsp7eWH95ZHRNXl9nTbl7Z2pszW/DeuXDxT02t/48R/Fo849Rf2/R7tp6dZj6tOV/dt5VZwednBYa+PU1TprbgjXK2OeFGju0+eecfLOfhkm7H6O6eE7zUZyb6xpQ5Uvi85+SOqaDVZZ5Un5/9RZPEtJhjnePlz/s8+N/TdH1LUGvqtpUlB/1tYj1x1PULHhTRbPldOyjUkse9V99t+Pl8MdEyWdN9DSw+A5J/wDW23wZWf8A+RY4j+lXf4uK0PgehTftNUq+3kv+zg2o/F9X9x11ra0bajGjb040qcVhRisY8jNUdOjTdSrOMIR6yk8JfEgNW4mo0uanYx9pNbe0kvdXourNfHp9Po43iOf3Y19Rq/ELbb7x9IhL311b2VF1bioorsu8n4JdzkNb1qtft0oZp0E9orv5t/kR13dV7mq6txUlOTfdkXeajCmuWjic9mn1S/e372KGr8Q5e77y19D4XFJ3nnPtZb25p0aTbknUe0Y9c+fp5/LviFrVZ1qjnN5f4eSLZNyeZNt+ZQ+dy5rZZ3l9FjxxjjaAAESQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAbNne1rf3U+ennPI309PAmLa7oXCbhJKSTbjJ4fn6/5M54E+LUXxdOiLJhrk6un6l1OpOnJShKUZLo08EHQ1K4p/7x+2WMJye6+Pfw3zsSdveW1fljCrGM5dY1MRx8Xt+ZoY9XS3XlKjfS2r05wnLTW68MKvFVYrbwfzJi01OxrLer7OT/pmvz6HI82S5dDQpqLV77s7Losdu20u+p01KKlFqSe6aeUzNCk8HAUbivQfNRqzg14NkjR4i1WnHl9tGaX90E380slumsp+6Nmbl8Myfstu7KNF+BlhRbXQ5WhxbcwSVW0pVPRuLNynxnSSXNpr+Fb/ACLFdVh7z9lO/h+qjpG/zdHC3fh9xmhbbdDmv+nFFdNMb9a3/wDJbU47q8uKOnUoP/iqOX3LBLGs09Y6/ZBPhuttPq7fN1sbbyLvq+OxwVXjPWZ5dOpSpJ/20k8fPJE32rajetu5vK1TPbmwvktiO3iOKI9GJlNj8E1Fp9O0RD0G/wBV0yy5o3F3TUl1hF8z+Syc7qPF0XFxsbZp52nUefuX6nJZbeHnOdu+5gr3dvRS56i3/pW7+X64KGbxO23WIhrabwXFXaZ3tP2SV/qN3ezcrmtKXgs4S9F0RHXFxRoxzUkltlLu/T9ojrjU6k48tKKhv9p9X+RoTlKcnKcnKT3bby2YubXzafR+rfw6OKxz+jau72pW92PuxXdbN+pqAGfa02neV2IiI2gABy9AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGWjc16KxTqyUc55eqz44exuUtVqqUfaUoyS68r5W/wAUvgiOB1W9q+rLm1K26wmlqtBxgnGak/t5WyfluzMru2lNRhXpzyk855UvL3sepz4LEazLHXminTY57Ohq3NCnJRlXpt4z7s1Jfdkp9bt9/wCdD5o58Hca6/shxOko6GFzb1KnJGtST8ZSSXzexSdzbwqKNSvTSfWSkpL/ANOWc+Dydbk9z2NLRN/xO1TkpSqNY91whlP54/fY1qmqybXJRXTD5nlP5YwRoIranLbukjDSOzYq3lzVUlKtJRn9qMfdT9UuprgEMzM9UoADwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB/9k=" alt="OnePilot" width="36" height="36" style="object-fit:contain;width:36px;height:36px;filter:drop-shadow(0 0 3px rgba(0,200,245,0.5));">
    OnePilot
  </div>
  <span class="badge">Lecture seule</span>
  <span class="meta">Expire le {expires}</span>
</div>
<div class="main">
  <div class="dashboard-title">{title}</div>
  <div class="widgets-grid" id="grid"></div>
</div>
<footer>Généré par OnePilot — Agent Conversationnel ERP · Ce lien expire le {expires}</footer>
<script>
const spec = {spec_escaped};
const COLORS = ['#006eb8','#2d72ff','#00a050','#d4a000','#c03030','#7c3aed'];
const grid = document.getElementById('grid');

function fmtVal(v) {{
  if(v===null||v===undefined) return '—';
  const n=parseFloat(v); if(isNaN(n)) return String(v);
  if(Math.abs(n)>=1e6) return (n/1e6).toFixed(2)+'M';
  if(Math.abs(n)>=1e3) return (n/1e3).toFixed(1)+'K';
  return n.toLocaleString('fr-FR',{{maximumFractionDigits:2}});
}}

(spec.widgets||[]).forEach((w,i) => {{
  const div = document.createElement('div'); div.className='widget';
  const ct = w.chart_type||''; const d = w.data||{{}};
  let inner = `<div class="widget-title">${{w.title||''}}</div>`;

  if(ct==='kpi_card'||ct==='kpi_delta') {{
    const dv=parseFloat(d.delta); const col=dv>=0?'#00d9a0':'#f07070'; const arr=dv>=0?'▲':'▼';
    inner += `<div class="kpi-card"><div class="kpi-value">${{fmtVal(d.value)}}</div>
      <div class="kpi-label">${{d.label||w.title||''}}</div>
      ${{!isNaN(dv)?`<div class="kpi-delta" style="color:${{col}}">${{arr}} ${{Math.abs(dv).toFixed(1)}}%</div>`:''}}
    </div>`;
  }} else if(ct==='table'||ct==='pivot') {{
    const headers=d.headers||[]; const rows=(d.rows||[]).slice(0,50);
    inner += `<div class="table-wrap"><table><thead><tr>${{headers.map(h=>`<th>${{h}}</th>`).join('')}}</tr></thead>
      <tbody>${{rows.map(r=>`<tr>${{r.map(c=>`<td>${{c??''}}</td>`).join('')}}</tr>`).join('')}}</tbody></table></div>`;
  }} else {{
    inner += `<div class="chart-container"><canvas id="c${{i}}"></canvas></div>`;
  }}
  div.innerHTML = inner; grid.appendChild(div);

  if(!['kpi_card','kpi_delta','table','pivot'].includes(ct)) {{
    const labels=d.labels||[]; const ds=d.datasets||[];
    const ctx=document.getElementById('c'+i)?.getContext('2d'); if(!ctx) return;
    const isBar=ct==='bar'||ct==='bar_horizontal'||ct==='column';
    const isPie=ct==='pie'||ct==='doughnut';
    new Chart(ctx, {{
      type: isPie?ct:'bar'===ct?'bar':ct==='bar_horizontal'?'bar':ct==='line'||ct==='area'?'line':'bar',
      data: {{
        labels,
        datasets: ds.map((s,j)=>({{
          label:s.label||'',
          data:s.data||[],
          backgroundColor:isPie?COLORS:COLORS[j%COLORS.length]+'33',
          borderColor:isPie?COLORS:COLORS[j%COLORS.length],
          borderWidth:isPie?0:2,
          fill:ct==='area',
          tension:0.4,
        }}))
      }},
      options:{{
        indexAxis: ct==='bar_horizontal'?'y':'x',
        responsive:true, maintainAspectRatio:false,
        plugins:{{legend:{{labels:{{color:'#7ea4be',font:{{size:10}}}}}}}},
        scales: isPie?{{}}:{{
          x:{{ticks:{{color:'#4a6a82',font:{{size:9}}}},grid:{{color:'#e8f0f8'}}}},
          y:{{ticks:{{color:'#4a6a82',font:{{size:9}}}},grid:{{color:'#e8f0f8'}}}}
        }}
      }}
    }});
  }}
}});
</script>
</body>
</html>"""
        return HTMLResponse(html)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/share/{token}")
async def delete_share_link(token: str):
    """Supprime un lien de partage."""
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.execute("DELETE FROM shared_dashboards WHERE token=$1", token)
        return {"status": "ok"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── FAVORIS ─────────────────────────────────────────────────────────────────

class FavoriteRequest(BaseModel):
    dashboard_id: str
    title:        str
    spec_json:    str
    source_id:    Optional[str] = None
    conv_id:      Optional[str] = None
    user_id:      str = "admin"

@app.post("/api/favorites")
async def add_favorite(req: FavoriteRequest):
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.execute("""
                INSERT INTO user_favorites (user_id, dashboard_id, title, spec_json, source_id, conv_id)
                VALUES ($1,$2,$3,$4,$5,$6)
                ON CONFLICT (user_id, dashboard_id)
                DO UPDATE SET title=EXCLUDED.title, spec_json=EXCLUDED.spec_json,
                              source_id=EXCLUDED.source_id, conv_id=EXCLUDED.conv_id, created_at=NOW()
            """, req.user_id, req.dashboard_id, req.title, req.spec_json, req.source_id, req.conv_id)
        return {"status": "ok", "dashboard_id": req.dashboard_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/favorites/{dashboard_id}")
async def remove_favorite(dashboard_id: str, user_id: str = "admin"):
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            await conn.execute("DELETE FROM user_favorites WHERE user_id=$1 AND dashboard_id=$2", user_id, dashboard_id)
        return {"status": "ok"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/favorites")
async def get_favorites(user_id: str = "admin"):
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT dashboard_id, title, spec_json, source_id, conv_id, created_at
                FROM user_favorites WHERE user_id=$1 ORDER BY created_at DESC
            """, user_id)
        return {"favorites": [dict(r) for r in rows]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/dashboard/export/pptx-multi", tags=["Dashboard"])
async def export_dashboard_pptx_multi(req: PptxMultiRequest):
    """Génère un PPT multi-slides : 1 slide par dashboard."""
    import json, sys, os
    from fastapi.responses import Response
    try:
        specs = json.loads(req.specs_json)
        sys.path.insert(0, os.path.dirname(__file__))
        from pptx_generator import generate_pptx_multi
        pptx_bytes = generate_pptx_multi(specs)
        title_safe = req.title.replace(" ","_")[:40]
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{title_safe}.pptx"'}
        )
    except Exception as e:
        logger.error(f"[PPTX-MULTI] Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


class PptxImageRequest(BaseModel):
    image_base64: str
    title:        str = "Dashboard"
    subtitle:     Optional[str] = None

@app.post("/api/dashboard/export/pptx-image", tags=["Dashboard"])
async def export_dashboard_pptx_image(req: PptxImageRequest):
    """Génère un PPT avec le dashboard capturé comme image pleine page."""
    import base64, io
    from fastapi.responses import Response
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        img_bytes = base64.b64decode(req.image_base64)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # ── Slide de couverture ──
        s0 = prs.slides.add_slide(prs.slide_layouts[6])
        s0.background.fill.solid()
        s0.background.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
        # Barre header
        hdr = s0.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = RGBColor(0x00,0x6E,0xB8)
        hdr.line.fill.background()
        # Titre
        tb = s0.shapes.add_textbox(Inches(0.4), Inches(0.07), Inches(9), Inches(0.38))
        tf = tb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = req.title; r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.name = "Calibri"
        # Sous-titre/date
        if req.subtitle:
            tb2 = s0.shapes.add_textbox(Inches(9.5), Inches(0.12), Inches(3.5), Inches(0.28))
            tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
            r2 = p2.add_run(); r2.text = req.subtitle
            r2.font.size = Pt(9); r2.font.color.rgb = RGBColor(0xCC,0xEA,0xFF); r2.font.name = "Calibri"
        # Image dashboard pleine page
        img_stream = io.BytesIO(img_bytes)
        s0.shapes.add_picture(img_stream, Inches(0), Inches(0.5), Inches(13.33), Inches(6.8))
        # Footer
        ft = s0.shapes.add_textbox(Inches(0), Inches(7.18), Inches(13.33), Inches(0.28))
        tft = ft.text_frame; pft = tft.paragraphs[0]; pft.alignment = PP_ALIGN.CENTER
        rft = pft.add_run(); rft.text = "OnePilot — Agent Conversationnel ERP"
        rft.font.size = Pt(8); rft.font.color.rgb = RGBColor(0x9A,0xB4,0xC8)
        rft.font.name = "Calibri"; rft.font.italic = True

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        title_safe = req.title.replace(" ","_").replace("/","_")[:40]
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{title_safe}.pptx"'}
        )
    except Exception as e:
        logger.error(f"[PPTX-IMG] Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


class PptxExportRequest(BaseModel):
    spec_json: str
    title:     Optional[str] = None

@app.post("/api/dashboard/export/pptx", tags=["Dashboard"])
async def export_dashboard_pptx(req: PptxExportRequest):
    """Génère un fichier PowerPoint depuis un spec de dashboard (python-pptx)."""
    import json, sys, os
    from fastapi.responses import Response
    try:
        spec = json.loads(req.spec_json)
        if req.title:
            spec["title"] = req.title
        # Import du générateur python-pptx
        sys.path.insert(0, os.path.dirname(__file__))
        from pptx_generator import generate_pptx
        pptx_bytes = generate_pptx(spec)
        title_safe = (spec.get("title","dashboard")).replace(" ","_").replace("/","_")[:40]
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{title_safe}.pptx"'}
        )
    except HTTPException:
        raise
    except ImportError as e:
        logger.error(f"[PPTX] Import error (python-pptx manquant?): {e}")
        raise HTTPException(status_code=500, detail=f"Module manquant: {e}")
    except Exception as e:
        import traceback
        logger.error(f"[PPTX] Export error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/dashboard/generate", tags=["Dashboard"])
async def generate_dashboard(req: DashboardRequest):
    """
    Génère un dashboard interactif complet depuis une question NL.
    §2.4.3 — Pipeline : NLU → SQL → Data → Spec → Chart.js widgets
    """
    pool  = await get_pg_pool()
    redis = await get_redis()

    # ── Cache Redis ───────────────────────────────────────────
    cache_key = f"onepilot:dashboard:{req.source_id}:{abs(hash(req.question))}"
    try:
        cached = await redis.get(cache_key)
        if cached:
            data = json_module.loads(cached)
            data["cached"] = True
            return data
    except Exception:
        pass

    source = await get_source(UUID(req.source_id))
    if not source:
        raise HTTPException(404, "Source introuvable")

    # ── Schéma ────────────────────────────────────────────────
    schema: dict = {}
    try:
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT se.name AS t, ef.name AS f
                FROM source_entities se
                JOIN entity_fields ef ON ef.entity_id = se.id
                WHERE se.source_id = $1 AND se.is_visible = TRUE
                ORDER BY se.name, ef.position
                LIMIT 5000
            """, UUID(req.source_id))
        for r in rows:
            schema.setdefault(r["t"], []).append(r["f"])
    except Exception as e:
        logger.warning(f"[Dashboard] Schema error: {e}")

    # ── NLU ───────────────────────────────────────────────────
    nlu   = get_nlu_pipeline()
    known = list(schema.keys())
    slots = nlu.process(req.question, None, known)

    # ── Dashboard Generator ───────────────────────────────────
    try:
        from .dashboard_engine import get_dashboard_generator
        generator = get_dashboard_generator()
        spec = await generator.generate(
            question           = req.question,
            slots              = slots,
            schema             = schema,
            source_id          = req.source_id,
            pg_pool            = pool,
            redis              = redis,
            connector_factory  = ConnectorFactory,
        )
        result = spec.to_dict()
        result["cached"] = False

        # Cache 5 minutes
        try:
            await redis.setex(cache_key, 300, json_module.dumps(result, default=str))
        except Exception:
            pass

        return result

    except Exception as e:
        logger.error(f"[Dashboard] Generation error: {e}", exc_info=True)
        raise HTTPException(500, f"Erreur génération dashboard: {e}")


@app.post("/dashboard/data", tags=["Dashboard"])
async def execute_dashboard_sql(req: DashboardDataRequest):
    """
    Exécute un SQL pour un widget dashboard et retourne les données.
    Permet le refresh individuel d'un widget.
    """
    pool = await get_pg_pool()
    source = await get_source(UUID(req.source_id))
    if not source:
        raise HTTPException(404, "Source introuvable")
    try:
        src_dict = source.model_dump()
        # ── Fetch password depuis connection_secrets ──────────────
        try:
            async with pool.acquire() as _c:
                _sec = await _c.fetchrow(
                    "SELECT secret_value FROM connection_secrets "
                    "WHERE source_id=$1 AND secret_key='password'",
                    UUID(req.source_id)
                )
                if _sec:
                    src_dict["password"] = _sec["secret_value"]
        except Exception:
            pass
        # ─────────────────────────────────────────────────────────
        connector = ConnectorFactory.create(src_dict)
        rows = connector.execute_query(req.sql)
        if not rows:
            return {"rows": [], "columns": [], "count": 0}
        cols = list(rows[0].keys()) if rows and isinstance(rows[0], dict) else []
        return {
            "rows":    rows[:req.limit],
            "columns": cols,
            "count":   len(rows),
        }
    except Exception as e:
        logger.error(f"[Dashboard/data] {e}")
        raise HTTPException(500, str(e))


@app.get("/dashboard/templates", tags=["Dashboard"])
async def list_dashboard_templates():
    """Retourne les templates de dashboards pré-configurés."""
    return {
        "templates": [
            {
                "id":          "sales_overview",
                "name":        "Vue d'ensemble Ventes",
                "description": "CA, commandes, top clients, évolution mensuelle",
                "question":    "Dashboard des ventes avec évolution mensuelle et top clients",
                "icon":        "bar",
            },
            {
                "id":          "top_entities",
                "name":        "Top Entités",
                "description": "Classement des entités par valeur",
                "question":    "Top 10 par montant total",
                "icon":        "bar_horizontal",
            },
            {
                "id":          "trend_analysis",
                "name":        "Analyse de tendance",
                "description": "Évolution temporelle d'une métrique",
                "question":    "Évolution mensuelle des montants",
                "icon":        "line",
            },
            {
                "id":          "distribution",
                "name":        "Répartition",
                "description": "Distribution des données par catégorie",
                "question":    "Répartition par catégorie",
                "icon":        "pie",
            },
        ]
    }


# ── Sauvegarde et restauration des dashboards par conversation ────────────
class DashboardSaveRequest(BaseModel):
    conv_id:      str
    dashboard_id: str
    spec:         dict

@app.post("/dashboard/save", tags=["Dashboard"])
async def save_dashboard(req: DashboardSaveRequest):
    """
    Sauvegarde le spec JSON d'un dashboard pour une conversation.
    Permet la restauration lors du rechargement de la conv.
    """
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            # Essai 1 : UPSERT avec ON CONFLICT
            try:
                await conn.execute("""
                    INSERT INTO conversation_dashboards
                        (conv_id, dashboard_id, spec_json, updated_at)
                    VALUES ($1, $2, $3, NOW())
                    ON CONFLICT (conv_id, dashboard_id)
                    DO UPDATE SET
                        spec_json  = EXCLUDED.spec_json,
                        updated_at = NOW()
                """,
                    req.conv_id,
                    req.dashboard_id,
                    json_module.dumps(req.spec, default=str),
                )
            except Exception as e_upsert:
                # Fallback : DELETE + INSERT si la contrainte unique est absente
                logger.warning(f"[Dashboard/save] UPSERT failed ({e_upsert}), trying DELETE+INSERT")
                await conn.execute(
                    "DELETE FROM conversation_dashboards WHERE conv_id=$1 AND dashboard_id=$2",
                    req.conv_id, req.dashboard_id
                )
                await conn.execute("""
                    INSERT INTO conversation_dashboards
                        (conv_id, dashboard_id, spec_json, updated_at)
                    VALUES ($1, $2, $3, NOW())
                """,
                    req.conv_id,
                    req.dashboard_id,
                    json_module.dumps(req.spec, default=str),
                )
        return {"status": "saved", "conv_id": req.conv_id, "dashboard_id": req.dashboard_id}
    except Exception as e:
        logger.error(f"[Dashboard/save] {e}")
        raise HTTPException(500, f"Erreur sauvegarde dashboard : {e}")


@app.get("/dashboard/conv/{conv_id}", tags=["Dashboard"])
async def get_dashboards_for_conv(conv_id: str):
    """
    Retourne tous les dashboards sauvegardés pour une conversation.
    Appelé au chargement d'une conv pour restaurer les visualisations.
    """
    try:
        pool = await get_pg_pool()
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT dashboard_id, spec_json, updated_at
                FROM conversation_dashboards
                WHERE conv_id = $1
                ORDER BY updated_at ASC
            """, conv_id)
        if not rows:
            raise HTTPException(404, "Aucun dashboard pour cette conversation")
        dashboards = []
        for r in rows:
            try:
                spec = json_module.loads(r["spec_json"])
                dashboards.append(spec)
            except Exception:
                pass
        if not dashboards:
            raise HTTPException(404, "Aucun dashboard valide pour cette conversation")
        # Retourner le dernier dashboard (comportement legacy) + liste complète
        return {
            "dashboards": dashboards,
            "latest":     dashboards[-1],
            "count":      len(dashboards),
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[Dashboard/conv] {e}")
        raise HTTPException(500, f"Erreur récupération dashboards : {e}")


# ══════════════════════════════════════════════════════════════
# §2.3.3.C — PLAN EXECUTION (parallélisation asyncio.gather)
# ══════════════════════════════════════════════════════════════

class PlanExecuteRequest(BaseModel):
    plan_id:    Optional[str] = None
    question:   str
    source_ids: List[str]

@app.post("/nlu/execute-plan", tags=["NLU"])
async def execute_query_plan(req: PlanExecuteRequest):
    """
    Exécute un plan cross-source avec parallélisation réelle.
    Les steps parallel=True sont exécutés via asyncio.gather().
    §2.3.3C — Universal Query Planner execution.
    """
    import time as _time
    pool  = await get_pg_pool()
    redis = await get_redis()

    sources, schemas = [], {}
    for sid in req.source_ids[:5]:
        try:
            src = await get_source(UUID(sid))
            if not src: continue
            sources.append(src.model_dump())
            async with pool.acquire() as conn:
                rows = await conn.fetch("""
                    SELECT se.name AS t, ef.name AS f
                    FROM source_entities se
                    JOIN entity_fields ef ON ef.entity_id = se.id
                    WHERE se.source_id = $1 AND se.is_visible = TRUE
                    LIMIT 500
                """, UUID(sid))
            schema = {}
            for r in rows:
                schema.setdefault(r["t"], []).append(r["f"])
            schemas[sid] = schema
        except Exception as e:
            logger.warning(f"[ExecutePlan] Source {sid}: {e}")

    if not sources:
        raise HTTPException(400, "Aucune source valide")

    nlu   = get_nlu_pipeline()
    known = [t for s in schemas.values() for t in s.keys()]
    slots = nlu.process(req.question, None, known)

    planner = UniversalQueryPlanner()
    plan    = planner.plan(slots, sources, schemas)

    t0 = _time.time()

    async def _execute_step(step):
        cache_key = f"onepilot:step:{abs(hash(step.query))}"
        try:
            cached = await redis.get(cache_key)
            if cached:
                return step.step_id, {"data": json_module.loads(cached), "cached": True}
        except Exception:
            pass

        if step.action == "sql":
            src = next((s for s in sources if str(s.get("id","")) == step.source_id), None)
            if not src:
                return step.step_id, {"error": f"Source {step.source_id} introuvable"}
            try:
                connector = ConnectorFactory.create(src)
                loop = asyncio.get_event_loop()
                rows = await loop.run_in_executor(None, connector.execute_query, step.query)
                data = rows[:500] if rows else []
                try:
                    await redis.setex(cache_key, step.cache_ttl, json_module.dumps(data, default=str))
                except Exception:
                    pass
                return step.step_id, {"data": data, "sql": step.query, "cached": False}
            except Exception as e:
                return step.step_id, {"error": str(e), "sql": step.query}
        return step.step_id, {"skipped": True}

    executed = set()
    all_results = {}
    remaining = list(plan.steps)
    for _ in range(10):
        if not remaining:
            break
        ready = [s for s in remaining if all(d in executed for d in s.depends_on)]
        if not ready:
            break
        step_results = await asyncio.gather(*[_execute_step(s) for s in ready], return_exceptions=True)
        for item in step_results:
            if isinstance(item, Exception):
                continue
            sid, result = item
            all_results[sid] = result
            executed.add(sid)
        remaining = [s for s in remaining if s.step_id not in executed]

    ms = int((_time.time() - t0) * 1000)

    merged_data = []
    if plan.merge_strategy == "join":
        merge_key = planner._find_merge_key(slots, schemas)
        data_sets = [r.get("data", []) for r in all_results.values() if "data" in r]
        if len(data_sets) >= 2:
            base = {row.get(merge_key): row for row in data_sets[0] if isinstance(row, dict)}
            for row in data_sets[1]:
                if isinstance(row, dict):
                    key = row.get(merge_key)
                    merged_data.append({**base[key], **row} if key and key in base else row)
        elif data_sets:
            merged_data = data_sets[0]
    else:
        for r in all_results.values():
            if "data" in r:
                merged_data = r.get("data", [])
                break

    return {
        "question":       req.question,
        "plan_steps":     len(plan.steps),
        "executed_steps": len(executed),
        "merge_strategy": plan.merge_strategy,
        "duration_ms":    ms,
        "results":        all_results,
        "merged_data":    merged_data[:100],
        "total_rows":     len(merged_data),
    }


class BatchAPIRequest(BaseModel):
    source_id:   str
    question:    str
    total_pages: int = 3

@app.post("/api/batch-query", tags=["Query"])
async def batch_api_query(req: BatchAPIRequest):
    """
    Exécute plusieurs pages d'une requête API en parallèle.
    §2.3.3B — Batch requests avec asyncio.gather.
    """
    import httpx as _httpx
    source = await get_source(UUID(req.source_id))
    if not source:
        raise HTTPException(404, "Source introuvable")
    nlu   = get_nlu_pipeline()
    slots = nlu.process(req.question, None, [])
    from .query_engine import APIQueryBuilder
    bld   = APIQueryBuilder()
    pages = bld.build_batch(slots, [], source.base_url or "", req.total_pages)

    async def _fetch(p):
        try:
            async with _httpx.AsyncClient(timeout=15) as client:
                resp = await client.get(p["url"])
                return {"page": p["page"], "data": resp.json(), "status": resp.status_code}
        except Exception as e:
            return {"page": p["page"], "error": str(e)}

    results = await asyncio.gather(*[_fetch(p) for p in pages])
    all_data = []
    for r in results:
        d = r.get("data")
        if isinstance(d, list): all_data.extend(d)
        elif isinstance(d, dict) and "value" in d: all_data.extend(d["value"])

    return {"pages_fetched": len(results), "total_rows": len(all_data), "data": all_data[:500]}


class ExplainRequest(BaseModel):
    sql:       str
    source_id: str
    dialect:   str = "mssql"

@app.post("/sql/explain", tags=["Query"])
async def explain_sql(req: ExplainRequest):
    """
    Analyse et optimise un SQL — explain plan + WITH(NOLOCK) + index hints.
    §2.3.3A — SQL Optimizer.
    """
    from .query_engine import get_sql_optimizer
    pool = await get_pg_pool()
    schema = {}
    try:
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT se.name AS t, ef.name AS f
                FROM source_entities se JOIN entity_fields ef ON ef.entity_id = se.id
                WHERE se.source_id = $1 AND se.is_visible = TRUE LIMIT 500
            """, UUID(req.source_id))
        for r in rows:
            schema.setdefault(r["t"], []).append(r["f"])
    except Exception as e:
        logger.warning(f"[Explain] {e}")

    result = get_sql_optimizer().optimize(req.sql, schema, req.dialect)
    return {
        "original_sql":  req.sql,
        "optimized_sql": result["sql"],
        "explain_plan":  result["explain_plan"],
        "index_hints":   result["index_hints"],
        "warnings":      result["warnings"],
    }


# ══════════════════════════════════════════════════════════════
# §2.3.4 — CLARIFICATION
# ══════════════════════════════════════════════════════════════

class ClarificationRequest(BaseModel):
    conversation_id: str
    slot_key:        str
    value:           str
    original_question: str
    source_id:       Optional[str] = None

@app.post("/nlu/clarify", tags=["NLU"])
async def apply_clarification(req: ClarificationRequest):
    """
    Applique une réponse de clarification et régénère le SQL.
    """
    nlu      = get_nlu_pipeline()
    resolver = AmbiguityResolver()
    context  = get_context(req.conversation_id)

    known_entities = []
    known_fields   = {}
    if req.source_id:
        try:
            pool = await get_pg_pool()
            async with pool.acquire() as conn:
                rows = await conn.fetch(
                    "SELECT name FROM source_entities WHERE source_id=$1 AND is_visible=TRUE ORDER BY name LIMIT 2000",
                    UUID(req.source_id)
                )
                known_entities = [r["name"] for r in rows]
        except Exception:
            pass

    # Reprocess la question originale
    slots = nlu.process(req.original_question, context, known_entities)

    # Applique la clarification
    slots = resolver.apply_clarification(slots, req.slot_key, req.value)

    # Enregistre le choix pour apprentissage §2.3.4
    try:
        from .preference_learner import PreferenceLearner
        pool = await get_pg_pool()
        learner = PreferenceLearner(pg_pool=pool)
        await learner.record_choice(
            user_id   = req.conversation_id or "default",
            source_id = req.source_id or "",
            slot_key  = req.slot_key,
            value     = req.value,
            question  = req.original_question,
        )
        logger.info(f"[Clarify] Préférence enregistrée — slot={req.slot_key} value={req.value}")
    except Exception as e:
        logger.warning(f"[Clarify] Preference record error: {e}")

    # Regénère avec les slots corrigés directement
    if req.source_id and not slots.ambiguities:
        source_id = UUID(req.source_id)
        pool = await get_pg_pool()

        # Charge le schéma
        async with pool.acquire() as conn:
            rows = await conn.fetch("""
                SELECT se.name AS table_name, ef.name AS field_name
                FROM source_entities se
                JOIN entity_fields ef ON ef.entity_id = se.id
                WHERE se.source_id = $1 AND se.is_visible = TRUE
                LIMIT 20000
            """, source_id)
        schema = {}
        for r in rows:
            schema.setdefault(r["table_name"], []).append(r["field_name"])

        # Applique _fix_slots avec le bon intent
        slots.raw_text = req.original_question
        slots = await _fix_slots(slots, req.original_question, schema, known_entities,
                                  source_id_str=req.source_id)

        # Dialecte
        source = await get_source(source_id)
        dialect_map = {"mssql":"mssql","sage_100":"mssql","mysql":"mysql",
                       "postgresql":"postgresql","sqlite":"sqlite"}
        dialect = dialect_map.get(
            source.connector_type.value if source and hasattr(source.connector_type,'value')
            else "mssql", "mssql"
        )

        # Génère le SQL avec les vrais slots
        sql_gen = SQLGenerator()
        result  = sql_gen.generate(slots, schema, dialect)

        return {
            "sql":         result["sql"],
            "explanation": result["explanation"],
            "warnings":    result.get("validation",{}).get("warnings",[]),
            "intent":      slots.intent,
            "tables":      slots.table_names,
            "dialect":     dialect,
            "needs_clarification": False,
            "response_ms": 0,
        }

    return {
        "slots_updated": True,
        "intent":        slots.intent,
        "tables":        slots.table_names,
        "ambiguities":   slots.ambiguities,
    }


# ══════════════════════════════════════════════════════════════
# STATS NLU
# ══════════════════════════════════════════════════════════════

@app.get("/nlu/stats", tags=["NLU"])
async def nlu_stats():
    """Statistiques d'utilisation du NLU."""
    pool = await get_pg_pool()
    try:
        async with pool.acquire() as conn:
            rows = await conn.fetch("SELECT * FROM nlu_intent_stats")
            total = await conn.fetchval("SELECT COUNT(*) FROM nlu_query_log")
        return {
            "total_queries": total,
            "by_intent": [dict(r) for r in rows],
        }
    except Exception as e:
        return {"total_queries": 0, "by_intent": [], "error": str(e)}


@app.delete("/nlu/context/{conversation_id}", tags=["NLU"])
async def clear_nlu_context(conversation_id: str):
    """Efface le contexte NLU d'une conversation."""
    clear_context(conversation_id)
    return {"success": True, "conversation_id": conversation_id}


# ══════════════════════════════════════════════════════════════
# SERVEUR FICHIERS STATIQUES UI
# ══════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════
# SPRINT 14 — FORECASTING
# ══════════════════════════════════════════════════════════════

class ForecastRequest(BaseModel):
    source_id:   str
    question:    str  = "prévoir les ventes"
    granularity: str  = "weekly"   # 'weekly' ou 'monthly'
    horizon:     int  = 12         # nombre de périodes à prévoir

@app.post("/forecast", tags=["Forecasting"])
async def forecast_endpoint(req: ForecastRequest):
    """
    Lance le pipeline de forecasting complet (SARIMA + Prophet + LSTM + Ensemble).
    Retourne les prévisions + chart_data pour Chart.js dans chat.html.
    """
    try:
        from .forecast_agent import run_forecast, detect_entity_and_sql, SQL_TEMPLATES

        # 1. Récupérer la source
        source = await get_source(UUID(req.source_id))
        if not source:
            raise HTTPException(404, "Source introuvable")

        pool = await get_pg_pool()

        # 2. Construire src_dict avec password
        src_dict = source.model_dump()
        try:
            async with pool.acquire() as _c:
                _sec = await _c.fetchrow(
                    "SELECT secret_value FROM connection_secrets "
                    "WHERE source_id=$1 AND secret_key='password'",
                    UUID(req.source_id)
                )
                if _sec:
                    src_dict["password"] = _sec["secret_value"]
        except Exception:
            pass

        # 3. Détecter le SQL adapté à la question
        sql = detect_entity_and_sql(
            source.name or "", req.question, req.granularity
        )

        # 4. Extraire les données via ConnectorFactory
        loop = asyncio.get_event_loop()
        connector = ConnectorFactory.create(src_dict)
        rows = await loop.run_in_executor(None, connector.execute_query, sql)
        if not rows:
            raise HTTPException(422, "Aucune donnée disponible pour le forecasting")

        # 5. Lancer le pipeline forecasting
        result = await run_forecast(
            source_id   = req.source_id,
            question    = req.question,
            rows        = rows,
            granularity = req.granularity,
            horizon     = req.horizon,
        )

        return {
            "success":    True,
            "best_model": result["best_model"],
            "mae":        result["mae"],
            "duration_s": result["duration_s"],
            "chart_data": result["chart_data"],
            "ranking":    result["ranking"],
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[/forecast] {e}")
        raise HTTPException(500, str(e))


import pathlib
UI_DIR = pathlib.Path(__file__).parent.parent / "ui"
app.mount("/", StaticFiles(directory=str(UI_DIR), html=True), name="ui")


# ══════════════════════════════════════════════════════════════
# ENTRYPOINT (dev uniquement — en prod: uvicorn api.main:app)
# ══════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)